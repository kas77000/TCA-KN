#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
TCA deck builder — KIC & NPS
============================================

Single-file, portable. Copy this one file to the target machine, drop the
order export next to it, and run:

    python tca_deck.py --data symboldetailstable_full_KIC.xlsx

Outputs, all into --out (default ./output):

    KIC_TCA_deck.pptx     the deck
    charts/*.png          every exhibit as a standalone 200dpi PNG,
                          so you can lift them into the corporate template
    tables.xlsx           every underlying table, one sheet each
    dark_regression.txt   full regression output (if statsmodels is present)
    run_log.txt           sanity report + reference check

To smoke-test without the real file (generates synthetic KIC-shaped data):

    python tca_deck.py --sample

Requires: pandas, numpy, matplotlib, python-pptx, openpyxl.
Optional: statsmodels (for the controlled dark regression; degrades gracefully).

------------------------------------------------------------------------------
METHOD NOTE — why the dark section is built the way it is
------------------------------------------------------------------------------
A raw dark-vs-lit comparison is misleading. Dark fills land disproportionately
in liquid, tight-spread, large-cap names that were easier to trade anyway, so a
naive split credits dark with a difficulty effect it did not cause.

This script therefore:
  1. compares ONLY within dark-capable algos (derived from the data), so venue
     is not confounded with algo choice;
  2. splits "no dark fill" vs "any dark fill" — a genuine control group, since
     Dark % is populated as 0 for orders that got none;
  3. repeats that split WITHIN buckets of order difficulty (%ADV, spread,
     market, size), so the gap cannot be a composition artefact;
  4. runs a notional-weighted regression with size / spread / participation /
     duration controls and market fixed effects, and reports the equal-weighted
     fit alongside it.

Even after all that the finding is associational: dark fills happen when contra
flow exists, and the presence of contra flow may itself signal an easier order.
The deck says "orders that achieved dark fills executed better, and this holds
after controls" — never "dark causes X bps".
"""
from __future__ import annotations

import argparse
import re
import sys
import textwrap
import warnings
from pathlib import Path

warnings.filterwarnings("ignore")

try:
    import numpy as np
    import pandas as pd
except ImportError as e:  # pragma: no cover
    sys.exit(f"Missing dependency: {e.name}\n  pip install pandas numpy")

# ===========================================================================
# 1. CONFIGURATION — this is the only block you should need to edit
# ===========================================================================

CURRENCY = "USD"
# Shown on the cover. Left blank so the repo carries no firm branding;
# set it locally if you want it on the slide.
BROKER   = ""

# ---------------------------------------------------------------------------
# CLIENT REGISTRY  —  pick one with --client
# ---------------------------------------------------------------------------
# Everything that differs between clients lives here. The report-sourced tables
# (venue segmentation, industry) are transcribed from each client's
# post-trade report because they are NOT in the order-level export; update
# them when the period changes.
#
#   algos        which algos to analyse. None = every algo in the file.
#   dark_story   run the dark section at all. Off for a client that barely
#                uses dark - four slides of "you did 1% of this" is not advice.
#   exclude_markets
#                markets the client does not trade and does not want reviewed.
#                A stray line in the report - a mis-tagged listing, an ADR -
#                is still a row on every country exhibit, and the client sees
#                a market they never traded. Dropped everywhere, and what was
#                dropped is written to the run log.
#   reference    published headline figures, recomputed and diffed on every run
#                so a mis-mapped column or an inverted sign is caught early.
CLIENTS = {
    "KIC": {
        "name": "KIC",
        "code": "KIC",
        "period": "5 January – 21 August 2026",
        "algos": None,
        "algo_order": ["VWAP", "IIS", "PROG"],
        "dark_story": True,
        "dark_markets": ["Hong Kong", "Japan", "Australia"],
        "venue_segments": {
            "IIS":   {"Auction": 28.2, "Visible Post": 0.0,  "Visible Take": 71.8, "Dark": 0.0},
            "PROG":  {"Auction": 0.0,  "Visible Post": 0.0,  "Visible Take": 100.0, "Dark": 0.0},
            "VWAP":  {"Auction": 23.4, "Visible Post": 39.9, "Visible Take": 21.9, "Dark": 14.8},
            "TOTAL": {"Auction": 23.8, "Visible Post": 36.2, "Visible Take": 26.7, "Dark": 13.4},
        },
        # algo performance, straight off the report's Performance Summary By
        # Algorithm. (algo, orders, exec value, %weight, part%, adv%, spread,
        # benchmark, impact bps, weighted impact bps)
        "algo_report": [
            ("VWAP", 1270, 457_503_439, 90.6, 2.1, 1.4, 8.6, "Order PVWAP",    6.5, 5.9),
            ("IIS",   774,  46_460_964,  9.2, 0.7, 1.2, 8.4, "Day Close Prm",  0.0, 0.0),
            ("PROG",   30,   1_217_989,  0.2, 0.6, 0.1, 20.8, "Arrival Price", 119.6, 0.3),
        ],
        # ---------------------------------------------------------------
        # Report tables. Every row is the same 10-tuple:
        #   (name, orders, shares, exec_value, weight_pct, period_part,
        #    adv_pct, spread_bps, bps, contrib_bps)
        # period_part is None where the report does not carry that column.
        # ---------------------------------------------------------------
        "totals": {"orders": 2074, "shares": 234_817_014,
                   "notional": 505_182_392, "period_part": 1.9, "adv_pct": 1.4,
                   "spread_bps": 8.6, "bps": 6.2, "pnl_ccy": 313_500},
        "country": [
            ("Japan",         235,   9_693_103, 160_058_861, 31.7, None, 1.4,  3.6,  -1.3, -0.4),
            ("Hong Kong",     434,  73_978_891, 136_664_201, 27.1, None, 1.3, 10.6,  14.7,  4.0),
            ("Taiwan",        166,  16_659_534,  75_970_221, 15.0, None, 1.0, 15.4,   9.5,  1.4),
            ("India",         486,  11_331_565,  68_224_945, 13.5, None, 1.8,  3.8,  -4.4, -0.6),
            ("Stock Connect", 535,  10_784_768,  32_910_676,  6.5, None, 0.3,  4.2,   0.9,  0.1),
            ("Indonesia",      58, 102_291_655,  16_648_022,  3.3, None, 4.5, 26.1,  31.0,  1.0),
            ("Malaysia",       63,   6_244_399,   6_071_868,  1.2, None, 2.8, 26.7,  62.5,  0.7),
            ("Thailand",       42,   2_725_457,   3_808_854,  0.8, None, 0.5, 33.7, -19.9, -0.2),
            ("Australia",      18,     166_073,   2_008_070,  0.4, None, 0.1,  5.8,  -0.8, -0.0),
            ("Philippines",    24,     722_805,   1_338_151,  0.3, None, 2.8,  7.4,  55.2,  0.1),
            ("Singapore",       9,     128_600,     971_684,  0.2, None, 0.2,  9.7,   0.6,  0.0),
            ("New Zealand",     4,      90_164,     506_838,  0.1, None, 2.2, 18.3,  -2.2, -0.0),
        ],
        "marketcap": [
            ("Large $10-100B", 1277, 119_676_862, 324_093_275, 64.2, None, 1.1,  8.2,  7.7,  4.9),
            ("Mid $2-10B",      609,  98_903_018, 108_411_582, 21.5, None, 3.0, 11.2,  6.2,  1.3),
            ("Mega >$100B",     183,  14_243_091,  72_027_046, 14.3, None, 0.3,  6.7, -0.4, -0.1),
            ("Small $0.3-2B",     5,   1_994_043,     650_489,  0.1, None, 0.6, 10.1, 10.1,  0.0),
        ],
        "adv": [
            ("0-1%",    1831, 81_682_459, 280_901_115, 56.0, 1.0,  0.4,  7.9,  5.6, 3.1),
            ("1-3%",     180, 85_477_868, 155_316_114, 31.0, 2.5,  1.7,  8.8,  2.7, 0.8),
            ("3-5%",      42, 38_619_623,  44_481_654,  9.0, 3.6,  4.1, 10.8, 19.7, 1.7),
            ("5-10%",     18, 25_985_264,  22_578_625,  4.0, 5.9,  5.9, 11.6,  9.7, 0.4),
            ("10-25%",     3,  3_051_800,   1_904_883,  0.0, 11.7, 15.7, 11.6, 25.0, 0.1),
        ],
        "side": [
            ("BUY",  1365, 113_356_784, 269_638_873, 53.4, 1.8, 1.3, 8.7, 8.4, 4.5),
            ("SELL",  709, 121_460_230, 235_543_520, 46.6, 2.0, 1.5, 8.6, 3.7, 1.7),
        ],
        # (market, side, % of total value, bps) - from the country/side
        # breakdown. Market totals implied here reconcile with "country" above.
        # (market, auction%, post%, take%, dark%) - rows sum to 100
        "venue_country": [
            ("Hong Kong",     13.9, 54.4,  10.4, 21.3),
            ("Japan",         39.9, 32.9,   3.3, 23.9),
            ("Taiwan",        29.0, 46.5,  24.5,  0.0),
            ("India",          0.0,  0.0, 100.0,  0.0),
            ("Stock Connect", 33.0, 57.8,   9.2,  0.0),
            ("Indonesia",      2.6,  0.0,  97.0,  0.5),
            ("Malaysia",       0.0,  0.0, 100.0,  0.0),
            ("Thailand",      61.9, 19.6,  18.5,  0.0),
            ("Philippines",   70.3, 27.3,   2.4,  0.0),
            ("Australia",     27.6,  9.2,  51.4, 11.8),
            ("Singapore",      0.0,  0.0, 100.0,  0.0),
            ("New Zealand",    0.0,  0.0, 100.0,  0.0),
        ],
        "country_side": [
            ("Hong Kong",     "BUY", 14.2,  30.2), ("Hong Kong",     "SELL", 12.8,  -2.6),
            ("Japan",         "BUY", 16.3,   0.2), ("Japan",         "SELL", 15.3,  -3.0),
            ("Taiwan",        "BUY",  8.4,  13.1), ("Taiwan",        "SELL",  6.7,   4.9),
            ("India",         "BUY",  7.2, -17.3), ("India",         "SELL",  6.3,  10.5),
            ("Stock Connect", "BUY",  4.6,   1.9), ("Stock Connect", "SELL",  1.9,  -1.5),
            ("Indonesia",     "BUY",  1.4,  15.4), ("Indonesia",     "SELL",  1.9,  42.3),
            ("Malaysia",      "BUY",  0.3,  -0.4), ("Malaysia",      "SELL",  0.9,  82.6),
            ("Thailand",      "BUY",  0.2,  -4.7), ("Thailand",      "SELL",  0.5, -27.7),
            ("Philippines",   "BUY",  0.0, -13.3), ("Philippines",   "SELL",  0.2,  68.3),
            ("Australia",     "BUY",  0.4,  -0.5), ("Australia",     "SELL",  0.0, -12.2),
            ("Singapore",     "BUY",  0.2,   0.6),
            ("New Zealand",   "BUY",  0.1,  -2.2),
        ],
        # industry: (name, issues, %weight, impact bps, weighted impact bps)
        "industry": [
            ("Consumer Non-cyclical", 141,  9.4,   8.6,  0.8),
            ("Utilities",              37,  1.6,   8.7,  0.1),
            ("Energy",                 48,  2.3, -11.4, -0.3),
            ("Basic Materials",       107,  7.7,   5.4,  0.4),
            ("Industrial",            187, 17.5,   2.5,  0.4),
            ("Technology",            118, 14.3,   7.3,  1.0),
            ("Communications",         67, 11.8,  29.0,  3.4),
            ("Financial",             205, 16.6,  -6.4, -1.1),
            ("Diversified",             5,  0.2,  62.9,  0.1),
            ("Consumer Cyclical",     131, 18.5,   6.2,  1.1),
        ],
        "reference": {
            "orders": 2074, "notional": 505_182_392, "shares": 234_817_014,
            "impact_bps": 6.2, "spread_bps": 8.6, "adv_pct": 1.4,
            "dark_share_pct": 13.4,
        },
    },
    "NPS": {
        "name": "NPS",
        "code": "NPS",
        "period": "9 January – 25 August 2026",
        # NPS also runs PROG and SDMA, but they are a rounding error and the
        # review is scoped to the two algos that carry the flow.
        "algos": ["IIS", "VWAP"],
        "algo_order": ["VWAP", "IIS"],
        # 1.0% of value in dark across the whole book - there is no dark story
        # to tell, so the section is switched off and the deck concentrates on
        # the IIS vs VWAP choice instead.
        "dark_story": False,
        # NPS trades APAC only. The single USA line in the report is a stray -
        # one order, USD 38k - and putting a market they do not trade on a
        # client exhibit invites a question with no good answer.
        "exclude_markets": ["USA"],
        "dark_markets": ["Hong Kong", "Japan", "Australia"],
        "venue_segments": {
            "IIS":   {"Auction": 89.0, "Visible Post": 0.0,  "Visible Take": 11.0, "Dark": 0.0},
            "VWAP":  {"Auction": 7.5,  "Visible Post": 67.9, "Visible Take": 23.5, "Dark": 1.1},
            "TOTAL": {"Auction": 17.2, "Visible Post": 59.7, "Visible Take": 22.1, "Dark": 1.0},
        },
        "algo_report": [
            ("VWAP", 3761, 962_121_454, 87.9, 2.7, 1.9,  4.3, "Order PVWAP",   -1.7, -1.5),
            ("IIS",   274, 130_978_627, 12.0, 0.5, 0.5,  5.4, "Day Close Prm",  0.0,  0.0),
            ("PROG",   14,   1_375_612,  0.1, 0.5, 0.4, 12.3, "Arrival Price", 100.3,  0.1),
            ("SDMA",    3,     370_068,  0.0, 0.6, 0.1, 21.2, "Arrival Price", -41.5, -0.0),
        ],
        # ---------------------------------------------------------------
        # Report tables. Every row is the same 10-tuple:
        #   (name, orders, shares, exec_value, weight_pct, period_part,
        #    adv_pct, spread_bps, bps, contrib_bps)
        # period_part is None where the report does not carry that column.
        # ---------------------------------------------------------------
        "totals": {"orders": 4052, "shares": 87_472_636,
                   "notional": 1_094_845_761, "period_part": 2.4, "adv_pct": 1.7,
                   "spread_bps": 4.4, "bps": -1.4, "pnl_ccy": -153_278},
        "country": [
            ("Japan",        2061, 37_115_601, 755_028_595, 69.0, None, 1.8,  3.5,  -1.1, -0.7),
            ("Australia",     497,  7_639_685, 135_661_180, 12.4, None, 2.0,  4.1,  -0.7, -0.1),
            ("Hong Kong",     535, 13_774_357,  72_057_394,  6.6, None, 0.8,  6.8,  -1.8, -0.1),
            ("India",         298, 10_913_335,  63_476_005,  5.8, None, 2.4,  4.2,   1.5,  0.1),
            ("Taiwan",        235,  6_169_000,  36_650_859,  3.3, None, 0.2, 14.5, -11.3, -0.4),
            ("Singapore",     168,  3_240_432,  21_504_716,  2.0, None, 0.9, 10.3,  -9.2, -0.2),
            ("Stock Connect", 150,  2_071_016,   4_897_573,  0.4, None, 0.0,  5.7,   1.6,  0.0),
            ("Malaysia",       59,  1_609_500,   3_260_927,  0.3, None, 1.6, 17.4,  -0.0, -0.0),
            ("Indonesia",      27,  4_665_600,   1_127_942,  0.1, None, 0.4, 22.2,  -4.9, -0.0),
            ("New Zealand",    10,    123_253,     872_356,  0.1, None, 2.2, 15.9,  -2.8, -0.0),
            ("Philippines",    11,    150_250,     270_107,  0.0, None, 0.5,  6.7,  -0.9, -0.0),
            ("USA",             1,        607,      38_107,  0.0, None, 0.2,  4.6,   2.2,  0.0),
        ],
        "marketcap": [
            ("Large $10-100B", 3072, 65_894_658, 771_943_372, 70.5, None, 2.0, 4.4,  0.1,  0.1),
            ("Mega >$100B",     498, 10_732_155, 262_458_438, 24.0, None, 0.4, 4.1, -5.9, -1.4),
            ("Mid $2-10B",      482, 10_845_823,  60_443_952,  5.5, None, 4.0, 6.6, -1.5, -0.1),
        ],
        "adv": [
            ("0-1%",   3822, 60_286_511, 683_670_707, 62.0,  0.6,  0.3, 4.6, -2.7, -1.7),
            ("1-3%",    164, 13_720_140, 211_323_382, 19.0,  3.0,  1.8, 4.2,  0.3,  0.1),
            ("3-5%",     33,  4_179_241,  72_143_438,  7.0,  6.6,  4.0, 4.0, -2.6, -0.2),
            ("5-10%",    28,  8_113_156, 105_474_814, 10.0,  7.8,  6.9, 4.1,  4.8,  0.5),
            ("10-25%",    5,  1_173_588,  22_233_421,  2.0, 14.3, 12.4, 4.0, -4.3, -0.1),
        ],
        "side": None,   # the side summary itself was not captured
        "venue_country": [
            ("Japan",         17.7, 76.7,   4.7,  0.9),
            ("Australia",     11.2, 15.8,  72.0,  1.0),
            ("Hong Kong",     37.0, 51.7,   8.1,  3.3),
            ("India",          0.0,  0.0, 100.0,  0.0),
            ("Taiwan",        34.0, 30.8,  35.2,  0.0),
            ("Singapore",      0.0,  2.0,  98.0,  0.0),
            ("Stock Connect", 15.3, 74.7,  10.1,  0.0),
            ("Malaysia",       0.0,  0.0, 100.0,  0.0),
            ("Indonesia",      0.0,  1.3,  98.7,  0.0),
            ("New Zealand",    0.0, 12.7,  87.3,  0.0),
            ("Philippines",    0.0, 93.2,   6.8,  0.0),
            ("USA",            0.0,  0.0, 100.0,  0.0),
        ],
        "country_side": [
            ("Japan",         "BUY", 38.4,  -3.7), ("Japan",         "SELL", 30.5,   2.3),
            ("Australia",     "BUY",  8.4,  -0.5), ("Australia",     "SELL",  4.0,  -1.0),
            ("Hong Kong",     "BUY",  5.7,  -2.4), ("Hong Kong",     "SELL",  0.9,   1.8),
            ("India",         "BUY",  2.4,   2.0), ("India",         "SELL",  3.4,   1.2),
            ("Taiwan",        "BUY",  3.3, -11.3),
            ("Singapore",     "BUY",  1.3, -13.8), ("Singapore",     "SELL",  0.7,  -0.3),
            ("Stock Connect", "BUY",  0.4,   4.3), ("Stock Connect", "SELL",  0.1, -17.1),
            ("Malaysia",      "BUY",  0.3,  -0.0),
            ("Indonesia",     "BUY",  0.1,  -4.9),
            ("New Zealand",   "BUY",  0.1,  -2.8),
            ("Philippines",   "BUY",  0.0,  -0.9),
            ("USA",           "BUY",  0.0,   2.2),
        ],

        "industry": [
            ("Consumer Non-cyclical",  91, 12.5,   2.8,  0.4),
            ("Utilities",              21,  2.0,  -2.8, -0.1),
            ("Energy",                 23,  2.0,  12.1,  0.2),
            ("Basic Materials",        52,  8.4,  -1.2, -0.1),
            ("Industrial",             96, 18.9,   1.9,  0.4),
            ("Technology",             56, 12.7, -10.2, -1.3),
            ("Financial",             160, 21.4,  -2.6, -0.6),
            ("Communications",         38,  4.8,  -2.3, -0.1),
            ("Diversified",             1,  0.1,  -8.8,  0.0),
            ("Consumer Cyclical",      98, 17.2,  -1.5, -0.3),
        ],
        "reference": {
            "orders": 4052, "notional": 1_094_845_761, "shares": 87_472_636,
            "impact_bps": -1.4, "spread_bps": 4.4, "adv_pct": 1.7,
            "dark_share_pct": 1.0,
        },
    },
}

DEFAULT_CLIENT = "KIC"

# Populated from CLIENTS[--client] at startup. Defaults keep the module
# importable and let the sample run stand alone.
CLIENT_NAME   = "KIC"
CLIENT_CODE   = "KIC"
PERIOD_LABEL  = "5 January – 21 August 2026"
ALGOS_STUDIED = None
DARK_STORY    = True
INDUSTRY_REPORT = None
ALGO_REPORT     = None
REPORT          = {}      # country / marketcap / adv / side / totals

# --- Column mapping -------------------------------------------------------
# Header names are matched case-insensitively and ignoring spaces, dots,
# underscores and brackets, so "Dark %", "dark%", "DARK_%" all resolve.
# Each entry may be a single name or a list of acceptable alternatives.
COLUMNS = {
    "date":            ["Date"],
    "strike_time":     ["Strike Time"],
    "end_time":        ["End Time"],
    "duration_min":    ["Dur (mins)", "Dur(mins)", "Duration"],
    "symbol":          ["Symbol"],
    "side":            ["Side"],
    "order_type":      ["Order Type"],
    "order_qty":       ["Order Qty"],
    "shares_exec":     ["Shares Exec"],
    "notional":        ["Value Exec"],
    "dark_qty":        ["Dark Qty"],
    "dark_pct":        ["Dark %"],
    "dark_value":      ["Dark Value"],
    "weight":          ["Weight"],
    "adv_pct":         ["% Adv", "%Adv"],
    "period_part":     ["Period Part"],
    "avg_price":       ["Avg. Price (local)", "Avg Price (local)"],
    "fx":              ["FX"],
    "algo":            ["ClientAlgo", "Client Algo"],
    "spread_bps":      ["Hist Spread"],
    "period_spread":   ["Period Spread"],
    "market":          ["Market"],
    # benchmark reference prices - needed to tell a self-referential benchmark
    # (algo executes AT the benchmark) from an unpopulated one
    "px_arrival":      ["Arrival Price"],
    "px_pvwap":        ["PVWAP Price"],
    "px_vwap":         ["VWAP Price"],
    "px_close":        ["Close Price"],
    # benchmark slippages, in bps
    "slip_arrival":    ["Arrival ImpBps"],
    "slip_pvwap":      ["PVWAP ImpBps"],
    "slip_vwap":       ["VWAP ImpBps"],
    "slip_close":      ["Close ImpBps"],
    # benchmark slippages, in currency
    "amt_arrival":     ["Arrival ImpAmt"],
    "amt_pvwap":       ["PVWAP ImpAmt"],
    "amt_vwap":        ["VWAP ImpAmt"],
    "amt_close":       ["Close ImpAmt"],
}

# Columns without which the deck cannot be built. Beyond these the loader
# needs at least one slippage column and at least one dark column, but any of
# several will do - so those are checked as groups, after the derivations have
# had their chance (dark_pct can be rebuilt from Dark Value or Dark Qty).
REQUIRED = ["notional", "algo"]
REQUIRED_ONE_OF = {
    "a benchmark slippage": ["slip_pvwap", "slip_arrival", "slip_vwap",
                             "slip_close"],
    "a dark measure": ["dark_pct", "dark_value", "dark_qty"],
}

# The export writes a merged "Symbol Details" banner on row 1 and the real
# header on row 2. Set to 0 if your file has the header on row 1.
HEADER_ROW = 1          # 0-indexed -> row 2 of the sheet
SHEET      = 0

# --- Conventions ----------------------------------------------------------
# In this export a POSITIVE number is a SAVING (the order beat its
# benchmark) and a NEGATIVE number is a COST. This matches the report, where
# the negative country and industry figures are the ones printed in red.
#
# Everything downstream keeps that convention, and the deck presents it that
# way too, so every figure reconciles with the post-trade report the client
# already has. Flip this only if your export uses the opposite sign - the
# reference check re-derives it from the data and shouts on a mismatch.
POSITIVE_IS_SAVING = True

BUY_VALUES = {"B", "BUY", "BOT", "1"}

# Mixed benchmark, mirroring the mandated benchmark per algo in the report.
PRIMARY_BENCHMARK = {
    "VWAP": "slip_pvwap",     # Order PVWAP
    "IIS":  "slip_close",     # Day Close Prm
    "PROG": "slip_arrival",   # Arrival Price
}
DEFAULT_BENCHMARK = "slip_pvwap"

BENCHMARK_LABELS = {
    "slip_arrival": "vs Arrival",
    "slip_pvwap":   "vs Interval PVWAP",
    "slip_vwap":    "vs Day VWAP",
    "slip_close":   "vs Close",
}
MATRIX_BENCHMARKS = ["slip_arrival", "slip_pvwap", "slip_vwap", "slip_close"]

# Algos with fewer than this many orders are shown but flagged as not
# statistically quotable. PROG (n=30, 0.2% of notional) trips this.
MIN_N_QUOTABLE = 100
ALGO_ORDER = ["VWAP", "IIS", "PROG"]

# --- Buckets --------------------------------------------------------------
ADV_BUCKETS = [-0.001, 1, 3, 5, 10, 25, 1e9]
ADV_LABELS  = ["0-1%", "1-3%", "3-5%", "5-10%", "10-25%", "25%+"]

DARK_BUCKETS = [-0.001, 0.0001, 10, 25, 50, 100.001]
DARK_LABELS  = ["0% (none)", "0-10%", "10-25%", "25-50%", "50%+"]

# An algo counts as dark-capable if at least this share of its orders got a
# non-zero dark fill. Set DARK_ALGOS to a list to override the derivation.
DARK_ALGOS = None
DARK_MIN_NONZERO_SHARE = 0.05

# Dark is only offered in these markets. Every dark comparison is restricted to
# them, because an order in a market with no dark access never had the choice -
# including it would compare markets rather than venues, and would flatter or
# damn dark purely on where the flow happened to sit.
# Names are matched loosely, so "Hong Kong" and "HongKong" both work.
# Set to None to use every market in the file.
DARK_MARKETS = ["Hong Kong", "Japan", "Australia"]

# Markets to drop from the review entirely. Populated from the client config.
EXCLUDE_MARKETS = []
# Below this share of zero-dark orders, fall back to a median split.
DARK_ZERO_SPLIT_THRESHOLD = 0.15
DARK_CONTROLS = ["adv_bucket", "spread_tercile", "size_tercile", "market"]

# --- Statistics -----------------------------------------------------------
WINSOR       = (0.01, 0.99)   # None to disable
BOOTSTRAP_N  = 2000
SEED         = 7
MIN_N_FOR_CI = 8
MIN_N_MARKET = 10             # markets below this are pooled into "Other"
COUNTRY_MIN_SHARE = 0.005

# --- Reference check ------------------------------------------------------
# Headline figures read off the published post-trade report for this period.
# The script recomputes each from the order file and reports the difference, so
# a mis-mapped column or a wrong sign convention is caught before any slide is
# built. Set to None to skip.
REFERENCE = {
    "orders":        2074,
    "notional":      505_182_392,
    "shares":        234_817_014,
    "impact_bps":    6.2,
    "spread_bps":    8.6,
    "adv_pct":       1.4,
    "dark_share_pct": 13.4,
}
REFERENCE_TOL = {"orders": 0, "notional": 0.005, "shares": 0.005,
                 "impact_bps": 0.5, "spread_bps": 0.5, "adv_pct": 0.3,
                 "dark_share_pct": 1.0}

# --- Palette --------------------------------------------------------------
# Validated with the six checks (lightness band, chroma floor, CVD Delta E,
# normal-vision Delta E, contrast) on surface #fcfcfb. The venue set was
# re-stepped off amber onto indigo: orange/amber scored 13.7 on normal vision,
# below the 15 floor. Final venue set: worst CVD 9.2, worst normal 24.0 — pass.
# The green carries a contrast WARN (2.74 vs 3.0), which obliges the direct
# labels and the table that accompany every chart using it.
# Sampled from the corporate title slide: azure band, charcoal panel, light
# grey accent, white ground. Neutrals are pulled toward that grey rather than
# the warm off-white used before, so charts sit on the same page as the
# template.
SURFACE     = "#ffffff"
INK         = "#111315"
INK_SECOND  = "#4a4d50"
INK_MUTED   = "#8a8d90"
GRID        = "#e6e7e8"
BASELINE    = "#cccdcc"

# Diverging — cost vs savings. Neutral gray midpoint, never a hue.
# Corporate azure carries "good"; the red is stepped a little deeper than the
# previous one so it holds its contrast against a pure white ground.
# Validated against white: azure lightness 0.656, chroma 0.147, contrast 3.13;
# the pair separates 34.4 for normal vision and 25.7 under protan/deutan.
SAVE_COLOR  = "#2b98e4"
COST_COLOR  = "#c62828"
NEUTRAL     = "#f2f3f4"

# Categorical — venue segments. Fixed order, never cycled.
VENUE_COLORS = {
    "Auction":      "#2b98e4",
    "Visible Post": "#1baf7a",
    "Visible Take": "#eb6834",
    "Dark":         "#4a3aa7",
}
VENUE_ORDER = ["Auction", "Visible Post", "Visible Take", "Dark"]

# Colour follows the entity: dark stays indigo everywhere it appears.
DARK_GROUP_COLORS = {"No dark": "#2b98e4", "Any dark": "#4a3aa7",
                     "Low dark": "#2a78d6", "High dark": "#4a3aa7"}

FIGSIZE = (10.0, 5.6)
DPI     = 200


# ===========================================================================
# 2. LOADING AND NORMALISATION
# ===========================================================================

LOG: list = []


def log(msg: str = "") -> None:
    """Print, then keep for run_log.txt.

    Windows consoles often run cp1252, where an em-dash raises
    UnicodeEncodeError and would kill the run. Degrade the console copy to
    ASCII if needed; run_log.txt is always written as full UTF-8.
    """
    msg = str(msg)
    try:
        print(msg)
    except UnicodeEncodeError:
        enc = getattr(sys.stdout, "encoding", None) or "ascii"
        print(msg.encode(enc, errors="replace").decode(enc, errors="replace"))
    LOG.append(msg)


def _norm_key(s) -> str:
    """Header -> comparison key: lowercase, alphanumerics only."""
    return re.sub(r"[^a-z0-9]", "", str(s).lower())


def resolve_columns(df: pd.DataFrame):
    """Map logical name -> actual column, tolerant of spacing and case."""
    lookup = {}
    for actual in df.columns:
        lookup.setdefault(_norm_key(actual), actual)

    resolved, missing = {}, []
    for logical, candidates in COLUMNS.items():
        if isinstance(candidates, str):
            candidates = [candidates]
        for cand in candidates:
            hit = lookup.get(_norm_key(cand))
            if hit is not None:
                resolved[logical] = hit
                break
        else:
            missing.append(logical)
    return resolved, missing


def _to_num(s: pd.Series) -> pd.Series:
    """Coerce to float, tolerating thousands separators, %, parens and blanks."""
    if pd.api.types.is_numeric_dtype(s):
        return s.astype(float)
    t = (s.astype(str)
           .str.strip()
           .str.replace(",", "", regex=False)
           .str.replace("%", "", regex=False)
           .str.replace(r"^\((.*)\)$", r"-\1", regex=True)
           .replace({"": None, "-": None, "nan": None, "None": None,
                     "#N/A": None, "N/A": None, "NULL": None}))
    return pd.to_numeric(t, errors="coerce")


def _as_percent(s: pd.Series, name: str) -> pd.Series:
    """Percent fields may arrive as 0-1 fractions or 0-100. Detect and align."""
    v = _to_num(s)
    finite = v[np.isfinite(v)]
    if len(finite) and finite.max() <= 1.5 and (finite > 0).any():
        log(f"  {name}: looks like a 0-1 fraction (max {finite.max():.3f})"
            f" -> scaled x100")
        return v * 100
    return v


def load_orders(path: Path) -> pd.DataFrame:
    """Read the symbol-details export into a normalised frame."""
    log(f"Reading {path}")
    if path.suffix.lower() in (".xlsx", ".xlsm", ".xls"):
        raw = pd.read_excel(path, sheet_name=SHEET, header=HEADER_ROW)
    else:
        raw = pd.read_csv(path, header=HEADER_ROW, encoding="utf-8-sig")

    raw = raw.loc[:, ~raw.columns.astype(str).str.match(r"^Unnamed")]
    raw = raw.dropna(how="all")

    cols, missing = resolve_columns(raw)
    log(f"  {len(raw):,} rows x {len(raw.columns)} columns; "
        f"resolved {len(cols)}/{len(COLUMNS)} known fields")
    if missing:
        log(f"  not found (ignored unless required): {', '.join(missing)}")
    hard = [c for c in REQUIRED if c not in cols]
    for label, group in REQUIRED_ONE_OF.items():
        if not any(c in cols for c in group):
            hard.append(f"{label} (one of: {', '.join(group)})")
    if hard:
        sys.exit(
            "\nFATAL - required columns could not be matched:\n  - "
            + "\n  - ".join(hard)
            + "\n\nHeaders actually present:\n  "
            + "\n  ".join(str(c) for c in raw.columns)
            + "\n\nEdit COLUMNS at the top of this script, or set HEADER_ROW "
              "if the header is not on row 2.")

    df = pd.DataFrame(index=raw.index)
    for logical, actual in cols.items():
        df[logical] = raw[actual]

    for c in ["order_qty", "shares_exec", "notional", "dark_qty", "dark_value",
              "duration_min", "avg_price", "fx", "spread_bps", "period_spread",
              "px_arrival", "px_pvwap", "px_vwap", "px_close",
              "slip_arrival", "slip_pvwap", "slip_vwap", "slip_close",
              "amt_arrival", "amt_pvwap", "amt_vwap", "amt_close"]:
        if c in df:
            df[c] = _to_num(df[c])
    for c in ["dark_pct", "adv_pct", "period_part", "weight"]:
        if c in df:
            df[c] = _as_percent(df[c], c)

    if "algo" in df:
        df["algo"] = df["algo"].astype(str).str.strip().str.upper()
    if "market" in df:
        df["market"] = df["market"].astype(str).str.strip()
    if "side" in df:
        up = df["side"].astype(str).str.strip().str.upper()
        df["side"] = np.where(up.isin(BUY_VALUES), "BUY", "SELL")
    if "date" in df:
        df["date"] = pd.to_datetime(df["date"], errors="coerce")

    before = len(df)
    df = df[np.isfinite(df["notional"]) & (df["notional"] > 0)].copy()
    if len(df) < before:
        log(f"  dropped {before - len(df):,} rows with no executed notional")

    # --- dark share --------------------------------------------------------
    if "dark_pct" not in df or df["dark_pct"].isna().all():
        if "dark_value" in df and "notional" in df:
            df["dark_pct"] = 100 * df["dark_value"] / df["notional"]
            log("  dark_pct derived from Dark Value / Value Exec")
        elif "dark_qty" in df and "shares_exec" in df:
            df["dark_pct"] = 100 * df["dark_qty"] / df["shares_exec"]
            log("  dark_pct derived from Dark Qty / Shares Exec")
        else:
            df["dark_pct"] = 0.0
    df["dark_pct"] = df["dark_pct"].fillna(0.0).clip(0, 100)

    if "dark_value" not in df or df["dark_value"].isna().all():
        df["dark_value"] = df["notional"] * df["dark_pct"] / 100
    df["dark_value"] = df["dark_value"].fillna(0.0)

    df["has_dark"] = df["dark_pct"] > 0
    df["dark_bucket"] = pd.cut(df["dark_pct"], DARK_BUCKETS, labels=DARK_LABELS)

    if "adv_pct" in df:
        df["adv_bucket"] = pd.cut(df["adv_pct"], ADV_BUCKETS, labels=ADV_LABELS)

    # --- mandated benchmark per algo --------------------------------------
    df["slip_primary"] = np.nan
    df["benchmark"] = DEFAULT_BENCHMARK
    for algo, bench in PRIMARY_BENCHMARK.items():
        m = df["algo"] == algo
        if m.any() and bench in df:
            df.loc[m, "slip_primary"] = df.loc[m, bench]
            df.loc[m, "benchmark"] = bench
    fallback = DEFAULT_BENCHMARK if DEFAULT_BENCHMARK in df else next(
        (c for c in REQUIRED_ONE_OF["a benchmark slippage"] if c in df), None)
    if fallback:
        gap = df["slip_primary"].isna() & df[fallback].notna()
        if gap.any():
            df.loc[gap, "slip_primary"] = df.loc[gap, fallback]
            if fallback != DEFAULT_BENCHMARK:
                log(f"  {DEFAULT_BENCHMARK} not present -> using {fallback} "
                    "for orders with no mapped benchmark")

    if not POSITIVE_IS_SAVING:
        for c in ["slip_arrival", "slip_pvwap", "slip_vwap", "slip_close",
                  "slip_primary"]:
            if c in df:
                df[c] = -df[c]

    # positive = money saved against the benchmark
    df["pnl_ccy"] = df["slip_primary"] * df["notional"] / 1e4
    return df.reset_index(drop=True)


def tercile(s: pd.Series, labels=("Low", "Mid", "High")) -> pd.Series:
    try:
        return pd.qcut(s, 3, labels=list(labels), duplicates="drop")
    except (ValueError, TypeError):
        return pd.Series(pd.NA, index=s.index)


def add_control_buckets(df: pd.DataFrame) -> pd.DataFrame:
    if "spread_bps" in df:
        df["spread_tercile"] = tercile(df["spread_bps"])
    if "duration_min" in df:
        df["dur_tercile"] = tercile(df["duration_min"])
    df["size_tercile"] = tercile(df["notional"])
    return df


def derive_dark_algos(df: pd.DataFrame) -> list:
    """Which algos actually route to dark, read off the data."""
    if DARK_ALGOS is not None:
        return list(DARK_ALGOS)
    share = df.groupby("algo")["has_dark"].mean()
    algos = sorted(share[share >= DARK_MIN_NONZERO_SHARE].index.tolist())
    log(f"  dark-capable algos (>= {DARK_MIN_NONZERO_SHARE:.0%} of orders "
        f"with a dark fill): " + (", ".join(algos) if algos else "NONE"))
    for a, s in share.sort_values(ascending=False).items():
        log(f"    {a:<8} {s:6.1%} of orders got a dark fill")
    return algos


# ===========================================================================
# 3. SANITY AND REFERENCE CHECK
# ===========================================================================

def sanity_report(df: pd.DataFrame) -> None:
    log("")
    log("-" * 74)
    log("SANITY REPORT")
    log("-" * 74)
    log(f"  orders              {len(df):,}")
    log(f"  executed notional   {CURRENCY} {df['notional'].sum():,.0f}")
    if "shares_exec" in df:
        log(f"  executed shares     {df['shares_exec'].sum():,.0f}")
    if "date" in df and df["date"].notna().any():
        log(f"  date range          {df['date'].min():%Y-%m-%d} to "
            f"{df['date'].max():%Y-%m-%d}")
    log(f"  algos               " + ", ".join(
        f"{a} ({n:,})" for a, n in df["algo"].value_counts().items()))
    if "market" in df:
        log(f"  markets             {df['market'].nunique()}")

    # sign convention cross-check: buys and sells should not be systematically
    # opposite once side-adjusted. If they are, the file is not side-adjusted.
    if "side" in df and df["side"].nunique() == 2:
        b = wmean(df.loc[df.side == "BUY", "slip_primary"],
                  df.loc[df.side == "BUY", "notional"])
        s = wmean(df.loc[df.side == "SELL", "slip_primary"],
                  df.loc[df.side == "SELL", "notional"])
        log(f"  cost by side        BUY {b:+.1f} bps / SELL {s:+.1f} bps")
        if np.isfinite(b) and np.isfinite(s) and b * s < 0 and min(abs(b), abs(s)) > 3:
            log("  ** WARNING: buys and sells have opposite signs. The export "
                "may not be side-adjusted.")

    miss = {c: df[c].isna().mean() for c in
            ["slip_arrival", "slip_pvwap", "slip_vwap", "slip_close",
             "adv_pct", "spread_bps", "duration_min"] if c in df}
    bad = {c: v for c, v in miss.items() if v > 0.02}
    if bad:
        log("  missing values      " + ", ".join(
            f"{c} {v:.0%}" for c, v in sorted(bad.items(), key=lambda x: -x[1])))

    # A benchmark scoring a flat zero can mean one of two very different
    # things. Test which, rather than guessing.
    rec = reconcile_algo(df, ALGO_REPORT)
    if not rec.empty:
        log("")
        log("  PER-ALGO RECONCILIATION  (order file vs published report)")
        for algo, r in rec.iterrows():
            bad = (abs(r["orders_diff"]) > 0 or abs(r["notional_diff_%"]) > 0.5
                   or abs(r["bps_diff"]) > 0.5)
            log(f"    {algo:<8} orders {int(r['orders_got']):>6,} vs "
                f"{int(r['orders_pub']):>6,}   "
                f"notional {r['notional_diff_%']:+6.2f}%   "
                f"bps {r['bps_got']:+7.2f} vs {r['bps_pub']:+6.2f}"
                + ("   <== MISMATCH" if bad else ""))
        if (rec["bps_diff"].abs() > 0.5).any():
            log("      -> a per-algo gap that the totals hide usually means the "
                "benchmark mapping is wrong for that algo")

    deg = benchmark_degeneracy(df)
    flagged = deg[deg["verdict"] != "normal"] if not deg.empty else deg
    if len(flagged):
        log("")
        log("  BENCHMARK CHECK")
        for _, r in flagged.iterrows():
            log(f"    {r['algo']:<8} vs {r['benchmark']:<13} "
                f"{r['zero_slippage_%']:.0f}% of orders score exactly zero; "
                f"exec price equals benchmark price on "
                f"{r['exec_price_equals_bench_%']:.0f}%")
            if "construction" in r["verdict"]:
                log(f"      -> {r['algo']} executes AT its benchmark, so the "
                    "zero is arithmetic, not performance.")
                log(f"      -> This benchmark cannot rank {r['algo']}. Judge it "
                    "on arrival instead: the question for a close strategy is "
                    "whether waiting for the close paid, not whether it got "
                    "the close.")
            else:
                log(f"      -> exec and benchmark prices DIFFER, so the zero is "
                    "not arithmetic. The field looks genuinely unpopulated - "
                    "raise it before publishing.")


# which reference price belongs to which slippage column
BENCH_PRICE = {"slip_arrival": "px_arrival", "slip_pvwap": "px_pvwap",
               "slip_vwap": "px_vwap", "slip_close": "px_close"}


def benchmark_degeneracy(df: pd.DataFrame) -> pd.DataFrame:
    """Is an algo's benchmark self-referential?

    An algo that executes at the close, measured against the close, scores
    exactly zero on every order. That is arithmetic, not performance - and it
    is very different from a benchmark field nobody populated. Comparing the
    average execution price against the benchmark price separates the two.
    """
    rows = []
    for algo, sub_df in df.groupby("algo", observed=True):
        bench = PRIMARY_BENCHMARK.get(algo, DEFAULT_BENCHMARK)
        pxcol = BENCH_PRICE.get(bench)
        if bench not in sub_df or "avg_price" not in sub_df or pxcol not in sub_df:
            continue
        slip = sub_df[bench]
        zero_share = float(slip.fillna(0).abs().le(1e-9).mean())

        px, ex = sub_df[pxcol], sub_df["avg_price"]
        ok = np.isfinite(px) & np.isfinite(ex) & (px > 0)
        if ok.sum() == 0:
            continue
        # relative difference, so it works across price levels and currencies
        rel = (ex[ok] - px[ok]).abs() / px[ok]
        match_share = float((rel <= 1e-6).mean())

        if zero_share >= 0.95:
            verdict = ("executes AT the benchmark - score is zero by "
                       "construction" if match_share >= 0.95
                       else "benchmark field looks UNPOPULATED")
        else:
            verdict = "normal"
        rows.append({"algo": algo, "benchmark": bench, "n": len(sub_df),
                     "zero_slippage_%": 100 * zero_share,
                     "exec_price_equals_bench_%": 100 * match_share,
                     "verdict": verdict})
    return pd.DataFrame(rows)


def reference_check(df: pd.DataFrame) -> pd.DataFrame:
    """Recompute the published headline figures and diff them.

    Catches a mis-mapped column or an inverted sign before any slide is built.
    """
    if not REFERENCE:
        return pd.DataFrame()
    got = {
        "orders":         float(len(df)),
        "notional":       float(df["notional"].sum()),
        "shares":         float(df["shares_exec"].sum()) if "shares_exec" in df else np.nan,
        "impact_bps":     wmean(df["slip_primary"], df["notional"]),
        "spread_bps":     wmean(df["spread_bps"], df["notional"]) if "spread_bps" in df else np.nan,
        "adv_pct":        wmean(df["adv_pct"], df["notional"]) if "adv_pct" in df else np.nan,
        "dark_share_pct": 100 * df["dark_value"].sum() / df["notional"].sum(),
    }
    rows = []
    for k, expected in REFERENCE.items():
        actual = got.get(k, np.nan)
        tol = REFERENCE_TOL.get(k, 0.01)
        if not np.isfinite(actual):
            verdict = "n/a"
            diff = np.nan
        else:
            diff = actual - expected
            ok = (abs(diff) <= tol) if tol >= 1 or k in (
                "impact_bps", "spread_bps", "adv_pct", "dark_share_pct") \
                else (abs(diff) <= abs(expected) * tol)
            if k == "orders":
                ok = abs(diff) <= 0
            verdict = "ok" if ok else "MISMATCH"
        rows.append({"metric": k, "published": expected, "recomputed": actual,
                     "diff": diff, "check": verdict})
    t = pd.DataFrame(rows)

    log("")
    log("-" * 74)
    log("REFERENCE CHECK  (recomputed vs the published post-trade report)")
    log("-" * 74)
    for _, r in t.iterrows():
        pub = f"{r['published']:,.4g}" if np.isfinite(r["published"]) else "-"
        act = f"{r['recomputed']:,.4g}" if np.isfinite(r["recomputed"]) else "-"
        flag = "" if r["check"] == "ok" else f"   <== {r['check']}"
        log(f"  {r['metric']:<16} published {pub:>16}   recomputed {act:>16}{flag}")
    if (t["check"] == "MISMATCH").any():
        log("")
        log("  ** One or more headline figures do not reconcile. Check the")
        log("     COLUMNS mapping and POSITIVE_IS_SAVING before using it.")
        # Only cry "inverted" when the magnitudes actually line up. Opposite
        # signs alone are weak evidence: a book sitting near zero flips for
        # all sorts of reasons that have nothing to do with the convention.
        if np.isfinite(got["impact_bps"]) and REFERENCE.get("impact_bps"):
            ratio = (abs(got["impact_bps"]) / abs(REFERENCE["impact_bps"])
                     if REFERENCE["impact_bps"] else np.nan)
            if (got["impact_bps"] * REFERENCE["impact_bps"] < 0
                    and np.isfinite(ratio) and 0.5 <= ratio <= 2.0):
                log("     The sign is INVERTED -> set POSITIVE_IS_SAVING "
                    f"= {not POSITIVE_IS_SAVING}")
    return t


# ===========================================================================
# 4. WEIGHTED STATISTICS
# ===========================================================================

def winsorize(s: pd.Series, limits=WINSOR) -> pd.Series:
    if limits is None:
        return s
    s = pd.Series(s).astype(float)
    lo, hi = s.quantile(limits[0]), s.quantile(limits[1])
    return s.clip(lo, hi)


def wmean(values, weights) -> float:
    v = np.asarray(values, dtype=float)
    w = np.asarray(weights, dtype=float)
    m = np.isfinite(v) & np.isfinite(w) & (w > 0)
    if m.sum() == 0:
        return np.nan
    return float(np.sum(v[m] * w[m]) / np.sum(w[m]))


def boot_ci(values, weights, n=BOOTSTRAP_N, alpha=0.05):
    """Percentile bootstrap CI for the notional-weighted mean."""
    rng = np.random.default_rng(SEED)
    v = np.asarray(values, dtype=float)
    w = np.asarray(weights, dtype=float)
    m = np.isfinite(v) & np.isfinite(w) & (w > 0)
    v, w = v[m], w[m]
    if len(v) < 3:
        return (np.nan, np.nan)
    idx = rng.integers(0, len(v), size=(n, len(v)))
    draws = (v[idx] * w[idx]).sum(1) / w[idx].sum(1)
    return (float(np.quantile(draws, alpha / 2)),
            float(np.quantile(draws, 1 - alpha / 2)))


def beat_rate(values) -> float:
    """Share of orders that beat their benchmark (positive = saving)."""
    v = pd.Series(values).dropna()
    return float((v > 0).mean() * 100) if len(v) else np.nan


def _agg(sub: pd.DataFrame, value_col="slip_primary", ci=True) -> dict:
    s = winsorize(sub[value_col])
    lo, hi = boot_ci(s, sub["notional"]) if (ci and len(sub) >= MIN_N_FOR_CI) \
        else (np.nan, np.nan)
    out = {
        "n": len(sub),
        "notional_m": sub["notional"].sum() / 1e6,
        "bps": wmean(s, sub["notional"]),
        "lo": lo, "hi": hi,
        "beat_rate": beat_rate(sub[value_col]),
        "pnl_ccy": float((sub[value_col] * sub["notional"] / 1e4).sum()),
    }
    for c, key in [("adv_pct", "adv_pct"), ("spread_bps", "spread_bps"),
                   ("dark_pct", "dark_pct"), ("period_part", "period_part")]:
        if c in sub:
            out[key] = wmean(sub[c], sub["notional"])
    if "duration_min" in sub:
        out["median_dur_min"] = float(sub["duration_min"].median())
    return out


def group_table(df: pd.DataFrame, by: str, value_col="slip_primary",
                order=None, min_n=1) -> pd.DataFrame:
    """Weighted stats per group, plus each group's contribution to the total."""
    rows = []
    for key, sub in df.groupby(by, observed=True):
        if pd.isna(key) or len(sub) < min_n:
            continue
        r = {by: str(key)}
        r.update(_agg(sub, value_col))
        rows.append(r)
    if not rows:
        return pd.DataFrame()
    t = pd.DataFrame(rows).set_index(by)
    total_notional = df["notional"].sum()
    t["weight_pct"] = 100 * t["notional_m"] * 1e6 / total_notional
    # contribution to the programme-level bps: sums to the headline figure
    t["contrib_bps"] = t["bps"] * t["weight_pct"] / 100
    if order:
        t = t.reindex([o for o in order if o in t.index])
    else:
        t = t.sort_values("notional_m", ascending=False)
    return t.round(3)


def dark_markets_present(df: pd.DataFrame) -> list:
    """The dark-enabled markets that actually have executions in this file.

    Keeps the deck from naming a market the client did not trade this period.
    Returned in descending order of value traded, so the biggest is named first.
    """
    if not DARK_MARKETS or "market" not in df:
        return []
    keys = {_norm_key(m): m for m in DARK_MARKETS}
    e = df[df["market"].map(lambda v: _norm_key(v) in keys)]
    if e.empty:
        return []
    by_val = e.groupby("market")["notional"].sum().sort_values(ascending=False)
    # report the client's own spelling from the file
    return [str(m) for m in by_val.index]


def _dark_share_eligible(df: pd.DataFrame) -> float:
    """Dark as a share of value in the markets where dark is actually offered.

    The programme-wide share understates usage, because most markets have no
    dark to use in the first place.
    """
    if not DARK_MARKETS or "market" not in df:
        return np.nan
    keys = {_norm_key(m) for m in DARK_MARKETS}
    e = df[df["market"].map(lambda v: _norm_key(v) in keys)]
    if e.empty or e["notional"].sum() <= 0:
        return np.nan

    return float(100 * e["dark_value"].sum() / e["notional"].sum())


REPORT_COLS = ["name", "n", "shares", "notional", "weight_pct", "period_part",
               "adv_pct", "spread_bps", "bps", "contrib_bps"]


def table_from_report(rows, label="", total_bps=None) -> pd.DataFrame:
    """Turn a transcribed report table into the shape the charts expect.

    These figures are typed in by hand from the report, so each one is checked
    on load: the weights must sum to 100 and the contributions must sum to the
    published headline. A transcription slip shows up here rather than in a
    client deck.
    """
    if not rows:
        return pd.DataFrame()
    t = pd.DataFrame(list(rows), columns=REPORT_COLS)
    t["notional_m"] = t["notional"] / 1e6
    t["pnl_ccy"] = t["bps"] * t["notional"] / 1e4
    for c in ["lo", "hi", "beat_rate", "dark_pct", "median_dur_min"]:
        t[c] = np.nan
    t = t.set_index("name")
    t.index.name = label or "group"

    w = float(t["weight_pct"].sum())
    if abs(w - 100.0) > 1.5:
        log(f"  ** CHECK {label}: weights sum to {w:.1f}%, not 100. "
            "Re-read the transcribed table.")
    if total_bps is not None and np.isfinite(total_bps):
        c = float(t["contrib_bps"].sum())
        if abs(c - total_bps) > 0.6:
            log(f"  ** CHECK {label}: contributions sum to {c:+.2f} bps but the "
                f"published total is {total_bps:+.2f}. Re-read the table.")
    return t.drop(columns=["notional"]).round(3)


def headline_from_report(tot: dict) -> dict:
    """Headline KPIs straight off the report."""
    d = {
        "orders": int(tot["orders"]),
        "notional": float(tot["notional"]),
        "shares": float(tot.get("shares", np.nan)),
        "bps": float(tot["bps"]),
        "lo": np.nan, "hi": np.nan,
        "pnl_ccy": float(tot.get("pnl_ccy",
                                 tot["bps"] * tot["notional"] / 1e4)),
        "beat_rate": np.nan,
        "adv_pct": float(tot.get("adv_pct", np.nan)),
        "spread_bps": float(tot.get("spread_bps", np.nan)),
        "period_part": float(tot.get("period_part", np.nan)),
    }
    return d


def algo_table_from_report(rows, algos=None) -> pd.DataFrame:
    """The report's algo table, shaped like the computed one.

    Same columns as group_table(df, "algo") so every downstream chart, table
    and sentence works unchanged. `weight_pct` is renormalised when the review
    is scoped to a subset of algos, so the percentages still sum to 100.
    """
    if not rows:
        return pd.DataFrame()
    t = pd.DataFrame(rows, columns=["algo", "n", "notional", "weight_pct",
                                    "period_part", "adv_pct", "spread_bps",
                                    "benchmark", "bps", "contrib_bps"])
    if algos:
        t = t[t["algo"].isin(algos)]
        if t.empty:
            return t
    t["notional_m"] = t["notional"] / 1e6
    # blanks for the fields the report does not carry, so the shape matches
    for c in ["lo", "hi", "beat_rate", "dark_pct", "median_dur_min"]:
        t[c] = np.nan
    t["pnl_ccy"] = t["bps"] * t["notional"] / 1e4
    t = t.set_index("algo").drop(columns=["notional"])
    order = [a for a in ALGO_ORDER if a in t.index]
    order += [a for a in t.index if a not in order]
    return t.reindex(order).round(3)


def reconcile_algo(df: pd.DataFrame, rows) -> pd.DataFrame:
    """Recomputed vs published, per algo.

    The old reference check compared totals only, so two offsetting per-algo
    errors would pass. This one cannot.
    """
    if not rows or df.empty:
        return pd.DataFrame()
    pub = {r[0]: {"orders": r[1], "notional": r[2], "bps": r[8]} for r in rows}
    out = []
    for algo, sub_df in df.groupby("algo", observed=True):
        if algo not in pub:
            continue
        got_bps = wmean(winsorize(sub_df["slip_primary"]), sub_df["notional"])
        out.append({
            "algo": algo,
            "orders_pub": pub[algo]["orders"], "orders_got": len(sub_df),
            "notional_pub": pub[algo]["notional"],
            "notional_got": float(sub_df["notional"].sum()),
            "bps_pub": pub[algo]["bps"], "bps_got": got_bps,
        })
    t = pd.DataFrame(out)
    if t.empty:
        return t
    t["orders_diff"] = t["orders_got"] - t["orders_pub"]
    t["notional_diff_%"] = 100 * (t["notional_got"] / t["notional_pub"] - 1)
    t["bps_diff"] = t["bps_got"] - t["bps_pub"]
    return t.set_index("algo").round(2)


def headline(df: pd.DataFrame) -> dict:
    s = winsorize(df["slip_primary"])
    lo, hi = boot_ci(s, df["notional"])
    d = {
        "orders": len(df),
        "notional": float(df["notional"].sum()),
        "shares": float(df["shares_exec"].sum()) if "shares_exec" in df else np.nan,
        "bps": wmean(s, df["notional"]),
        "lo": lo, "hi": hi,
        "pnl_ccy": float(df["pnl_ccy"].sum()),
        "beat_rate": beat_rate(df["slip_primary"]),
        "dark_share": 100 * df["dark_value"].sum() / df["notional"].sum(),
        "dark_order_share": 100 * float(df["has_dark"].mean()),
        "markets": int(df["market"].nunique()) if "market" in df else np.nan,
        "dark_share_eligible": _dark_share_eligible(df),
        "symbols": int(df["symbol"].nunique()) if "symbol" in df else np.nan,
    }
    for c in ["adv_pct", "spread_bps", "period_part"]:
        if c in df:
            d[c] = wmean(df[c], df["notional"])
    return d


def benchmark_matrix(df: pd.DataFrame) -> pd.DataFrame:
    """Every algo against every benchmark.

    The mandated benchmark is only one view. An algo can look fine against the
    benchmark it is measured on and poor against arrival - that gap is the
    client's timing decision, and it belongs in the conversation.
    """
    benches = [b for b in MATRIX_BENCHMARKS if b in df]
    rows = []
    for algo, sub in df.groupby("algo", observed=True):
        r = {"algo": algo, "n": len(sub),
             "notional_m": sub["notional"].sum() / 1e6,
             "mandated": BENCHMARK_LABELS.get(
                 PRIMARY_BENCHMARK.get(algo, DEFAULT_BENCHMARK), "-")}
        for b in benches:
            r[BENCHMARK_LABELS[b]] = wmean(winsorize(sub[b]), sub["notional"])
        rows.append(r)
    if not rows:
        return pd.DataFrame()
    t = pd.DataFrame(rows).set_index("algo")
    tot = {"n": len(df), "notional_m": df["notional"].sum() / 1e6,
           "mandated": "Mixed"}
    for b in benches:
        tot[BENCHMARK_LABELS[b]] = wmean(winsorize(df[b]), df["notional"])
    t.loc["TOTAL"] = tot
    order = [a for a in ALGO_ORDER if a in t.index]
    order += [a for a in t.index if a not in order and a != "TOTAL"] + ["TOTAL"]
    return t.reindex(order).round(2)


def venue_mix(df: pd.DataFrame) -> pd.DataFrame:
    """Dark share of executed value per algo.

    The export carries dark explicitly; auction / visible post / visible take
    are not in the order file, so this table reports the split the order-level
    data can actually support (dark vs lit) and the deck sources the four-way
    venue segmentation from the report page where it is available.
    """
    rows = []
    for algo, sub in df.groupby("algo", observed=True):
        n_val = sub["notional"].sum()
        rows.append({"algo": algo, "n": len(sub), "notional_m": n_val / 1e6,
                     "Dark": 100 * sub["dark_value"].sum() / n_val,
                     "Lit": 100 * (1 - sub["dark_value"].sum() / n_val),
                     "orders_with_dark_pct": 100 * float(sub["has_dark"].mean())})
    if not rows:
        return pd.DataFrame()
    t = pd.DataFrame(rows).set_index("algo")
    t.loc["TOTAL"] = {
        "n": len(df), "notional_m": df["notional"].sum() / 1e6,
        "Dark": 100 * df["dark_value"].sum() / df["notional"].sum(),
        "Lit": 100 * (1 - df["dark_value"].sum() / df["notional"].sum()),
        "orders_with_dark_pct": 100 * float(df["has_dark"].mean())}
    order = [a for a in ALGO_ORDER if a in t.index]
    order += [a for a in t.index if a not in order and a != "TOTAL"] + ["TOTAL"]
    return t.reindex(order).round(2)


def pool_small_markets(df: pd.DataFrame) -> pd.DataFrame:
    """Fold markets below the notional threshold into 'Other'."""
    if "market" not in df:
        return df
    share = df.groupby("market")["notional"].sum() / df["notional"].sum()
    keep = set(share[share >= COUNTRY_MIN_SHARE].index)
    df = df.copy()
    df["market_grp"] = np.where(df["market"].isin(keep), df["market"], "Other")
    return df


# ===========================================================================
# 5. DARK ANALYSIS
# ===========================================================================

def dark_frame(df: pd.DataFrame, dark_algos: list) -> pd.DataFrame:
    """Dark-capable orders only, with the control buckets and the split label.

    Restricting to dark-capable algos is what stops the comparison from being a
    proxy for algo choice: IIS does no dark at all, so including it would make
    'no dark' mean 'IIS' rather than 'a VWAP order that got no dark fill'.
    """
    if not dark_algos:
        return df.iloc[0:0].copy()
    d = df[df["algo"].isin(dark_algos)].copy()
    if d.empty:
        return d

    if DARK_MARKETS and "market" in d:
        keys = {_norm_key(m) for m in DARK_MARKETS}
        keep = d["market"].map(lambda v: _norm_key(v) in keys)
        dropped = int((~keep).sum())
        d = d[keep].copy()
        log(f"  dark markets only ({', '.join(DARK_MARKETS)}): kept "
            f"{len(d):,} orders, dropped {dropped:,} in markets where dark is "
            "not offered")
        if d.empty:
            log("  ** no orders in the dark-enabled markets - check the "
                "spelling in DARK_MARKETS against the Market column")
            return d

    zero_share = float((d["dark_pct"] <= 0).mean())
    if zero_share >= DARK_ZERO_SPLIT_THRESHOLD:
        d["dark_group"] = np.where(d["dark_pct"] > 0, "Any dark", "No dark")
        d.attrs["groups"] = ["No dark", "Any dark"]
        d.attrs["split_note"] = (
            f"{'/'.join(dark_algos)} orders only. We compare the ones that got "
            f"no dark fill ({zero_share:.0%} of them) against the ones that "
            "got some.")
    else:
        med = float(d["dark_pct"].median())
        d["dark_group"] = np.where(d["dark_pct"] >= med, "High dark", "Low dark")
        d.attrs["groups"] = ["Low dark", "High dark"]
        d.attrs["split_note"] = (
            f"{'/'.join(dark_algos)} orders only, split at the middle of the "
            f"range ({med:.0f}% dark). Too few orders got no dark at all to "
            "compare against.")
    return d


def dark_sanity(d: pd.DataFrame, published_bps=None,
                value_col="slip_primary") -> bool:
    """Is the dark section fit to show?

    The dark subset is measured from the order file while the headline comes
    from the published report. If the two are far apart, the order-level
    slippage column is not the same quantity as the report's - wrong sign,
    wrong benchmark, or a handful of orders dominating - and no chart built on
    it should reach a client. Returns True when the section looks usable.
    """
    if d.empty:
        return False
    got = wmean(winsorize(d[value_col]), d["notional"])
    log("")
    log("  DARK SECTION SANITY")
    log(f"    dark-capable subset: {len(d):,} orders, "
        f"{CURRENCY} {d['notional'].sum()/1e6:,.0f}m, "
        f"{got:+.1f} bps (from the order file)")

    # the orders doing the most damage, by money not by bps
    d = d.assign(_m=d[value_col] * d["notional"] / 1e4)
    worst = d.nsmallest(5, "_m")
    tot = float(d["_m"].sum())
    if len(worst) and tot != 0:
        share = 100 * float(worst["_m"].sum()) / tot if tot else np.nan
        log(f"    top 5 losing orders account for {share:.0f}% of the subset's "
            "money impact:")
        for _, r in worst.iterrows():
            sym = r.get("symbol", "?")
            log(f"      {str(sym):<12} {r[value_col]:+9.1f} bps  "
                f"{CURRENCY} {r['_m']:>12,.0f}  "
                f"{r.get('market', '?')}")

    ok = True
    if published_bps is not None and np.isfinite(published_bps) and np.isfinite(got):
        if abs(got - published_bps) > 15:
            ok = False
            log("")
            log(f"    ** The subset is {got:+.1f} bps but the published book is "
                f"{published_bps:+.1f} bps.")
            log("    ** A gap that size is a measurement problem, not a "
                "finding. Do NOT present the dark charts until it is")
            log("       resolved. Check, in order: the sign of the order-level "
                "slippage column; whether it is the same benchmark")
            log("       the report used; and whether the outliers above are "
                "real orders or bad rows.")
    return ok


def dark_zero_vs_any(d: pd.DataFrame, value_col="slip_primary") -> pd.DataFrame:
    """The headline dark number: no dark fill vs any, within dark-capable algos."""
    if d.empty:
        return pd.DataFrame()
    rows = []
    for grp, sub in d.groupby("dark_group", observed=True):
        r = {"group": grp, "mean_dark_pct": wmean(sub["dark_pct"], sub["notional"])}
        r.update(_agg(sub, value_col))
        rows.append(r)
    t = pd.DataFrame(rows).set_index("group")
    return t.reindex([g for g in d.attrs.get("groups", []) if g in t.index]).round(2)


def dark_controlled(d: pd.DataFrame, control: str,
                    value_col="slip_primary") -> pd.DataFrame:
    """The same split, repeated inside each bucket of order difficulty.

    If the gap survives here it is not a composition artefact.
    """
    if d.empty or control not in d:
        return pd.DataFrame()
    rows = []
    for (ctrl, grp), sub in d.groupby([control, "dark_group"], observed=True):
        if pd.isna(ctrl):
            continue
        rows.append({"control": str(ctrl), "group": grp, "n": len(sub),
                     "notional_m": sub["notional"].sum() / 1e6,
                     "bps": wmean(winsorize(sub[value_col]), sub["notional"])})
    return pd.DataFrame(rows).round(2)


def dark_ladder(d: pd.DataFrame, value_col="slip_primary") -> pd.DataFrame:
    """Cost by dark-participation bucket."""
    if d.empty:
        return pd.DataFrame()
    rows = []
    for b, sub in d.groupby("dark_bucket", observed=True):
        r = {"bucket": str(b)}
        r.update(_agg(sub, value_col))
        rows.append(r)
    t = pd.DataFrame(rows).set_index("bucket")
    return t.reindex([b for b in DARK_LABELS if b in t.index]).round(2)


def dark_completion(d: pd.DataFrame) -> pd.DataFrame:
    """Does more dark cost completion? The standing objection to 'use more dark'."""
    if d.empty:
        return pd.DataFrame()
    rows = []
    for b, sub in d.groupby("dark_bucket", observed=True):
        r = {"bucket": str(b), "n": len(sub),
             "notional_m": sub["notional"].sum() / 1e6}
        if "period_part" in sub:
            r["participation_pct"] = wmean(sub["period_part"], sub["notional"])
        if "duration_min" in sub:
            r["median_dur_min"] = float(sub["duration_min"].median())
        if "order_qty" in sub and "shares_exec" in sub:
            q = sub["order_qty"].replace(0, np.nan)
            r["fill_rate_pct"] = wmean(100 * sub["shares_exec"] / q, sub["notional"])
        rows.append(r)
    t = pd.DataFrame(rows).set_index("bucket")
    return t.reindex([b for b in DARK_LABELS if b in t.index]).round(2)


def dark_by_market(d: pd.DataFrame, value_col="slip_primary") -> pd.DataFrame:
    """Availability and benefit, per market.

    This is what turns 'use more dark' into advice you can act on. Where only a
    few percent of orders can get a dark fill at all, there is nothing to
    increase, and the recommendation has to be something else.
    """
    if d.empty or "market" not in d:
        return pd.DataFrame()
    rows = []
    for mkt, sub in d.groupby("market", observed=True):
        if len(sub) < MIN_N_MARKET:
            continue
        nod = sub[~sub["has_dark"]]
        anyd = sub[sub["has_dark"]]
        rows.append({
            "market": str(mkt), "n": len(sub),
            "notional_m": sub["notional"].sum() / 1e6,
            "avail_orders_pct": 100 * float(sub["has_dark"].mean()),
            "dark_share_pct": 100 * sub["dark_value"].sum() / sub["notional"].sum(),
            "n_no_dark": len(nod), "n_any_dark": len(anyd),
            "bps_no_dark": wmean(winsorize(nod[value_col]), nod["notional"])
            if len(nod) else np.nan,
            "bps_any_dark": wmean(winsorize(anyd[value_col]), anyd["notional"])
            if len(anyd) else np.nan,
        })
    if not rows:
        return pd.DataFrame()
    t = pd.DataFrame(rows).set_index("market")
    t["benefit_bps"] = t["bps_no_dark"] - t["bps_any_dark"]
    return t.sort_values("notional_m", ascending=False).round(2)


def dark_regression(d: pd.DataFrame, value_col="slip_primary"):
    """cost ~ dark_pct + difficulty controls + market fixed effects.

    Weighted by executed notional (the client's real dollar exposure); the
    equal-weighted fit is reported alongside. If the two disagree materially the
    dark effect is concentrated in one size band and the deck should say so.

    Returns (text, coef_per_10pp, pvalue, coef_equal_weighted) or None.
    """
    try:
        import statsmodels.formula.api as smf
    except ImportError:
        log("  statsmodels not installed -> skipping the dark regression "
            "(pip install statsmodels)")
        return None
    if d.empty or len(d) < 60:
        log("  too few dark-capable orders for a regression")
        return None

    reg = d.copy()
    reg["y"] = winsorize(reg[value_col])
    reg["log_notional"] = np.log(reg["notional"].clip(lower=1))
    terms = ["dark_pct", "log_notional"]
    for c in ["spread_bps", "adv_pct", "period_part", "duration_min"]:
        if c in reg and reg[c].notna().sum() > 0.8 * len(reg):
            terms.append(c)

    # Fixed effects as explicit dummies: patsy would evaluate C(market) in this
    # namespace, and market levels are not valid Python identifiers.
    for fe_col in ["market", "algo", "side"]:
        if fe_col in reg and reg[fe_col].notna().sum() > 0.8 * len(reg) \
                and reg[fe_col].nunique() > 1:
            fe = pd.get_dummies(reg[fe_col].astype(str), prefix=fe_col[:3],
                                drop_first=True, dtype=float)
            clean, seen = [], {}
            for col in fe.columns:
                base = re.sub(r"\W+", "_", str(col)).strip("_") or "lvl"
                seen[base] = seen.get(base, 0) + 1
                clean.append(base if seen[base] == 1 else f"{base}_{seen[base]}")
            fe.columns = clean
            reg = pd.concat([reg, fe], axis=1)
            terms += clean

    reg = reg.dropna(subset=["y", "notional"] + terms)
    if len(reg) < 60:
        log("  too few complete rows for the dark regression")
        return None

    formula = "y ~ " + " + ".join(terms)
    try:
        wls = smf.wls(formula, data=reg, weights=reg["notional"]).fit(cov_type="HC3")
        ols = smf.ols(formula, data=reg).fit(cov_type="HC3")
    except Exception as e:
        log(f"  regression failed: {e}")
        return None

    coef = float(wls.params["dark_pct"]) * 10
    coef_ew = float(ols.params["dark_pct"]) * 10
    pval = float(wls.pvalues["dark_pct"])
    text = ("Dependent variable: execution cost in bps vs the mandated "
            "benchmark (positive = cost).\n"
            "Sample: dark-capable algos only. dark_pct is in percentage "
            "points, so the reported\neffect is per +10pp of dark "
            "participation.\n\n"
            "NOTIONAL-WEIGHTED (WLS, weight = executed notional) - headline\n"
            + wls.summary().as_text()
            + "\n\n\nEQUAL-WEIGHTED (OLS, one order = one observation) - "
              "cross-check\n" + ols.summary().as_text())
    log(f"  dark effect: {coef:+.2f} bps per +10pp dark (notional-weighted, "
        f"p={pval:.3g}); {coef_ew:+.2f} equal-weighted; n={len(reg):,}")
    if np.isfinite(coef) and np.isfinite(coef_ew) and coef * coef_ew < 0:
        log("  ** weighted and equal-weighted fits disagree in sign - the "
            "effect is concentrated in one size band. Say so in the deck.")
    return text, coef, pval, coef_ew


def worst_orders(df: pd.DataFrame, n=12) -> pd.DataFrame:
    """The orders that lost the most money against benchmark."""
    cols = [c for c in ["date", "symbol", "market", "algo", "side", "notional",
                        "adv_pct", "period_part", "duration_min", "spread_bps",
                        "dark_pct", "slip_primary", "pnl_ccy"] if c in df]
    # the biggest losses, i.e. the most negative
    return df.nsmallest(n, "pnl_ccy")[cols].round(2)


# ===========================================================================
# 6. CHARTS
# ===========================================================================
# Every exhibit is written as a standalone 200dpi PNG so it can be lifted
# straight into the corporate template.
#
# Design rules applied throughout: one axis (never a second y-scale); a
# diverging pair with a neutral midpoint for cost/savings; categorical hues
# assigned in fixed order and never cycled; colour follows the entity, so dark
# is indigo wherever it appears; a legend whenever two or more series are
# present, plus direct labels; recessive grid and axes; values in ink, never in
# the series colour.

import matplotlib
matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import Patch

# Set from the client config at startup (see CLIENTS). The order-level export
# carries dark only, so auction / visible post / visible take cannot be
# recomputed here and have to come from the report.
VENUE_SEGMENTS_REPORT = CLIENTS[DEFAULT_CLIENT]["venue_segments"]

# positive = saving, so savings run to the RIGHT and costs to the LEFT
COST_NOTE = "← costs you"
SAVE_NOTE = "saves you →"


def apply_style():
    matplotlib.rcParams.update({
        "figure.facecolor": SURFACE, "axes.facecolor": SURFACE,
        "savefig.facecolor": SURFACE, "savefig.dpi": DPI,
        "font.family": "sans-serif",
        "font.sans-serif": ["Segoe UI", "Calibri", "DejaVu Sans"],
        "font.size": 10.5,
        "text.color": INK, "axes.labelcolor": INK_SECOND,
        "xtick.color": INK_MUTED, "ytick.color": INK_MUTED,
        "axes.edgecolor": BASELINE, "axes.linewidth": 0.8,
        "axes.spines.top": False, "axes.spines.right": False,
        "grid.color": GRID, "grid.linewidth": 0.8, "grid.linestyle": "-",
        "legend.frameon": False, "figure.constrained_layout.use": True,
    })


def _finish(ax, title, subtitle=None, xlabel=None, ylabel=None, axis="x"):
    pad = 10 + 13 * (subtitle.count("\n") + 1) if subtitle else 8
    ax.set_title(title, loc="left", fontsize=14.5, color=INK, pad=pad)
    if subtitle:
        ax.text(0, 1.02, subtitle, transform=ax.transAxes, fontsize=9.5,
                color=INK_SECOND, va="bottom", linespacing=1.35)
    if xlabel:
        ax.set_xlabel(xlabel, fontsize=9.5)
    if ylabel:
        ax.set_ylabel(ylabel, fontsize=9.5)
    ax.grid(axis=axis, zorder=0)
    ax.set_axisbelow(True)


def _sign_colors(vals):
    """Positive is a saving (blue); negative is a cost (red)."""
    return [SAVE_COLOR if (np.isfinite(v) and v >= 0) else COST_COLOR
            for v in vals]


def _save(fig, name, outdir: Path) -> Path:
    outdir.mkdir(parents=True, exist_ok=True)
    path = outdir / f"{name}.png"
    fig.savefig(path, bbox_inches="tight")
    plt.close(fig)
    log(f"  saved {path.name}")
    return path


def _hbar_labels(ax, ys, vals, fmt="{:+.1f}", pad=None, floor=None):
    """Direct-label every bar, in ink, outside the bar end.

    `floor` is the smallest magnitude the format can express. Below it the
    label is written "under 1k" rather than rounded to "-0k", which reads as
    a real zero and makes a rounding artefact look like a finding.
    """
    if not len(vals):
        return
    span = max(abs(np.nanmax(vals)), abs(np.nanmin(vals)), 1e-9)
    pad = pad if pad is not None else span * 0.035
    for y, v in zip(ys, vals):
        if not np.isfinite(v):
            continue
        txt = ("under 1k" if floor is not None and abs(v) < floor
               else fmt.format(v))
        ax.text(v + (pad if v >= 0 else -pad), y, txt,
                va="center", ha="left" if v >= 0 else "right",
                fontsize=9.5, color=INK)


def _zero_line(ax, vertical=True):
    (ax.axvline if vertical else ax.axhline)(0, color=BASELINE, lw=1.1, zorder=2)


def _axis_note(ax, y=-0.13):
    ax.text(0, y, COST_NOTE, transform=ax.transAxes, fontsize=9,
            color=COST_COLOR, ha="left", va="top")
    ax.text(1, y, SAVE_NOTE, transform=ax.transAxes, fontsize=9,
            color=SAVE_COLOR, ha="right", va="top")


# --- Exhibit 1 — performance by algo ---------------------------------------
def chart_algo(t: pd.DataFrame, outdir: Path, hl: dict) -> Path:
    fig, ax = plt.subplots(figsize=(10.0, 4.6))
    t = t.iloc[::-1]
    ys = np.arange(len(t))
    vals = t["bps"].values
    quotable = (t["n"] >= MIN_N_QUOTABLE).values

    # Scale to the algos that carry real weight. A 30-order algo at +119 bps
    # would otherwise set the axis and flatten everything that matters.
    ref = vals[quotable] if quotable.any() else vals
    span = float(np.nanmax(np.abs(ref))) if len(ref) else 1.0
    span = max(span, 1.0)
    lim_hi, lim_lo = span * 1.55, -span * 0.55
    if np.nanmin(vals) < 0:
        lim_lo = min(lim_lo, float(np.nanmin(ref)) * 1.55 if len(ref) else lim_lo)

    drawn = np.clip(vals, lim_lo * 0.97, lim_hi * 0.97)
    clipped = ~np.isclose(drawn, vals)
    colors = _sign_colors(vals)
    bars = ax.barh(ys, drawn, height=0.55, color=colors, zorder=3)
    for b, is_clip, ok in zip(bars, clipped, quotable):
        if not ok:
            b.set_alpha(0.45)
            b.set_hatch("///")
            b.set_edgecolor(SURFACE)
        if is_clip:
            b.set_hatch("///")

    if t["lo"].notna().any():
        lo = t["bps"] - t["lo"]
        hi = t["hi"] - t["bps"]
        err = np.vstack([np.abs(lo.fillna(0)), np.abs(hi.fillna(0))])
        ax.errorbar(vals, ys, xerr=err, fmt="none", ecolor=INK_MUTED,
                    elinewidth=1.1, capsize=3, zorder=4)

    labels = []
    for i, r in t.iterrows():
        tail = " · too few to quote" if r["n"] < MIN_N_QUOTABLE else ""
        labels.append(f"{i}\n{int(r.n):,} orders · "
                      f"{f_money(r.notional_m * 1e6)}{tail}")
    ax.set_yticks(ys, labels, fontsize=10)

    # label every bar with its TRUE value, at the drawn position
    pad = span * 0.05
    for y, v, d in zip(ys, vals, drawn):
        if not np.isfinite(v):
            continue
        ax.text(d + (pad if d >= 0 else -pad), y, f"{v:+.1f}",
                va="center", ha="left" if d >= 0 else "right",
                fontsize=9.5, color=INK)
    _zero_line(ax)
    ax.set_xlim(lim_lo, lim_hi)
    if clipped.any():
        ax.text(0.0, -0.20, "Hatched bars run past the axis; the figure beside "
                "them is the true value.", transform=ax.transAxes, fontsize=8.5,
                color=INK_MUTED, ha="left", va="top")
    _finish(ax, "Performance by algorithm",
            xlabel="bps vs benchmark (positive = saving)")
    if ALGO_REPORT:
        ax.text(1.0, -0.28, "Source: Performance Summary By Algorithm.",
                transform=ax.transAxes, fontsize=8.5, color=INK_MUTED,
                ha="right", va="top")
    _axis_note(ax)
    return _save(fig, "01_algo_performance", outdir)


# --- Exhibit 2 — where the cost sits ---------------------------------------
def chart_attribution(t: pd.DataFrame, outdir: Path, by="market",
                      total_bps=None) -> Path:
    """Money made and lost, group by group.

    A market can be expensive per order and irrelevant to the total. Drawing
    the money rather than a contribution in basis points removes the one
    number a reader has to be taught, and the bars still add to the headline.
    """
    t = t[t["pnl_ccy"].notna()].copy()
    t = t.sort_values("pnl_ccy")
    fig, ax = plt.subplots(figsize=(10.0, 5.8))
    ys = np.arange(len(t))
    vals = (t["pnl_ccy"] / 1e3).values
    ax.barh(ys, vals, height=0.6, color=_sign_colors(vals), zorder=3)
    labels = [f"{i}\n{f_money(r.notional_m * 1e6)} traded · {r.bps:+.1f} bps"
              for i, r in t.iterrows()]
    ax.set_yticks(ys, labels, fontsize=9.5)
    _hbar_labels(ax, ys, vals, fmt="{:+,.0f}k", floor=0.5)
    _zero_line(ax)
    span = float(np.nanmax(np.abs(vals))) if len(vals) else 1.0
    ax.set_xlim(min(-span * 0.45, float(np.nanmin(vals)) * 1.5), span * 1.5)
    _finish(ax, f"Money made and lost, by {by}",
            xlabel=f"{CURRENCY} thousands better or worse than benchmark")
    _axis_note(ax, y=-0.10)
    return _save(fig, f"02_attribution_{by}", outdir)


# --- Exhibit 2b — executed notional by market -------------------------------
def chart_notional_by_country(df: pd.DataFrame, outdir: Path,
                              dark_markets: list = None,
                              report: pd.DataFrame = None) -> Path:
    """How much you actually trade, and where.

    Volumes come from the report where it is available. Bars are coloured by
    whether dark is offered in that market, which shows at a glance how much of
    the book the dark section can even apply to.
    """
    if report is not None and not report.empty:
        g = pd.DataFrame({"n": report["n"],
                          "notional": report["notional_m"] * 1e6})
        g = g.sort_values("notional", ascending=False)
    elif "market" in df:
        g = (df.groupby("market")
               .agg(n=("notional", "size"), notional=("notional", "sum"))
               .sort_values("notional", ascending=False))
    else:
        return None
    if g.empty:
        return None

    keys = {_norm_key(m) for m in (dark_markets or [])}
    is_dark = [_norm_key(m) in keys for m in g.index]
    total = g["notional"].sum()

    fig, ax = plt.subplots(figsize=(10.0, 2.2 + 0.34 * len(g)))
    ys = np.arange(len(g))[::-1]
    colors = [VENUE_COLORS["Dark"] if d else SAVE_COLOR for d in is_dark]
    ax.barh(ys, g["notional"] / 1e6, height=0.58, color=colors, zorder=3)
    for y, r in zip(ys, g.itertuples()):
        m = r.notional / 1e6
        # sub-10m markets would otherwise all read "0m"
        amt = f"{m:,.0f}m" if m >= 10 else (f"{m:,.1f}m" if m >= 0.1
                                            else f"{r.notional / 1e3:,.0f}k")
        pct = 100 * r.notional / total
        ax.text(m * 1.015 + float(g["notional"].max()) / 1e6 * 0.004, y,
                f"{amt}   ({int(r.n):,} orders, "
                + (f"{pct:.0f}%)" if pct >= 1 else f"{pct:.1f}%)"),
                va="center", fontsize=9.5, color=INK_SECOND)
    ax.set_yticks(ys, g.index, fontsize=10)
    ax.set_xlim(0, float(g["notional"].max()) / 1e6 * 1.42)
    if any(is_dark) and not all(is_dark):
        ax.legend(handles=[
            Patch(facecolor=VENUE_COLORS["Dark"], label="Dark available"),
            Patch(facecolor=SAVE_COLOR, label="No dark offered")],
            loc="lower right", fontsize=9.5)
    _finish(ax, "Executed notional by market",
            xlabel=f"executed notional ({CURRENCY}m)")
    if report is not None and not report.empty:
        ax.text(1.0, -0.16, "Source: country breakdown.",
                transform=ax.transAxes, fontsize=8.5, color=INK_MUTED,
                ha="right", va="top")
    return _save(fig, "02b_notional_by_market", outdir)


# --- Exhibit 3 — cost against order difficulty ------------------------------
def chart_adv(t: pd.DataFrame, outdir: Path) -> Path:
    fig, ax = plt.subplots(figsize=(10.0, 4.8))
    xs = np.arange(len(t))
    vals = t["bps"].values
    ax.bar(xs, vals, width=0.55, color=_sign_colors(vals), zorder=3)
    lo = np.abs((t["bps"] - t["lo"]).fillna(0))
    hi = np.abs((t["hi"] - t["bps"]).fillna(0))
    ax.errorbar(xs, vals, yerr=np.vstack([lo, hi]), fmt="none",
                ecolor=INK_MUTED, elinewidth=1.1, capsize=3, zorder=4)
    for x, v, r in zip(xs, vals, t.itertuples()):
        if not np.isfinite(v):
            continue
        off = max(np.nanmax(np.abs(vals)) * 0.05, 0.3)
        ax.text(x, v + (off if v >= 0 else -off), f"{v:+.1f}",
                ha="center", va="bottom" if v >= 0 else "top",
                fontsize=10, color=INK)
    ax.set_xticks(xs, [f"{i}\n{int(r.n):,} orders\n"
                       f"{f_money(r.notional_m * 1e6)}"
                       for i, r in t.iterrows()], fontsize=9.5)
    _zero_line(ax, vertical=False)
    _finish(ax, "Performance by order size",
            xlabel="order size, as a % of what the share trades in a normal day",
            ylabel="bps vs benchmark (positive = saving)", axis="y")
    return _save(fig, "03_adv_curve", outdir)


# --- Exhibit 4 — venue mix --------------------------------------------------
def chart_venue(outdir: Path, order_venue: pd.DataFrame = None) -> Path:
    """Stacked venue segmentation by algo.

    Four categorical hues in fixed order. The set was validated for CVD and
    normal-vision separation; the green sits just under the contrast floor, so
    every segment is also direct-labelled and the numbers appear in the table
    beside it.
    """
    if VENUE_SEGMENTS_REPORT:
        rows = [a for a in ["IIS", "PROG", "VWAP", "TOTAL"]
                if a in VENUE_SEGMENTS_REPORT]
        data = VENUE_SEGMENTS_REPORT
        src = ("Source: venue segment breakdown, % of executed "
               "value.")
    else:
        rows = list(order_venue.index)
        data = {r: {"Dark": order_venue.loc[r, "Dark"],
                    "Visible Take": order_venue.loc[r, "Lit"]} for r in rows}
        src = "Source: order file (dark vs lit only)."

    fig, ax = plt.subplots(figsize=(10.0, 4.4))
    ys = np.arange(len(rows))[::-1]
    left = np.zeros(len(rows))
    for seg in VENUE_ORDER:
        vals = np.array([data[r].get(seg, 0.0) for r in rows], dtype=float)
        # 2px surface gap between stacked segments keeps the boundary readable
        ax.barh(ys, vals, left=left, height=0.55, color=VENUE_COLORS[seg],
                edgecolor=SURFACE, linewidth=1.6, zorder=3, label=seg)
        for y, v, l in zip(ys, vals, left):
            if v >= 6:
                ax.text(l + v / 2, y, f"{v:.0f}%", ha="center", va="center",
                        fontsize=9.5, color="white", fontweight="bold")
        left += vals
    ax.set_yticks(ys, rows, fontsize=11)
    ax.set_xlim(0, 100)
    ax.set_xticks(range(0, 101, 20), [f"{v}%" for v in range(0, 101, 20)])
    for lbl in ax.get_yticklabels():
        if lbl.get_text() == "TOTAL":
            lbl.set_fontweight("bold")
    ax.legend(handles=[Patch(facecolor=VENUE_COLORS[s], label=s)
                       for s in VENUE_ORDER],
              loc="upper center", bbox_to_anchor=(0.5, -0.14), ncol=4,
              fontsize=10)
    _finish(ax, "Venue mix by algorithm")
    ax.text(1.0, -0.26, src, transform=ax.transAxes, fontsize=8.5,
            color=INK_MUTED, ha="right", va="top")
    return _save(fig, "04_venue_mix", outdir)


# --- Exhibit 3b — industry breakdown ----------------------------------------
def chart_industry(rows: list, outdir: Path, total_bps=None,
                   notional=None) -> Path:
    """Contribution by industry, from the report.

    Sector is not in the order-level export, so this is transcribed rather than
    recomputed. Sorted by contribution, which is what the total is made of - a
    sector can be dreadful per order and irrelevant to the number.
    """
    if not rows:
        return None
    t = pd.DataFrame(rows, columns=["industry", "issues", "weight_pct",
                                    "bps", "contrib_bps"])
    # the report carries no notional per sector, so rebuild it from the weight
    t["notional"] = t["weight_pct"] / 100 * (notional or np.nan)
    t["pnl_ccy"] = t["bps"] * t["notional"] / 1e4
    money = bool(t["pnl_ccy"].notna().all())
    t = t.sort_values("pnl_ccy" if money else "contrib_bps")

    fig, ax = plt.subplots(figsize=(10.0, 5.8))
    ys = np.arange(len(t))
    vals = (t["pnl_ccy"] / 1e3).values if money else t["contrib_bps"].values
    ax.barh(ys, vals, height=0.6, color=_sign_colors(vals), zorder=3)
    labels = [f"{r.industry}\n"
              + (f"{f_money(r.notional)} traded · " if money
                 else f"{r.weight_pct:.0f}% of value · ")
              + f"{r.bps:+.1f} bps"
              for r in t.itertuples()]
    ax.set_yticks(ys, labels, fontsize=9.5)
    _hbar_labels(ax, ys, vals, fmt="{:+,.0f}k" if money else "{:+.2f}",
                 floor=0.5 if money else None)
    _zero_line(ax)
    span = float(np.nanmax(np.abs(vals))) if len(vals) else 1.0
    ax.set_xlim(min(-span * 0.45, float(np.nanmin(vals)) * 1.5), span * 1.5)
    _finish(ax, "Money made and lost, by industry",
            xlabel=(f"{CURRENCY} thousands better or worse than benchmark"
                    if money else "effect on the total (bps)"))
    _axis_note(ax, y=-0.10)
    return _save(fig, "03b_industry", outdir)


# --- Exhibit 5 — the headline dark comparison -------------------------------
def chart_dark_zero_vs_any(t: pd.DataFrame, outdir: Path, note: str) -> Path:
    fig, ax = plt.subplots(figsize=(10.0, 4.2))
    ys = np.arange(len(t))[::-1]
    vals = t["bps"].values
    colors = [DARK_GROUP_COLORS.get(g, SAVE_COLOR) for g in t.index]
    ax.barh(ys, vals, height=0.5, color=colors, zorder=3)
    lo = np.abs((t["bps"] - t["lo"]).fillna(0))
    hi = np.abs((t["hi"] - t["bps"]).fillna(0))
    ax.errorbar(vals, ys, xerr=np.vstack([lo, hi]), fmt="none",
                ecolor=INK_MUTED, elinewidth=1.2, capsize=4, zorder=4)
    ax.set_yticks(ys, [f"{i}\n{int(r.n):,} orders · "
                       f"{CURRENCY} {r.notional_m:,.0f}m"
                       for i, r in t.iterrows()], fontsize=10)
    _hbar_labels(ax, ys, vals)
    _zero_line(ax)
    span = np.nanmax(np.abs(vals)) if len(vals) else 1
    ax.set_xlim(min(-span * 0.4, np.nanmin(vals) * 1.5), span * 1.6)

    if len(t) == 2:
        # say which side won in words: a signed "gap" reads ambiguously when
        # the metric is a cost and the sign convention is already inverted
        # positive = saving, so the better group is the HIGHER bar
        gap = float(t["bps"].iloc[1] - t["bps"].iloc[0])
        winner = t.index[1] if gap > 0 else t.index[0]
        ax.text(0.99, 0.08,
                f"{winner.lower()} better by {abs(gap):.1f} bps",
                transform=ax.transAxes, ha="right", fontsize=12,
                color=INK, fontweight="bold")
    _finish(ax, "Performance: no dark vs any dark",
            xlabel="bps vs benchmark (positive = saving)")
    _axis_note(ax, y=-0.16)
    return _save(fig, "05_dark_zero_vs_any", outdir)


# --- Exhibit 6 — the same split, inside each difficulty bucket ---------------
def chart_dark_controlled(d: pd.DataFrame, controls: list, outdir: Path,
                          groups: list) -> Path:
    """Small multiples: the dark split repeated within each control.

    If the gap only exists in the pooled number and vanishes here, it was a
    composition artefact and the deck must not claim a dark benefit.
    """
    panels = []
    for c in controls:
        t = dark_controlled(d, c)
        if not t.empty and t["control"].nunique() >= 2:
            panels.append((c, t))
    if not panels:
        return None
    panels = panels[:4]
    ncol = 2 if len(panels) > 1 else 1
    nrow = int(np.ceil(len(panels) / ncol))
    # constrained_layout is disabled here: the shared legend sits in reserved
    # space below the grid, and the automatic layout collides with it.
    fig, axes = plt.subplots(nrow, ncol, figsize=(10.0, 3.25 * nrow),
                             squeeze=False, constrained_layout=False)
    nice = {"adv_bucket": "order size (% ADV)", "spread_tercile": "spread",
            "size_tercile": "order notional", "market": "market",
            "dur_tercile": "duration"}
    order_map = {"spread_tercile": ["Low", "Mid", "High"],
                 "size_tercile": ["Low", "Mid", "High"],
                 "dur_tercile": ["Low", "Mid", "High"],
                 "adv_bucket": ADV_LABELS}
    for k, (c, t) in enumerate(panels):
        ax = axes[k // ncol][k % ncol]
        cats = [x for x in (order_map.get(c) or
                            sorted(t["control"].unique()))
                if x in set(t["control"])]
        if c == "market":
            # drop markets too thin to read; a three-order bar is noise
            n_by = t.groupby("control")["n"].sum()
            share = (t.groupby("control")["notional_m"].sum()
                      .loc[n_by[n_by >= MIN_N_MARKET].index]
                      .sort_values(ascending=False))
            cats = list(share.head(6).index)
            if not cats:
                continue
        xs = np.arange(len(cats))
        w = 0.36
        for gi, g in enumerate(groups):
            sub = t[t["group"] == g].set_index("control")
            vals = [sub["bps"].get(x, np.nan) for x in cats]
            ax.bar(xs + (gi - 0.5) * w, vals, width=w * 0.92,
                   color=DARK_GROUP_COLORS.get(g, SAVE_COLOR),
                   edgecolor=SURFACE, linewidth=1.2, zorder=3, label=g)
        ax.set_xticks(xs, cats, fontsize=8.8)
        ax.axhline(0, color=BASELINE, lw=1.0, zorder=2)
        ax.set_title(f"by {nice.get(c, c)}", loc="left", fontsize=11,
                     color=INK_SECOND, pad=6)
        ax.grid(axis="y", zorder=0)
        ax.set_axisbelow(True)
        if k % ncol == 0:
            ax.set_ylabel("bps (positive = saving)", fontsize=9)
    for k in range(len(panels), nrow * ncol):
        axes[k // ncol][k % ncol].axis("off")
    handles = [Patch(facecolor=DARK_GROUP_COLORS.get(g, SAVE_COLOR), label=g)
               for g in groups]
    # Descriptive, not a verdict: on some books the gap narrows or reverses
    # inside the controls, and the title must not assert what the panels then
    # contradict. The deck states the verdict, driven by the numbers.
    fig.suptitle("No dark vs any dark, by order difficulty", x=0.012,
                 y=0.985, ha="left", fontsize=14.5, color=INK)
    fig.tight_layout(rect=(0, 0.075, 1, 0.945))
    fig.legend(handles=handles, loc="lower center", ncol=2, fontsize=10.5,
               bbox_to_anchor=(0.5, 0.005))
    return _save(fig, "06_dark_controlled", outdir)


# --- Exhibit 7 — availability and benefit, per market ------------------------
def chart_dark_by_market(t: pd.DataFrame, outdir: Path, groups: list) -> Path:
    """Two panels, one axis each. Availability first, because where dark is
    unavailable the benefit column is advice you cannot act on."""
    t = t.head(9).iloc[::-1]
    fig, axes = plt.subplots(1, 2, figsize=(10.6, 4.9), sharey=True)
    ys = np.arange(len(t))

    # faint guides so a market's pair of panels reads as one row
    for ax_ in axes:
        for y in ys:
            ax_.axhline(y, color=GRID, lw=0.7, zorder=1)

    ax = axes[0]
    ax.barh(ys, t["avail_orders_pct"].values, height=0.55,
            color=VENUE_COLORS["Dark"], zorder=3)
    for y, v in zip(ys, t["avail_orders_pct"].values):
        if np.isfinite(v):
            ax.text(v + 1.5, y, f"{v:.0f}%", va="center", fontsize=9.5,
                    color=INK)
    ax.set_yticks(ys, t.index, fontsize=10)
    ax.set_xlim(0, max(float(t["avail_orders_pct"].max()) * 1.28, 10))
    _finish(ax, "Orders with a dark fill", xlabel="% of orders")

    ax = axes[1]
    w = 0.36
    for gi, g in enumerate(groups):
        col = "bps_no_dark" if g.startswith(("No", "Low")) else "bps_any_dark"
        ax.barh(ys + (0.5 - gi) * w, t[col].values, height=w * 0.92,
                color=DARK_GROUP_COLORS.get(g, SAVE_COLOR),
                edgecolor=SURFACE, linewidth=1.2, zorder=3, label=g)
    ax.axvline(0, color=BASELINE, lw=1.1, zorder=2)
    _finish(ax, "Performance with and without dark",
            xlabel="bps vs benchmark (positive = saving)")
    ax.legend(loc="upper center", bbox_to_anchor=(0.5, -0.16), ncol=2,
              fontsize=10)
    thin = t[(t["n_no_dark"] < MIN_N_MARKET) | (t["n_any_dark"] < MIN_N_MARKET)]
    if len(thin):
        fig.text(0.5, -0.06,
                 "Too few orders to rely on: "
                 + ", ".join(map(str, thin.index)) + ".",
                 ha="center", fontsize=8.5, color=INK_MUTED)
    fig.suptitle("Dark availability and performance by market", x=0.008,
                 ha="left", fontsize=14.5, color=INK)
    return _save(fig, "07_dark_by_market", outdir)


# --- Exhibit 8 — does more dark cost completion? ----------------------------
def chart_dark_completion(t: pd.DataFrame, outdir: Path) -> Path:
    """The standing objection to 'route more to dark'. Flat panels mean the
    recommendation is free; a slope means it has a price and a ceiling."""
    metrics = [(c, lbl) for c, lbl in [
        ("fill_rate_pct", "Fill rate (%)"),
        ("participation_pct", "Participation rate (%)"),
        ("median_dur_min", "Median duration (min)")] if c in t]
    if not metrics or t.empty:
        return None
    fig, axes = plt.subplots(1, len(metrics), figsize=(10.6, 3.6),
                             squeeze=False)
    xs = np.arange(len(t))
    ramp = ["#a9d5f5", "#7cc0ef", "#4aa8e9", "#2b98e4", "#1c6ba3"]
    for k, (c, lbl) in enumerate(metrics):
        ax = axes[0][k]
        vals = t[c].values
        ax.bar(xs, vals, width=0.6, zorder=3,
               color=[ramp[min(i, len(ramp) - 1)] for i in range(len(t))])
        for x, v in zip(xs, vals):
            if np.isfinite(v):
                ax.text(x, v, f"{v:,.0f}", ha="center", va="bottom",
                        fontsize=9, color=INK)
        ax.set_xticks(xs, [str(i).replace(" (none)", "") for i in t.index],
                      fontsize=8.5, rotation=30, ha="right")
        ax.set_title(lbl, loc="left", fontsize=11, color=INK_SECOND, pad=6)
        ax.grid(axis="y", zorder=0)
        ax.set_axisbelow(True)
    fig.suptitle("Fill rate, participation and duration by dark usage",
                 x=0.008, ha="left", fontsize=14.5, color=INK)
    return _save(fig, "08_dark_completion", outdir)


# ===========================================================================
# 7. DECK
# ===========================================================================

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:  # pragma: no cover
    sys.exit("Missing dependency: python-pptx\n  pip install python-pptx")


def _rgb(hex_str: str) -> "RGBColor":
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


T_INK      = _rgb(INK)
T_SECOND   = _rgb(INK_SECOND)
T_MUTED    = _rgb(INK_MUTED)
T_COST     = _rgb(COST_COLOR)
T_SAVE     = _rgb(SAVE_COLOR)
T_DARK     = _rgb(VENUE_COLORS["Dark"])
T_WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
T_BLACK    = RGBColor(0x00, 0x00, 0x00)
T_BAND_A   = _rgb("#e8f2fb")     # light azure tint
T_BAND_B   = _rgb("#f7fbfe")
T_HEAD     = _rgb("#2b2e31")     # charcoal from the template's panel
T_ACCENT   = _rgb("#2b98e4")     # corporate azure

FONT = "Calibri"
SLIDE_W, SLIDE_H = Inches(13.333), Inches(7.5)
MARGIN = Inches(0.62)
FOOTER_Y = Inches(6.95)


def new_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _run(p, text, size=11, bold=False, color=T_INK, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    r.font.name = FONT
    return r


def title(slide, text, accent=None):
    box = slide.shapes.add_textbox(MARGIN, Inches(0.30), Inches(12.1), Inches(0.62))
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    _run(p, text, size=25, bold=True, color=T_HEAD)
    if accent:
        _run(p, "     " + accent, size=25, bold=False, color=T_MUTED)
    ln = slide.shapes.add_connector(1, MARGIN, Inches(0.99),
                                    MARGIN + Inches(7.9), Inches(0.99))
    ln.line.color.rgb = T_ACCENT
    ln.line.width = Pt(2.25)
    return box


def strapline(slide, text, top=Inches(1.06), color=None, size=12.5):
    box = slide.shapes.add_textbox(MARGIN, top, Inches(12.1), Inches(0.40))
    tf = box.text_frame
    tf.word_wrap = True
    _run(tf.paragraphs[0], text, size=size, bold=True, color=color or T_INK)
    return box


def bullets(slide, items, left, top, width, height, size=11.5, gap=8,
            marker="▪  "):
    """**bold** segments render bold; a leading '~' makes the line muted."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        muted = item.startswith("~")
        item = item[1:] if muted else item
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = Pt(gap)
        p.line_spacing = 1.14
        base = T_MUTED if muted else T_INK
        if marker:
            _run(p, marker, size=size, color=base)
        for j, seg in enumerate(item.split("**")):
            if seg:
                _run(p, seg, size=size, bold=(j % 2 == 1), color=base)
    return box


def note(slide, text, top=Inches(6.52), size=8.5, width=Inches(12.1)):
    box = slide.shapes.add_textbox(MARGIN, top, width, Inches(0.42))
    tf = box.text_frame
    tf.word_wrap = True
    _run(tf.paragraphs[0], text, size=size, color=T_MUTED)
    return box


_PAGE = [0]
IS_SAMPLE = [False]


def sample_stamp(slide):
    """Red banner on every slide when the data is synthetic."""
    if not IS_SAMPLE[0]:
        return
    bar = slide.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(0.26))
    bar.fill.solid()
    bar.fill.fore_color.rgb = T_COST
    bar.line.fill.background()
    bar.shadow.inherit = False
    tf = bar.text_frame
    tf.margin_top = Inches(0.01)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    _run(tf.paragraphs[0],
         "SAMPLE DATA - INVENTED NUMBERS, NOT REAL CLIENT DATA. "
         "DO NOT PRESENT.",
         size=10, bold=True, color=T_WHITE)


def footer(slide, page=None):
    sample_stamp(slide)
    if page is None:
        _PAGE[0] += 1
        page = _PAGE[0]
    box = slide.shapes.add_textbox(Inches(0.62), FOOTER_Y, Inches(8.0), Inches(0.3))
    _run(box.text_frame.paragraphs[0],
         f"{CLIENT_NAME} — Transaction Cost Analysis · {PERIOD_LABEL}",
         size=8.5, color=T_MUTED)
    pn = slide.shapes.add_textbox(Inches(12.5), FOOTER_Y, Inches(0.5), Inches(0.3))
    tf = pn.text_frame
    tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    _run(tf.paragraphs[0], str(page), size=8.5, color=T_MUTED)


def _cover_line(hl: dict) -> str:
    """The headline sentence for the cover notes, built in one piece."""
    verb = "beat the reference price by" if hl["bps"] >= 0 else            "came in below the reference price by"
    kept = "kept" if hl["bps"] >= 0 else "given away"
    return (f"{f_bps(hl['bps'])} bps means that, on average, you {verb} "
            f"{abs(hl['bps']) / 100:.3f}%. On {f_money(hl['notional'], 1)} "
            f"traded, that works out at about "
            f"{f_money(abs(hl['pnl_ccy']))} {kept}.")


def notes(slide, text: str) -> None:
    """Attach speaker notes.

    Written plainly on purpose: these are read by whoever presents the pack,
    who may not be the person who built it, and every term the slide uses gets
    explained once here.
    """
    tf = slide.notes_slide.notes_text_frame
    tf.text = text.strip()


def picture(slide, path, left, top, width=None, height=None):
    if path is None or not Path(path).exists():
        return None
    return slide.shapes.add_picture(str(path), left, top, width=width,
                                    height=height)


def table(slide, data, left, top, width, col_w=None, row_h=Inches(0.28),
          size=9.5, header_size=9.5, colored=()):
    """data[0] is the header. colored: {(row, col): 'cost'|'save'}."""
    rows, cols = len(data), len(data[0])
    shape = slide.shapes.add_table(rows, cols, left, top, width, row_h * rows)
    tbl = shape.table
    tbl.first_row = False
    if col_w:
        for i, w in enumerate(col_w):
            tbl.columns[i].width = w
    for r_i, row in enumerate(data):
        tbl.rows[r_i].height = row_h
        for c_i, val in enumerate(row):
            cell = tbl.cell(r_i, c_i)
            cell.text = ""
            cell.margin_left = Inches(0.05)
            cell.margin_right = Inches(0.05)
            cell.margin_top = Inches(0.01)
            cell.margin_bottom = Inches(0.01)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.fill.solid()
            if r_i == 0:
                cell.fill.fore_color.rgb = T_HEAD
            else:
                cell.fill.fore_color.rgb = T_BAND_A if r_i % 2 else T_BAND_B
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c_i == 0 else PP_ALIGN.RIGHT
            is_total = str(row[0]).strip().upper() in ("TOTAL", "PROGRAMME")
            colr = T_WHITE if r_i == 0 else T_INK
            tag = colored.get((r_i, c_i)) if isinstance(colored, dict) else None
            if tag == "cost":
                colr = T_COST
            elif tag == "save":
                colr = T_SAVE
            _run(p, str(val), size=header_size if r_i == 0 else size,
                 bold=(r_i == 0 or is_total), color=colr)
    return shape


def kpi(slide, left, top, w, h, big, caption, fill=None, big_size=25):
    box = slide.shapes.add_shape(1, left, top, w, h)
    box.fill.solid()
    box.fill.fore_color.rgb = fill or T_HEAD
    box.line.fill.background()
    box.shadow.inherit = False
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_top = Inches(0.09)
    tf.margin_left = Inches(0.06)
    tf.margin_right = Inches(0.06)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    _run(p, big, size=big_size, bold=True, color=T_WHITE)
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    _run(p2, caption, size=9.5, color=T_WHITE)
    return box


# --- formatting helpers -----------------------------------------------------
def f_bps(v, dp=1):
    return "—" if not np.isfinite(v) else f"{v:+.{dp}f}"


def f_money(v, dp=0):
    if not np.isfinite(v):
        return "—"
    a = abs(v)
    if a >= 1e9:
        return f"{CURRENCY} {v / 1e9:,.2f}bn"
    if a >= 1e6:
        return f"{CURRENCY} {v / 1e6:,.{dp}f}m"
    if a >= 1e3:
        return f"{CURRENCY} {v / 1e3:,.0f}k"
    return f"{CURRENCY} {v:,.0f}"


def f_pct(v, dp=1):
    return "—" if not np.isfinite(v) else f"{v:.{dp}f}%"


# ===========================================================================
# 8. NARRATIVE
# ===========================================================================
# Sign convention throughout: POSITIVE = saving, NEGATIVE = cost.
# Nothing here is a fixed claim. Every sentence is assembled from the computed
# numbers and dropped when the evidence for it is not in the data.

def build_narrative(ctx: dict) -> dict:
    hl = ctx["headline"]
    algo = ctx["algo_table"]
    mkt = ctx["market_table"]
    adv = ctx["adv_table"]
    dz = ctx["dark_zero_vs_any"]
    reg = ctx["regression"]
    dmkt = ctx["dark_by_market"]
    venue = ctx["venue_table"]

    dmk = ctx.get("dark_markets") or []
    dmk_txt = ", ".join(dmk) if dmk else "the dark-enabled markets"

    n = {"findings": [], "recs": [], "dark_verdict": [], "caveats": [],
         "rec_tags": {}}

    def add_rec(tag, text):
        n["recs"].append(text)
        n["rec_tags"][tag] = text

    ahead = np.isfinite(hl["bps"]) and hl["bps"] >= 0

    # --- the headline -----------------------------------------------------
    n["findings"].append(
        (f"Overall you are **{f_bps(hl['bps'])} bps ahead** of your "
         f"benchmarks, worth about {f_money(abs(hl['pnl_ccy']))}."
         if ahead else
         f"Overall you are **{f_bps(hl['bps'])} bps behind** your benchmarks, "
         f"which cost about {f_money(abs(hl['pnl_ccy']))}.")
        + (f" {hl['beat_rate']:.0f}% of orders beat their benchmark."
           if np.isfinite(hl.get("beat_rate", np.nan)) else ""))

    # --- which algo drives it ---------------------------------------------
    if not algo.empty and "contrib_bps" in algo:
        big = algo[algo["n"] >= MIN_N_QUOTABLE]
        if not big.empty:
            top = big["contrib_bps"].idxmax() if ahead else big["contrib_bps"].idxmin()
            r = big.loc[top]
            share = 100 * r["contrib_bps"] / hl["bps"] if hl["bps"] else np.nan
            n["findings"].append(
                f"**{top}** does {r['weight_pct']:.0f}% of your volume at "
                f"{f_bps(r['bps'])} bps, and that is where "
                + (f"**{abs(share):.0f}%** of the result comes from."
                   if np.isfinite(share) else "most of the result comes from."))
        drag = algo[(algo["bps"] < 0) & (algo["n"] >= MIN_N_QUOTABLE)]
        if len(drag):
            d0 = drag["contrib_bps"].idxmin()
            n["findings"].append(
                f"**{d0}** is the one losing money: {f_bps(algo.loc[d0, 'bps'])} "
                f"bps on {algo.loc[d0, 'weight_pct']:.0f}% of your volume.")

    small = algo[algo["n"] < MIN_N_QUOTABLE] if not algo.empty else algo
    if len(small):
        nm = ", ".join(f"{i} (only {int(r.n)} orders)" for i, r in small.iterrows())
        n["caveats"].append(
            f"{nm} — far too few to draw anything from. We have left them in "
            "so nothing is hidden, but no advice here rests on them.")

    # --- markets ----------------------------------------------------------
    if not mkt.empty and "contrib_bps" in mkt:
        good = mkt[mkt["contrib_bps"] > 0].sort_values("contrib_bps",
                                                       ascending=False)
        bad = mkt[mkt["contrib_bps"] < 0].sort_values("contrib_bps")
        if len(good):
            r = good.iloc[0]
            n["findings"].append(
                f"**{good.index[0]}** is doing the heavy lifting: "
                f"{r['weight_pct']:.0f}% of your value at {f_bps(r['bps'])} "
                f"bps, adding **{r['contrib_bps']:+.1f} bps** to the total.")
        if len(bad):
            r = bad.iloc[0]
            # what it would be worth to bring this market back to flat
            prize = abs(r["bps"]) * (r["notional_m"] * 1e6) / 1e4
            # only call something "the biggest drag" if it actually moves the
            # total; a market worth -0.2 bps does not
            threshold = max(0.3, 0.10 * abs(hl["bps"])) if np.isfinite(hl["bps"]) \
                else 0.3
            material = abs(r["contrib_bps"]) >= threshold
            n["findings"].append(
                f"**{bad.index[0]}** is your weakest market: "
                f"{f_bps(r['bps'])} bps on {r['weight_pct']:.0f}% of your "
                f"value, pulling **{r['contrib_bps']:+.1f} bps** off the total."
                + ("" if material else
                   " It is small enough that it barely moves the headline."))
            if material:
                add_rec("market",
                        f"**Look at {bad.index[0]} first.** It is the biggest "
                        f"single drag on your number at {f_bps(r['bps'])} bps. "
                        f"Bringing it back to flat is worth about "
                        f"{f_money(prize)}.")
            else:
                spread_thin = bad["contrib_bps"].sum()
                add_rec("market",
                        f"**No single market is hurting you much.** The worst "
                        f"is {bad.index[0]} at {f_bps(r['bps'])} bps, but it is "
                        f"only {r['weight_pct']:.0f}% of your value. All the "
                        f"losing markets together come to "
                        f"{spread_thin:+.1f} bps, so the gains here are in "
                        "housekeeping rather than one big fix.")

    # --- order size -------------------------------------------------------
    if not adv.empty and len(adv) >= 2:
        big = adv[adv["n"] >= 20] if (adv["n"] >= 20).any() else adv
        if len(big) >= 2:
            best = big["bps"].idxmax()
            worst = big["bps"].idxmin()
            first = big.index[0]
            if worst == first:
                n["findings"].append(
                    f"**Your small orders do worst.** The **{first}** bucket "
                    f"returns {f_bps(big.loc[first, 'bps'])} bps against "
                    f"{f_bps(big.loc[best, 'bps'])} for **{best}**, and it is "
                    f"{big.loc[first, 'weight_pct']:.0f}% of your value. "
                    "Small orders are being pushed too fast.")
                add_rec("small_orders",
                        f"**Slow down the small orders.** The {first} bucket "
                        "does worse than genuinely hard trades. Give them more "
                        "time and let them post instead of crossing the "
                        "spread.")
            else:
                n["findings"].append(
                    f"Best on **{best}** ({f_bps(big.loc[best, 'bps'])} bps), "
                    f"worst on **{worst}** ({f_bps(big.loc[worst, 'bps'])} "
                    "bps).")

    # --- venue ------------------------------------------------------------
    seg = VENUE_SEGMENTS_REPORT or {}
    algo_rows = {a: v for a, v in seg.items() if a != "TOTAL"}

    if not DARK_STORY and algo_rows:
        # No dark to talk about - the real difference between these algos is
        # where they put the volume: the auction, the passive queue, or across
        # the spread.
        auc = max(algo_rows, key=lambda a: algo_rows[a].get("Auction", 0))
        tak = max(algo_rows, key=lambda a: algo_rows[a].get("Visible Take", 0))
        pos = max(algo_rows, key=lambda a: algo_rows[a].get("Visible Post", 0))
        if algo_rows[auc].get("Auction", 0) >= 50:
            n["findings"].append(
                f"**{auc}** puts **{algo_rows[auc]['Auction']:.0f}%** of its "
                "volume into the auction. Its result is really a question of "
                "where the close printed, not how the order was worked.")
        if algo_rows[pos].get("Visible Post", 0) >= 40:
            n["findings"].append(
                f"**{pos}** posts **{algo_rows[pos]['Visible Post']:.0f}%** of "
                "its volume and earns the spread rather than paying it.")
        take_pct = algo_rows[tak].get("Visible Take", 0)
        if take_pct >= 20:
            half = (hl.get("spread_bps", np.nan) or np.nan) / 2
            cost_line = (f" At a {hl['spread_bps']:.1f} bps spread that is "
                         f"roughly {half:.1f} bps given away on every one."
                         if np.isfinite(half) else "")
            n["findings"].append(
                f"**{tak}** crosses the spread on **{take_pct:.0f}%** of what "
                f"it trades.{cost_line}")
            add_rec("venue",
                    f"**Get {tak} posting more and taking less.** It pays the "
                    f"spread on {take_pct:.0f}% of its volume. Even a small "
                    "shift from taking to posting drops straight through to "
                    "the number.")
        if algo_rows[auc].get("Auction", 0) >= 50:
            add_rec("algo_mix",
                    f"**Check that {auc} is the right tool.** An auction-heavy "
                    "strategy is the right answer when you genuinely want the "
                    "close. If you do not, working the order through the day "
                    f"the way {pos} does gives you far more control.")

    no_dark_algos = [a for a, v in algo_rows.items() if v.get("Dark", 0) < 1.0]
    if DARK_STORY and no_dark_algos:
        takers = [a for a in no_dark_algos if seg[a].get("Visible Take", 0) >= 50]
        nm = ", ".join(no_dark_algos)
        verb = "never trade" if len(no_dark_algos) > 1 else "never trades"
        line = f"**{nm}** {verb} in dark"
        if takers:
            t0 = takers[0]
            line += (f", and {t0} pays the spread on "
                     f"**{seg[t0]['Visible Take']:.0f}%** of what it trades")
        n["findings"].append(line + ".")
        add_rec("dark_on",
                f"**Switch dark on for {nm}.** Today these orders either cross "
                "the spread or wait for the auction — they never try the "
                f"midpoint. In {dmk_txt} that is free optionality you are "
                "not using.")

    if np.isfinite(hl.get("dark_share_eligible", np.nan)) and dmk:
        n["findings"].append(
            f"Dark is only available in **{dmk_txt}**. Across those markets "
            f"you already do **{hl['dark_share_eligible']:.0f}%** of your "
            f"value in dark, against {f_pct(hl['dark_share'])} of the whole "
            "book.")

    # --- did dark pay? ----------------------------------------------------
    if not dz.empty and len(dz) == 2:
        g0, g1 = dz.index[0], dz.index[1]
        b0, b1 = dz["bps"].iloc[0], dz["bps"].iloc[1]
        gap = b1 - b0                      # positive => the dark group did better
        better = g1 if gap > 0 else g0
        overlap = (np.isfinite(dz["hi"].iloc[1]) and np.isfinite(dz["lo"].iloc[0])
                   and not (dz["hi"].iloc[1] < dz["lo"].iloc[0]
                            or dz["lo"].iloc[1] > dz["hi"].iloc[0]))
        n["dark_verdict"].append(
            f"Orders with **{g1.lower()}** returned {f_bps(b1)} bps. Orders "
            f"with **{g0.lower()}** returned {f_bps(b0)} bps. That is "
            f"**{abs(gap):.1f} bps** better for **{better.lower()}**.")
        if overlap:
            n["dark_verdict"].append(
                "~On this many orders the gap could still be luck. Treat it as "
                "a steer, not a proven number.")
        n["dark_verdict"].append(
            f"This is like for like — same algos, same markets, same desk. "
            f"{int(dz['n'].iloc[0]):,} orders got no dark and "
            f"{int(dz['n'].iloc[1]):,} got some.")

    if reg:
        _, coef, pval, coef_ew = reg
        sig = "a solid result" if pval < 0.05 else "could still be luck"
        n["dark_verdict"].append(
            f"Once we adjust for order size, spread, speed, duration and "
            f"market, every extra **10% of dark** is worth **{coef:+.2f} bps** "
            f"({sig}, p={pval:.3f}).")
        if np.isfinite(coef_ew) and coef * coef_ew < 0:
            n["dark_verdict"].append(
                "~The answer flips depending on whether we weight by order "
                "size, so it holds for some order sizes and not others.")

    # --- already saturated? -------------------------------------------------
    if not dmkt.empty and "avail_orders_pct" in dmkt:
        avail = wmean(dmkt["avail_orders_pct"], dmkt["notional_m"])
        if np.isfinite(avail) and avail >= 70:
            n["findings"].append(
                f"**You already get dark on {avail:.0f}% of eligible orders.** "
                "In the markets where we offer it, almost every order is "
                "already interacting with the midpoint — there is very little "
                "headroom left to add.")
            n["dark_verdict"].append(
                f"~With {avail:.0f}% of eligible orders already getting a dark "
                "fill there is no clean no-dark group to compare against, so "
                "this is a more-versus-less comparison rather than a "
                "with-versus-without one.")
            # a "use more dark" recommendation is not available here
            n["rec_tags"].pop("dark_market", None)
            n["recs"] = [r for r in n["recs"] if "Push dark hardest" not in r]

    # --- where dark actually pays ------------------------------------------
    if not dmkt.empty:
        good = dmkt[(dmkt["benefit_bps"] > 0)
                    & (dmkt["n_any_dark"] >= MIN_N_MARKET)].sort_values(
            "benefit_bps", ascending=False)
        if len(good):
            nm = ", ".join(good.head(2).index)
            add_rec("dark_market",
                    f"**Push dark hardest in {nm}.** That is where orders "
                    "getting dark clearly beat the ones that did not.")
        low = dmkt[dmkt["avail_orders_pct"] < 25]
        if len(low):
            n["caveats"].append(
                "Even where dark is offered, only "
                + ", ".join(f"{i} ({r.avail_orders_pct:.0f}%)"
                            for i, r in low.head(3).iterrows())
                + " of orders actually get a fill — so there is a ceiling on "
                  "how much more you can route.")

    if dmk:
        n["caveats"].append(
            f"We only offer dark in {dmk_txt}, so the whole dark comparison is "
            "limited to those markets. Everywhere else the thing to change is "
            "timing and order size, not venue.")

    n["caveats"].append(
        "One thing to keep in mind: you only get a dark fill when someone is "
        "there on the other side, and that may mean the order was easier "
        "anyway. So the honest reading is that orders getting dark did better, "
        "and this holds up when we adjust for how hard they were — not that "
        "dark on its own caused it.")
    return n


# ===========================================================================
# 9. SLIDES
# ===========================================================================

def slide_cover(prs, ctx):
    s = new_slide(prs)
    hl = ctx["headline"]
    band = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(2.55))
    band.fill.solid()
    band.fill.fore_color.rgb = T_HEAD
    band.line.fill.background()
    band.shadow.inherit = False

    box = s.shapes.add_textbox(MARGIN, Inches(0.72), Inches(11.0), Inches(1.6))
    tf = box.text_frame
    tf.word_wrap = True
    _run(tf.paragraphs[0], "Transaction Cost Analysis", size=38, bold=True,
         color=T_WHITE)
    p = tf.add_paragraph()
    _run(p, f"{CLIENT_NAME}  ·  {PERIOD_LABEL}", size=17, color=T_WHITE)
    p2 = tf.add_paragraph()
    _run(p2, BROKER, size=11, color=T_WHITE)

    cards = [
        (f"{hl['orders']:,}", "orders"),
        (f_money(hl["notional"], 1), "executed"),
        (f_bps(hl["bps"]) + " bps", "vs benchmark"),
        (f_money(abs(hl["pnl_ccy"])),
         "saved" if hl["pnl_ccy"] >= 0 else "cost"),
        (f_pct(hl["dark_share"]), "executed in dark"),
    ]
    w, gap = Inches(2.34), Inches(0.10)
    left = MARGIN
    for i, (big, cap) in enumerate(cards):
        good = np.isfinite(hl["bps"]) and hl["bps"] >= 0
        fill = T_DARK if i == 4 else (
            (T_SAVE if good else T_COST) if i in (2, 3) else T_HEAD)
        size = 22 if len(big) <= 10 else 17
        kpi(s, left, Inches(2.95), w, Inches(1.15), big, cap, fill=fill,
            big_size=size)
        left += w + gap

    bullets(s, [
        f"We looked at **{hl['orders']:,} orders** across "
        f"**{hl['markets']} markets**"
        + (f" and {hl['symbols']:,} names" if np.isfinite(hl.get('symbols', np.nan)) else "")
        + f", worth {f_money(hl['notional'], 1)}.",
        (f"Overall you **beat** your benchmarks by **{f_bps(hl['bps'])} bps**, "
         f"worth about {f_money(abs(hl['pnl_ccy']))}."
         if hl["bps"] >= 0 else
         f"Overall you came in **{f_bps(hl['bps'])} bps** behind your "
         f"benchmarks, which cost about {f_money(abs(hl['pnl_ccy']))}."),
        f"**{f_pct(hl['dark_share'])}** of that traded in dark. That is what "
        "this review focuses on.",
    ], MARGIN, Inches(4.45), Inches(11.9), Inches(1.7), size=13)

    note(s, "A plus means you beat the benchmark and saved money. A minus "
            "means it cost you. Averages are weighted by how much was traded, "
            "so big orders count for more.")
    footer(s)


def slide_summary(prs, ctx, nar):
    s = new_slide(prs)
    title(s, "Executive Summary")
    hl = ctx["headline"]
    strapline(s, (f"You are {f_bps(hl['bps'])} bps ahead of benchmark "
                  f"({f_money(abs(hl['pnl_ccy']))}). "
                  if hl["bps"] >= 0 else
                  f"You are {f_bps(hl['bps'])} bps behind benchmark "
                  f"({f_money(abs(hl['pnl_ccy']))}). ")
                 + "Here is what is driving it, and where you are leaking.")

    box = s.shapes.add_textbox(MARGIN, Inches(1.60), Inches(6.25), Inches(0.35))
    _run(box.text_frame.paragraphs[0], "WHAT WE FOUND", size=11, bold=True,
         color=T_HEAD)
    bullets(s, nar["findings"][:5], MARGIN, Inches(1.98), Inches(6.25),
            Inches(3.4), size=11.5)

    x2 = Inches(7.15)
    box = s.shapes.add_textbox(x2, Inches(1.60), Inches(5.55), Inches(0.35))
    _run(box.text_frame.paragraphs[0], "WHAT WE RECOMMEND", size=11, bold=True,
         color=T_HEAD)
    bullets(s, nar["recs"][:4], x2, Inches(1.98), Inches(5.55), Inches(3.4),
            size=11.5)

    if nar["caveats"]:
        box = s.shapes.add_textbox(MARGIN, Inches(5.45), Inches(12.1), Inches(0.3))
        _run(box.text_frame.paragraphs[0], "WORTH KNOWING", size=10, bold=True,
             color=T_MUTED)
        bullets(s, ["~" + c for c in nar["caveats"][:2]], MARGIN, Inches(5.75),
                Inches(12.1), Inches(0.9), size=9.5, gap=3)
    footer(s)


def slide_where_you_trade(prs, ctx, nar):
    s = new_slide(prs)
    hl = ctx["headline"]
    dmk = ctx.get("dark_markets") or []
    title(s, "Where You Trade")
    strapline(s, f"{f_money(hl['notional'], 1)} across {hl['markets']} markets"
                 + (f". Dark is only offered in {', '.join(dmk)}." if dmk
                    else "."))
    picture(s, ctx["charts"].get("notional"), MARGIN, Inches(1.52),
            width=Inches(7.3))

    x = Inches(8.05)
    algo = ctx["algo_table"]
    if not algo.empty:
        box = s.shapes.add_textbox(x, Inches(1.58), Inches(4.7), Inches(0.3))
        _run(box.text_frame.paragraphs[0], "WHAT YOU USE",
             size=9.5, bold=True, color=T_HEAD)
        data = [["Algo", "Orders", "% value", "bps"]]
        colored = {}
        thin = False
        for i, (idx, r) in enumerate(algo.iterrows(), start=1):
            tag = ""
            if r["n"] < MIN_N_QUOTABLE:
                tag, thin = " *", True
            data.append([f"{idx}{tag}", f"{int(r['n']):,}",
                         f"{r['weight_pct']:.1f}%", f_bps(r["bps"])])
            colored[(i, 3)] = "save" if r["bps"] >= 0 else "cost"
        table(s, data, x, Inches(1.90), Inches(4.7), row_h=Inches(0.28),
              colored=colored)
        if thin:
            box = s.shapes.add_textbox(x, Inches(1.92) + Inches(0.28) * len(data),
                                       Inches(4.7), Inches(0.25))
            _run(box.text_frame.paragraphs[0],
                 f"* fewer than {MIN_N_QUOTABLE} orders - not a reliable number",
                 size=8, color=T_MUTED)

    lines = []
    if "market" in ctx["df"]:
        g = ctx["df"].groupby("market")["notional"].sum().sort_values(
            ascending=False)
        top3 = 100 * g.head(3).sum() / g.sum()
        lines.append(f"Your top three markets — **{', '.join(g.head(3).index)}** "
                     f"— are **{top3:.0f}%** of everything you trade.")
    if np.isfinite(hl.get("dark_share_eligible", np.nan)) and dmk:
        eligible = ctx["df"][ctx["df"]["market"].map(
            lambda v: _norm_key(v) in {_norm_key(m) for m in dmk})]
        share = 100 * eligible["notional"].sum() / ctx["df"]["notional"].sum()
        lines.append(f"**{share:.0f}%** of your value sits in markets where we "
                     "offer dark. That is the part the dark discussion later "
                     "can apply to.")
    if np.isfinite(hl.get("adv_pct", np.nan)):
        lines.append(f"Typical order is **{hl['adv_pct']:.1f}%** of the day's "
                     "volume, traded at "
                     f"**{f_pct(hl.get('period_part', np.nan))}** "
                     "participation.")
    bullets(s, lines, x, Inches(3.70), Inches(4.7), Inches(2.6), size=11)
    footer(s)


def slide_cost_sits(prs, ctx):
    s = new_slide(prs)
    title(s, "What Drives the Number")
    hl = ctx["headline"]
    strapline(s, "A market only moves your total if it is both good or bad "
                 "AND big. This shows both together.")
    picture(s, ctx["charts"].get("attribution"), MARGIN, Inches(1.52),
            width=Inches(7.15))

    algo = ctx["algo_table"]
    if not algo.empty:
        data = [["Algo", "Orders", "% value", "bps", "Contrib"]]
        colored = {}
        thin = False
        for i, (idx, r) in enumerate(algo.iterrows(), start=1):
            tag = ""
            if r["n"] < MIN_N_QUOTABLE:
                tag, thin = " *", True
            data.append([f"{idx}{tag}", f"{int(r['n']):,}",
                         f"{r['weight_pct']:.1f}%",
                         f_bps(r["bps"]), f_bps(r["contrib_bps"], 2)])
            colored[(i, 3)] = "save" if r["bps"] >= 0 else "cost"
            colored[(i, 4)] = "save" if r["contrib_bps"] >= 0 else "cost"
        data.append(["TOTAL", f"{hl['orders']:,}", "100.0%",
                     f_bps(hl["bps"]), f_bps(hl["bps"], 2)])
        table(s, data, Inches(7.95), Inches(1.62), Inches(4.75),
              col_w=[Inches(1.15), Inches(0.95), Inches(0.85), Inches(0.9),
                     Inches(0.9)], colored=colored)
        if thin:
            box = s.shapes.add_textbox(Inches(7.95), Inches(3.02),
                                       Inches(4.75), Inches(0.25))
            _run(box.text_frame.paragraphs[0],
                 f"* fewer than {MIN_N_QUOTABLE} orders - not a reliable number",
                 size=8, color=T_MUTED)

    bm = ctx["benchmark_matrix"]
    if not bm.empty:
        box = s.shapes.add_textbox(Inches(7.95), Inches(3.30), Inches(4.75),
                                   Inches(0.3))
        _run(box.text_frame.paragraphs[0],
             "HOW IT LOOKS AGAINST OTHER BENCHMARKS",
             size=9.5, bold=True, color=T_HEAD)
        bcols = [c for c in bm.columns if c.startswith("vs ")]
        data = [["Algo"] + [c.replace("vs ", "") for c in bcols]]
        colored = {}
        for i, (idx, r) in enumerate(bm.iterrows(), start=1):
            data.append([idx] + [f_bps(r[c]) for c in bcols])
            for j, c in enumerate(bcols, start=1):
                if np.isfinite(r[c]):
                    colored[(i, j)] = "save" if r[c] >= 0 else "cost"
        table(s, data, Inches(7.95), Inches(3.62), Inches(4.75),
              row_h=Inches(0.27), size=9, colored=colored)
        box = s.shapes.add_textbox(Inches(7.95), Inches(5.55), Inches(4.75),
                                   Inches(0.9))
        tf = box.text_frame
        tf.word_wrap = True
        _run(tf.paragraphs[0],
             "If a number looks good against its own benchmark but bad "
             "against arrival, that is about when the order was sent, not how "
             "it was traded.",
             size=8.5, color=T_MUTED)
    footer(s)


def slide_industry(prs, ctx, nar):
    if not INDUSTRY_REPORT:
        return
    s = new_slide(prs)
    title(s, "By Industry")
    strapline(s, "Which sectors help and which hurt.")
    picture(s, ctx["charts"].get("industry"), MARGIN, Inches(1.52),
            width=Inches(7.3))

    t = pd.DataFrame(INDUSTRY_REPORT, columns=["industry", "issues",
                                               "weight_pct", "bps",
                                               "contrib_bps"])
    best = t.loc[t["contrib_bps"].idxmax()]
    worst = t.loc[t["contrib_bps"].idxmin()]
    hl = ctx["headline"]

    lines = []
    if worst["contrib_bps"] < 0:
        share = (100 * worst["contrib_bps"] / hl["bps"]
                 if hl["bps"] and hl["bps"] < 0 else np.nan)
        lines.append(
            f"**{worst['industry']}** is your worst sector: "
            f"{worst['bps']:+.1f} bps on {worst['weight_pct']:.0f}% of value, "
            f"taking **{worst['contrib_bps']:+.1f} bps** off the total"
            + (f" — {abs(share):.0f}% of the shortfall."
               if np.isfinite(share) else "."))
    if best["contrib_bps"] > 0:
        lines.append(
            f"**{best['industry']}** is your best: {best['bps']:+.1f} bps on "
            f"{best['weight_pct']:.0f}% of value, adding "
            f"**{best['contrib_bps']:+.1f} bps**.")

    thin = t[t["issues"] < 10]
    if len(thin):
        lines.append("~" + ", ".join(
            f"{r.industry} ({int(r.issues)} names)" for r in thin.itertuples())
            + " — too few names to read anything into.")
    lines.append(
        "Sector is a useful cross-check on the market view: if the same names "
        "keep appearing, the issue is the stock, not the venue or the country.")
    bullets(s, lines, Inches(8.05), Inches(1.72), Inches(4.7), Inches(4.4),
            size=11)
    note(s, "Sector is not in the order-level file, so these figures are taken "
            "from the published industry breakdown for the same period.")
    footer(s)


def slide_difficulty(prs, ctx, nar):
    s = new_slide(prs)
    title(s, "Order Size")
    strapline(s, "How you do on small orders versus big ones.")
    picture(s, ctx["charts"].get("adv"), MARGIN, Inches(1.52),
            width=Inches(7.55))

    adv = ctx["adv_table"]
    lines = []
    if not adv.empty:
        heavy = adv["contrib_bps"].idxmax()
        lines.append(
            f"The **{heavy}** group adds "
            f"**{adv.loc[heavy, 'contrib_bps']:+.1f} bps** to your total — more "
            f"than any other — and it is "
            f"{adv.loc[heavy, 'weight_pct']:.0f}% of what you trade.")
        big = adv[adv["n"] >= 20]
        if len(big) >= 2:
            lines.append(
                f"Order by order, costs run from "
                f"{f_bps(big['bps'].min())} to {f_bps(big['bps'].max())} bps.")
    lines += [l for l in nar["findings"]
              if "small orders do worst" in l or "Best on **" in l]
    lines.append(
        "**Speed is the thing you control.** When an order is small next to "
        "the day's volume, there is normally room to take longer and post "
        "rather than cross the spread.")
    bullets(s, lines, Inches(8.30), Inches(1.70), Inches(4.45), Inches(4.2),
            size=11)
    note(s, "Groups are the order size as a % of the stock's average daily "
            "volume. Groups with fewer than 20 orders are not used above.")
    footer(s)


def slide_venue(prs, ctx, nar):
    s = new_slide(prs)
    hl = ctx["headline"]
    if DARK_STORY:
        title(s, "Dark", accent="how much you use, and who uses it")
        strapline(s, f"{f_pct(hl['dark_share'])} of what you trade goes "
                     "through dark — but it is all coming from one algo.")
    else:
        title(s, "How Your Algos Execute")
        strapline(s, "Where each algo actually puts your volume. This is the "
                     "clearest difference between them.")
    picture(s, ctx["charts"].get("venue"), MARGIN, Inches(1.55),
            width=Inches(7.55))

    venue = ctx["venue_table"]
    lines = []
    seg = VENUE_SEGMENTS_REPORT or {}
    algos = [a for a in ALGO_ORDER if a in seg] or \
            [a for a in seg if a != "TOTAL"]

    if DARK_STORY:
        for a in algos:
            d = seg[a].get("Dark", 0)
            t = seg[a].get("Visible Take", 0)
            if d < 1:
                lines.append(f"**{a}** — no dark at all. It pays the spread "
                             f"on **{t:.0f}%** of what it trades.")
            else:
                lines.append(f"**{a}** — **{d:.0f}%** dark. This is the only "
                             "algo trying the midpoint.")
        if not venue.empty and "orders_with_dark_pct" in venue:
            tot = venue.loc["TOTAL", "orders_with_dark_pct"] \
                if "TOTAL" in venue.index else np.nan
            if np.isfinite(tot):
                lines.append(f"Only **{tot:.0f}%** of orders get any dark at "
                             "all. For the ones that do reach it, there is "
                             "more there than the average suggests.")
        if nar["rec_tags"].get("dark_on"):
            lines.append(nar["rec_tags"]["dark_on"])
    else:
        for a in algos:
            parts = []
            for k, lbl in [("Auction", "in the auction"),
                           ("Visible Post", "posted"),
                           ("Visible Take", "crossing the spread")]:
                v = seg[a].get(k, 0)
                if v >= 5:
                    parts.append(f"**{v:.0f}%** {lbl}")
            if parts:
                lines.append(f"**{a}** — " + ", ".join(parts) + ".")
        for tag in ("venue", "algo_mix"):
            if nar["rec_tags"].get(tag):
                lines.append(nar["rec_tags"][tag])
    bullets(s, lines, Inches(8.30), Inches(1.72), Inches(4.45), Inches(4.2),
            size=11)
    note(s, "Venue split comes from the post-trade report. The order file "
            "only carries the dark share. All figures are % of value traded.")
    footer(s)


def slide_dark_effect(prs, ctx, nar):
    s = new_slide(prs)
    title(s, "Dark", accent="did it actually help?")
    dz = ctx["dark_zero_vs_any"]
    if dz.empty:
        strapline(s, "Nothing in this data traded in dark, so there is "
                     "nothing to compare.", color=T_MUTED)
        footer(s)
        return
    strapline(s, ctx["dark_split_note"])
    picture(s, ctx["charts"].get("dark_zero"), MARGIN, Inches(1.70),
            width=Inches(7.35))

    x = Inches(8.30)
    bullets(s, nar["dark_verdict"], x, Inches(1.75), Inches(4.45), Inches(3.1),
            size=11)
    reg = ctx["regression"]
    if reg:
        _, coef, pval, _ = reg
        kpi(s, x, Inches(4.95), Inches(2.10), Inches(1.05),
            f"{coef:+.2f}", "bps per 10% more dark", fill=T_DARK, big_size=21)
        kpi(s, x + Inches(2.30), Inches(4.95), Inches(2.15), Inches(1.05),
            f"p = {pval:.3f}",
            "solid result" if pval < 0.05 else "could be luck",
            fill=T_HEAD if pval < 0.05 else T_MUTED, big_size=19)
    box = s.shapes.add_textbox(x, Inches(6.15), Inches(4.45), Inches(0.75))
    tf = box.text_frame
    tf.word_wrap = True
    _run(tf.paragraphs[0], nar["caveats"][-1], size=8.5, color=T_MUTED)
    footer(s)


def slide_dark_controlled(prs, ctx, nar):
    """The controlled split gets a slide of its own.

    This is the exhibit the whole dark argument rests on, and it is the first
    thing a sceptical reader attacks. Shrinking it to fit alongside something
    else is how the argument gets lost.
    """
    s = new_slide(prs)
    if not ctx["charts"].get("dark_controlled"):
        return
    title(s, "Dark", accent="or were those orders just easier?")
    strapline(s, "Dark tends to fill on the easy trades. So we split the same "
                 "comparison by how hard each order was.")
    picture(s, ctx["charts"]["dark_controlled"], MARGIN, Inches(1.52),
            width=Inches(8.05))

    x = Inches(8.95)
    lines = [
        "Each panel repeats the same **no dark vs any dark** comparison, but "
        "only inside orders of a similar difficulty.",
        "If dark wins in **every** panel, the gap is real and not just easy "
        "orders getting easy fills.",
        "If it only wins in one or two, then dark helps in **those** "
        "situations and we should say exactly which.",
    ]
    dz = ctx["dark_zero_vs_any"]
    if not dz.empty and len(dz) == 2:
        d = ctx["dark_df"]
        wins = tot = 0
        for c in DARK_CONTROLS:
            t = dark_controlled(d, c)
            if t.empty:
                continue
            for ctrl, g in t.groupby("control"):
                g = g.set_index("group")
                if len(g) < 2:
                    continue
                a = g["bps"].get(ctx["dark_groups"][0], np.nan)
                b = g["bps"].get(ctx["dark_groups"][1], np.nan)
                if np.isfinite(a) and np.isfinite(b):
                    tot += 1
                    wins += int(b < a)
        if tot:
            lines.append(
                f"**On this data, dark comes out cheaper in {wins} of "
                f"{tot} groups.**")
    bullets(s, lines, x, Inches(1.80), Inches(3.85), Inches(4.3), size=11)
    footer(s)


def slide_dark_market(prs, ctx, nar):
    s = new_slide(prs)
    title(s, "Dark", accent="where you can actually get it")
    dmk = ctx.get("dark_markets") or []
    strapline(s, ("We only offer dark in " + ", ".join(dmk) + ", and how much "
                  "you get differs between them.") if dmk else
                 "How much dark you can get varies by market.")
    picture(s, ctx["charts"].get("dark_market"), MARGIN, Inches(1.55),
            width=Inches(7.5))

    # The completion numbers go in as a table: at this width the chart version
    # is too small to read. The PNG is still written to charts/ for lifting.
    comp = ctx["dark_completion"]
    if not comp.empty:
        box = s.shapes.add_textbox(Inches(8.30), Inches(4.30), Inches(4.5),
                                   Inches(0.3))
        _run(box.text_frame.paragraphs[0], "DOES MORE DARK SLOW YOU DOWN?",
             size=9.5, bold=True, color=T_HEAD)
        # participation is ~2% across the book, so 0dp would hide the signal
        cols = [("fill_rate_pct", "Fill %", 0),
                ("participation_pct", "Part %", 1),
                ("median_dur_min", "Mins", 0)]
        cols = [c for c in cols if c[0] in comp]
        data = [["Dark used", "Orders"] + [lbl for _, lbl, _ in cols]]
        for idx, r in comp.iterrows():
            data.append([str(idx).replace(" (none)", ""), f"{int(r['n']):,}"]
                        + [f"{r[c]:,.{dp}f}" if np.isfinite(r[c]) else "—"
                           for c, _, dp in cols])
        table(s, data, Inches(8.30), Inches(4.62), Inches(4.45),
              row_h=Inches(0.26), size=9)
        box = s.shapes.add_textbox(Inches(8.30), Inches(6.28), Inches(4.5),
                                   Inches(0.6))
        tf = box.text_frame
        tf.word_wrap = True
        _run(tf.paragraphs[0],
             "Flat down the columns means using more dark costs you nothing. "
             "A clear drop tells you where to stop.", size=8.5, color=T_MUTED)

    lines = []
    dmkt = ctx["dark_by_market"]
    if not dmkt.empty:
        hi = dmkt.sort_values("avail_orders_pct", ascending=False).head(2)
        lo = dmkt.sort_values("avail_orders_pct").head(2)
        lines.append("Easiest to get dark: **" + "**, **".join(
            f"{i} ({r.avail_orders_pct:.0f}%)" for i, r in hi.iterrows()) + "**.")
        lines.append("Hardest: **" + "**, **".join(
            f"{i} ({r.avail_orders_pct:.0f}%)" for i, r in lo.iterrows())
            + "**. There is simply not much there to use.")
    if nar["rec_tags"].get("dark_market"):
        lines.append(nar["rec_tags"]["dark_market"])
    # the "almost no dark available" caveat already appears on the summary
    # slide; repeating it here crowds the column
    bullets(s, lines, Inches(8.30), Inches(1.75), Inches(4.45), Inches(2.4),
            size=10.5)
    note(s, f"Markets with fewer than {MIN_N_MARKET} orders are left out. "
            "Order counts for each market are in the tables file — check them "
            "before quoting any single market on its own.")
    footer(s)


def slide_recs(prs, ctx, nar):
    s = new_slide(prs)
    title(s, "Recommendations")
    hl = ctx["headline"]
    strapline(s, "Biggest win first.")

    for i, r in enumerate(nar["recs"][:4]):
        y = Inches(1.65 + i * 1.13)
        num = s.shapes.add_shape(1, MARGIN, y, Inches(0.46), Inches(0.46))
        num.fill.solid()
        num.fill.fore_color.rgb = T_HEAD
        num.line.fill.background()
        num.shadow.inherit = False
        tf = num.text_frame
        tf.margin_top = Inches(0.02)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(tf.paragraphs[0], str(i + 1), size=15, bold=True, color=T_WHITE)
        bullets(s, [r], Inches(1.25), y - Inches(0.04), Inches(11.4),
                Inches(1.0), size=12, marker="")

    box = s.shapes.add_textbox(MARGIN, Inches(6.15), Inches(12.1), Inches(0.7))
    tf = box.text_frame
    tf.word_wrap = True
    _run(tf.paragraphs[0], "Next step:  ", size=11, bold=True, color=T_HEAD)
    _run(tf.paragraphs[0],
         "let us run the same analysis again next quarter and see if the "
         f"number moved. What matters is improving on your own "
         f"{f_bps(hl['bps'])} bps on the same kind of flow — not comparing "
         "against anyone else.", size=11, color=T_INK)
    footer(s)


def build_deck(ctx: dict, nar: dict, out: Path) -> Path:
    _PAGE[0] = 0
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    slide_cover(prs, ctx)
    slide_summary(prs, ctx, nar)
    slide_where_you_trade(prs, ctx, nar)
    slide_cost_sits(prs, ctx)
    slide_industry(prs, ctx, nar)
    slide_difficulty(prs, ctx, nar)
    slide_venue(prs, ctx, nar)
    if DARK_STORY:
        slide_dark_effect(prs, ctx, nar)
        slide_dark_controlled(prs, ctx, nar)
        slide_dark_market(prs, ctx, nar)
    slide_recs(prs, ctx, nar)
    if IS_SAMPLE[0]:
        out = out.with_name(out.stem + "_SAMPLE_DATA" + out.suffix)
    prs.save(str(out))
    log(f"  saved {out.name}  ({len(prs.slides._sldIdLst)} slides)")
    return out


# ===========================================================================
# 10. SYNTHETIC SAMPLE DATA
# ===========================================================================
# Shaped like the real export - same headers, same banner row, same algo mix,
# same market weights - so the pipeline can be smoke-tested before the client
# file is available. The numbers are invented; do not present them.

SAMPLE_MARKETS = [
    ("HongKong",      0.271, 10.6, "HK"), ("Japan",     0.317, 3.6,  "JP"),
    ("Taiwan",        0.150, 15.4, "TT"), ("India",     0.135, 3.8,  "IN"),
    ("Stock Connect", 0.065, 4.2,  "CH"), ("Indonesia", 0.033, 26.1, "IJ"),
    ("Malaysia",      0.012, 26.7, "MK"), ("Thailand",  0.008, 33.7, "TB"),
    ("Australia",     0.004, 5.8,  "AU"), ("Philippines", 0.003, 7.4, "PM"),
    ("Singapore",     0.002, 9.7,  "SP"),
]
SAMPLE_ALGOS = [("VWAP", 1270, 0.906), ("IIS", 774, 0.092), ("PROG", 30, 0.002)]


def make_sample(path: Path, n_target=2074, seed=SEED) -> Path:
    rng = np.random.default_rng(seed)
    rows = []
    names = [m[0] for m in SAMPLE_MARKETS]
    wts = np.array([m[1] for m in SAMPLE_MARKETS])
    wts = wts / wts.sum()

    for algo, n_orders, notional_share in SAMPLE_ALGOS:
        total_notional = 505_182_392 * notional_share
        for _ in range(n_orders):
            mi = rng.choice(len(SAMPLE_MARKETS), p=wts)
            mkt, _, spread, suf = SAMPLE_MARKETS[mi]
            notional = max(2_000.0,
                           rng.lognormal(np.log(total_notional / n_orders), 1.1))
            price = float(rng.uniform(0.5, 180))
            shares = notional / price
            adv = float(np.clip(rng.lognormal(np.log(0.9), 1.0), 0.01, 30))
            dur = float(np.clip(rng.normal(230, 130), 5, 400))
            spr = float(max(0.5, rng.normal(spread, spread * 0.35)))

            # dark only reachable inside VWAP, and more so in liquid markets
            if algo == "VWAP":
                reach = {"HongKong": .55, "Japan": .62, "Taiwan": .40,
                         "India": .30, "Stock Connect": .05}.get(mkt, .12)
                dark = float(np.clip(rng.beta(1.6, 5.0) * 100, 0, 90)) \
                    if rng.random() < reach else 0.0
            else:
                dark = 0.0

            # cost model: rises with adv and spread, falls a little with dark
            base = (0.9 + 0.55 * adv + 0.16 * spr
                    - 0.045 * dark - 0.004 * dur)
            noise = rng.normal(0, 16 + 0.8 * spr)
            cost = base + noise
            if algo == "PROG":
                cost += rng.normal(115, 60)

            side = "BUY" if rng.random() < 0.66 else "SELL"
            date = pd.Timestamp("2026-01-05") + pd.Timedelta(
                days=int(rng.integers(0, 228)))
            rows.append({
                "Date": date.strftime("%Y.%m.%d"),
                "Strike Time": "10:04:00", "End Time": "16:00:00",
                "Dur (mins)": round(dur, 1),
                "Symbol": f"{rng.integers(1, 9999):04d}.{suf}",
                "Side": side, "Order Type": "Market",
                "Order Qty": int(shares / max(rng.uniform(0.55, 1.0), .01)),
                "Shares Exec": int(shares),
                "Value Exec": round(notional, 2),
                "Dark Qty": int(shares * dark / 100),
                "Dark %": round(dark / 100, 4),
                "Dark Value": round(notional * dark / 100, 2),
                "Weight": round(notional / 505_182_392, 6),
                "% Adv": round(adv / 100, 5),
                "Period Part": round(float(np.clip(
                    rng.normal(0.019, 0.012), 0.001, 0.35)), 5),
                "Avg. Price (local)": round(price, 4),
                "FX": round(float(rng.uniform(0.0001, 1.0)), 6),
                "ClientAlgo": algo,
                "Hist Spread": round(spr, 2),
                "Period Spread": 0.0,
                "Market": mkt,
                "Arrival Price": round(price * 1.0001, 4),
                "Arrival ImpBps": round(cost + rng.normal(4, 22), 1),
                "Arrival ImpAmt": 0.0,
                "PVWAP Price": round(price, 4),
                "PVWAP ImpBps": round(cost, 1),
                "PVWAP ImpAmt": round(cost * notional / 1e4, 1),
                "VWAP Price": round(price, 4),
                "VWAP ImpBps": round(cost + rng.normal(0, 12), 1),
                "VWAP ImpAmt": 0.0,
                "Close Price": round(price, 4),
                "Close ImpBps": round(cost + rng.normal(0, 30), 1),
                "Close ImpAmt": 0.0,
            })

    df = pd.DataFrame(rows).sample(frac=1, random_state=seed).head(n_target)
    path.parent.mkdir(parents=True, exist_ok=True)

    # write the merged banner on row 1 and the header on row 2, as the tool does
    with pd.ExcelWriter(path, engine="openpyxl") as xl:
        df.to_excel(xl, sheet_name="symboldetailstable_full", index=False,
                    startrow=1)
        ws = xl.sheets["symboldetailstable_full"]
        ws.cell(row=1, column=1, value="Symbol Details")
    log(f"  wrote synthetic sample: {path}  ({len(df):,} rows)")
    log("  ** SYNTHETIC DATA — the numbers are invented. Do not present.")
    return path




# ===========================================================================
# 12. SIMPLE DECK  (--simple)
# ===========================================================================
# Built entirely from the published report tables in CLIENTS. No order file is
# read, so nothing here can be synthetic: every figure is the one the client
# already has in their own post-trade report.

def chart_marketcap(t: pd.DataFrame, outdir: Path) -> Path:
    """Contribution by market-cap band."""
    if t is None or t.empty:
        return None
    t = t.sort_values("pnl_ccy")
    fig, ax = plt.subplots(figsize=(10.0, 3.6))
    ys = np.arange(len(t))
    vals = (t["pnl_ccy"] / 1e3).values
    ax.barh(ys, vals, height=0.55, color=_sign_colors(vals), zorder=3)
    ax.set_yticks(ys, [f"{i}\n{f_money(r.notional_m * 1e6)} traded · "
                       f"{r.bps:+.1f} bps"
                       for i, r in t.iterrows()], fontsize=9.5)
    _hbar_labels(ax, ys, vals, fmt="{:+,.0f}k", floor=0.5)
    _zero_line(ax)
    span = float(np.nanmax(np.abs(vals))) if len(vals) else 1.0
    ax.set_xlim(min(-span * 0.5, float(np.nanmin(vals)) * 1.6), span * 1.6)
    _finish(ax, "Money made and lost, by company size",
            xlabel=f"{CURRENCY} thousands better or worse than benchmark")
    return _save(fig, "03c_marketcap", outdir)


VENUE_COLS = ["market", "Auction", "Visible Post", "Visible Take", "Dark"]


def venue_country(rows, mkt: pd.DataFrame = None):
    """Venue mix per market, joined to that market's result.

    Rows are transcribed by hand, so each is checked: the four venue shares
    must sum to 100. This is the table that decides whether a losing market
    needs to post more or to slow down - opposite instructions, and the
    aggregate venue split cannot tell them apart.
    """
    if not rows:
        return None
    t = pd.DataFrame(list(rows), columns=VENUE_COLS).set_index("market")
    bad = t[(t.sum(axis=1) - 100).abs() > 1.5]
    for m in bad.index:
        log(f"  ** CHECK venue_country: {m} sums to {t.loc[m].sum():.1f}%, "
            "not 100. Re-read that row.")
    if mkt is not None and not mkt.empty:
        t = t.join(mkt[["bps", "weight_pct", "spread_bps", "notional_m"]],
                   how="left")
    return t


def chart_venue_country(t: pd.DataFrame, outdir: Path, min_weight=1.0) -> Path:
    """Where each market's volume actually goes."""
    if t is None or t.empty:
        return None
    d = t[t.get("weight_pct", pd.Series(100, index=t.index)) >= min_weight]
    if d.empty:
        return None
    d = d.sort_values("weight_pct", ascending=True)
    fig, ax = plt.subplots(figsize=(10.0, 0.42 * len(d) + 2.0))
    ys = np.arange(len(d))
    left = np.zeros(len(d))
    for seg in VENUE_ORDER:
        vals = d[seg].values.astype(float)
        ax.barh(ys, vals, left=left, height=0.6, color=VENUE_COLORS[seg],
                edgecolor=SURFACE, linewidth=1.6, zorder=3, label=seg)
        for y, v, l in zip(ys, vals, left):
            if v >= 8:
                ax.text(l + v / 2, y, f"{v:.0f}", ha="center", va="center",
                        fontsize=9, color="white", fontweight="bold")
        left += vals
    labels = []
    for m, r in d.iterrows():
        bits = []
        if np.isfinite(r.get("spread_bps", np.nan)):
            bits.append(f"{r['spread_bps']:.1f} bps spread")
        if np.isfinite(r.get("bps", np.nan)):
            bits.append(f"{r['bps']:+.1f}")
        labels.append(f"{m}\n" + " · ".join(bits) if bits else str(m))
    ax.set_yticks(ys, labels, fontsize=9)
    ax.set_xlim(0, 100)
    ax.set_xticks(range(0, 101, 20), [f"{v}%" for v in range(0, 101, 20)])
    ax.legend(handles=[Patch(facecolor=VENUE_COLORS[x], label=x)
                       for x in VENUE_ORDER],
              loc="upper center", bbox_to_anchor=(0.5, -0.10), ncol=4,
              fontsize=10)
    _finish(ax, "How your shares changed hands, country by country")
    return _save(fig, "07_venue_by_market", outdir)


def side_lens(rows, total_notional, min_weight=2.0):
    """Where a market's result splits by side.

    Both books carry large buy/sell gaps inside a market that mostly cancel at
    the total, so the headline hides them. Working in money rather than basis
    points finds the single cell that matters.
    """
    if not rows or not np.isfinite(total_notional):
        return None
    t = pd.DataFrame(list(rows), columns=["market", "side", "weight_pct", "bps"])
    t["notional"] = t["weight_pct"] / 100 * total_notional
    t["pnl_ccy"] = t["bps"] * t["notional"] / 1e4
    t["label"] = t["market"] + " " + t["side"].str.lower() + "s"

    wide = t.pivot_table(index="market", columns="side", values="bps")
    wide["gap"] = wide.get("BUY", np.nan) - wide.get("SELL", np.nan)
    wt = t.groupby("market")["weight_pct"].sum()
    wide["weight_pct"] = wt
    big = wide[(wide["weight_pct"] >= min_weight) & wide["gap"].notna()]

    shown = t[t["weight_pct"] >= 1.0].sort_values("pnl_ccy")
    return {
        "table": t.sort_values("pnl_ccy"),
        "shown": shown,
        "worst": t.loc[t["pnl_ccy"].idxmin()] if len(t) else None,
        "best": t.loc[t["pnl_ccy"].idxmax()] if len(t) else None,
        "widest": (big["gap"].abs().idxmax() if len(big) else None),
        "wide": big,
    }


def chart_side(lens: dict, outdir: Path) -> Path:
    """Money made or lost in each market, split by side."""
    if not lens or lens["shown"].empty:
        return None
    t = lens["shown"]
    fig, ax = plt.subplots(figsize=(10.0, 0.34 * len(t) + 2.2))
    ys = np.arange(len(t))[::-1]
    vals = (t["pnl_ccy"] / 1e3).values
    ax.barh(ys, vals, height=0.62, color=_sign_colors(vals), zorder=3)
    ax.set_yticks(ys, [f"{r.label}   ({f_money(r.notional)} · "
                       f"{r.bps:+.1f} bps)"
                       for r in t.itertuples()], fontsize=9)
    _hbar_labels(ax, ys, vals, fmt="{:+,.0f}k", floor=0.5)
    _zero_line(ax)
    span = float(np.nanmax(np.abs(vals))) if len(vals) else 1.0
    ax.set_xlim(min(-span * 1.35, float(np.nanmin(vals)) * 1.35), span * 1.35)
    _finish(ax, "Money made and lost, by country and direction",
            xlabel=f"{CURRENCY} thousands better or worse than benchmark")
    _axis_note(ax, y=-0.09)
    return _save(fig, "06_side", outdir)

def spread_lens(mkt: pd.DataFrame, min_orders=50):
    """How much of the result is explained by the spread that was available?

    A passive / midpoint-seeking approach earns the spread. Where the spread is
    wide there is a lot to earn; where it is tight there is nothing, and impact
    and timing show through instead. Fitting result against spread says whether
    that is what is happening, and names the markets that miss the line.

    Returns None when there are too few markets to fit anything meaningful.
    """
    if mkt is None or mkt.empty or "spread_bps" not in mkt:
        return None
    t = mkt[(mkt["n"] >= min_orders) & mkt["spread_bps"].notna()
            & mkt["bps"].notna()].copy()
    if len(t) < 5:
        return None
    x, y = t["spread_bps"].values, t["bps"].values
    if np.ptp(x) < 1e-9:
        return None
    r = float(np.corrcoef(x, y)[0, 1])
    slope, intercept = np.polyfit(x, y, 1)
    t["predicted"] = slope * t["spread_bps"] + intercept
    t["vs_fit"] = t["bps"] - t["predicted"]
    breakeven = float(-intercept / slope) if slope else np.nan
    tight = t[t["spread_bps"] < breakeven] if np.isfinite(breakeven) else t.iloc[0:0]
    return {
        "table": t.sort_values("spread_bps"),
        "r": r, "slope": float(slope), "intercept": float(intercept),
        "breakeven": breakeven,
        "tight_share": float(tight["weight_pct"].sum()),
        "tight_names": list(tight.index),
        "tight_bps": (float(np.average(tight["bps"], weights=tight["weight_pct"]))
                      if len(tight) and tight["weight_pct"].sum() else np.nan),
        # rank the miss by money, not by bps: a 13 bps gap on a 3% market is
        # worth far less than an 12 bps gap on a 15% one, and the client can
        # only act on the second
        "underperformer": ((t["vs_fit"] * t["notional_m"]).idxmin()
                           if len(t) else None),
    }


def chart_spread_lens(lens: dict, outdir: Path) -> Path:
    """Result against the spread that was available, market by market."""
    if not lens:
        return None
    t = lens["table"]
    fig, ax = plt.subplots(figsize=(10.0, 5.4))

    xs = np.linspace(0, float(t["spread_bps"].max()) * 1.12, 50)
    ax.plot(xs, lens["slope"] * xs + lens["intercept"], color=BASELINE,
            lw=1.4, zorder=2)
    ax.axhline(0, color=BASELINE, lw=1.1, zorder=2)
    if np.isfinite(lens["breakeven"]):
        ax.axvline(lens["breakeven"], color=INK_MUTED, lw=1.0, ls=":", zorder=2)
        ax.text(lens["breakeven"], ax.get_ylim()[1],
                f"  break-even ≈ {lens['breakeven']:.1f} bps of spread",
                fontsize=9, color=INK_MUTED, va="top")

    sizes = 120 + 1400 * (t["weight_pct"] / max(t["weight_pct"].max(), 1e-9))
    ax.scatter(t["spread_bps"], t["bps"], s=sizes,
               color=[SAVE_COLOR if v >= 0 else COST_COLOR for v in t["bps"]],
               alpha=0.85, edgecolor=SURFACE, linewidth=1.6, zorder=3)
    # the tight-spread markets cluster, so stagger their labels instead of
    # stacking them on top of one another
    xr = max(float(t["spread_bps"].max() - t["spread_bps"].min()), 1e-9)
    yr = max(float(t["bps"].max() - t["bps"].min()), 1e-9)
    placed = []
    for name, r in t.iterrows():
        x, y = float(r["spread_bps"]), float(r["bps"])
        offs = [(0, 17), (0, -22), (0, 30), (0, -34), (0, 43)]
        pick = offs[0]
        for cand in offs:
            clash = any(abs(x - px) / xr < 0.09
                        and abs((y + cand[1] * yr / 380) - py) / yr < 0.07
                        for px, py in placed)
            if not clash:
                pick = cand
                break
        placed.append((x, y + pick[1] * yr / 380))
        ax.annotate(str(name), (x, y), textcoords="offset points",
                    xytext=pick, ha="center", fontsize=9.5, color=INK)
    _finish(ax, "Result against the buy/sell gap in each market",
            xlabel="average buy/sell gap in that market (bps)",
            ylabel="result in bps (higher is better)", axis="y")
    ax.text(1.0, -0.15,
            "Each bubble is a market, sized by how much you trade there. "
            f"They line up {abs(lens['r']) * 100:.0f}% of the way to a "
            "straight line.",
            transform=ax.transAxes, fontsize=9, color=INK_MUTED,
            ha="right", va="top")
    return _save(fig, "05_spread_lens", outdir)



def drop_markets() -> None:
    """Remove excluded markets from every report table, and say what went.

    The published totals are left alone: they are the client's own headline
    and must keep matching their report. Anything dropped is logged with its
    orders, value and effect, so the difference is never silent.
    """
    if not EXCLUDE_MARKETS:
        return
    rows = REPORT.get("country") or []
    gone = [r for r in rows if r[0] in EXCLUDE_MARKETS]
    if not gone:
        log(f"  exclude_markets: {_and(EXCLUDE_MARKETS)} not in the country "
            "table - nothing to drop")
        return
    n_ord = sum(r[1] for r in gone)
    log(f"  excluded from the review: {_and([r[0] for r in gone])} - "
        f"{n_ord:,} order{'' if n_ord == 1 else 's'}, "
        f"{f_money(sum(r[3] for r in gone))}, "
        f"{f_money(sum(r[3] * r[8] / 1e4 for r in gone))} of effect. "
        "The published totals are unchanged.")
    REPORT["country"] = [r for r in rows if r[0] not in EXCLUDE_MARKETS]
    for key in ("venue_country", "country_side"):
        if REPORT.get(key):
            REPORT[key] = [r for r in REPORT[key]
                           if r[0] not in EXCLUDE_MARKETS]


def analyse_report_only() -> dict:
    """Everything the simple deck needs, from the report tables alone."""
    log("")
    log("-" * 74)
    log("ANALYSIS  (published report figures only - no order file)")
    log("-" * 74)
    tot = REPORT.get("totals") or {}
    if not tot:
        sys.exit("FATAL - this client has no published totals in CLIENTS, so "
                 "--simple has nothing to build from.")

    drop_markets()
    hl = headline_from_report(tot)
    country = table_from_report(REPORT.get("country"), "market", tot.get("bps"))
    hl["markets"] = int(len(country)) if not country.empty else np.nan
    hl["symbols"] = np.nan
    seg_total = (VENUE_SEGMENTS_REPORT or {}).get("TOTAL", {})
    hl["dark_share"] = float(seg_total.get("Dark", np.nan))
    hl["dark_order_share"] = np.nan
    hl["dark_share_eligible"] = np.nan

    ctx = {
        "df": pd.DataFrame(),
        "headline": hl,
        "algo_table": algo_table_from_report(ALGO_REPORT, ALGOS_STUDIED),
        "market_table": country,
        "adv_table": table_from_report(REPORT.get("adv"), "order size",
                                       tot.get("bps")),
        "cap_table": table_from_report(REPORT.get("marketcap"), "market cap",
                                       tot.get("bps")),
        "side_table": table_from_report(REPORT.get("side"), "side",
                                        tot.get("bps")),
        "industry": INDUSTRY_REPORT,
        "dark_markets": list(DARK_MARKETS or []),
    }
    ctx["venue_country"] = venue_country(REPORT.get("venue_country"),
                                         ctx["market_table"])
    ctx["side_lens"] = side_lens(REPORT.get("country_side"), hl["notional"])
    if ctx["side_lens"] and ctx["side_lens"]["worst"] is not None:
        w = ctx["side_lens"]["worst"]
        log(f"  side lens: worst cell {w['label']} "
            f"{w['bps']:+.1f} bps on {w['weight_pct']:.1f}% of value "
            f"= {f_money(w['pnl_ccy'])}")
    ctx["spread_lens"] = spread_lens(ctx["market_table"])
    if ctx["spread_lens"]:
        L = ctx["spread_lens"]
        log(f"  spread lens: r = {L['r']:+.2f}, break-even "
            f"{L['breakeven']:.1f} bps of spread, "
            f"{L['tight_share']:.0f}% of value below it")
    log(f"  {hl['orders']:,} orders · {f_money(hl['notional'], 1)} · "
        f"{f_bps(hl['bps'])} bps · {hl['markets']} markets")
    return ctx


def make_simple_charts(ctx: dict, outdir: Path) -> dict:
    log("")
    log("-" * 74)
    log("CHARTS")
    log("-" * 74)
    apply_style()
    hl = ctx["headline"]
    c = {}
    if not ctx["algo_table"].empty:
        c["algo"] = chart_algo(ctx["algo_table"], outdir, hl)
    if not ctx["market_table"].empty:
        c["notional"] = chart_notional_by_country(
            pd.DataFrame(), outdir, ctx.get("dark_markets"),
            report=ctx["market_table"])
        c["attribution"] = chart_attribution(ctx["market_table"], outdir,
                                             by="market", total_bps=hl["bps"])
    if ctx["industry"]:
        c["industry"] = chart_industry(ctx["industry"], outdir, hl["bps"],
                                       notional=hl["notional"])
    if not ctx["adv_table"].empty:
        c["adv"] = chart_adv(ctx["adv_table"], outdir)
    if not ctx["cap_table"].empty:
        c["cap"] = chart_marketcap(ctx["cap_table"], outdir)
    c["venue"] = chart_venue(outdir)
    # only worth an exhibit where the relationship actually holds; on a book
    # where it does not, the chart would invite a story the data will not carry
    if ctx.get("venue_country") is not None:
        c["venue_country"] = chart_venue_country(ctx["venue_country"], outdir)
    if ctx.get("side_lens"):
        c["side"] = chart_side(ctx["side_lens"], outdir)
    L = ctx.get("spread_lens")
    if L and abs(L["r"]) >= 0.6:
        c["spread"] = chart_spread_lens(L, outdir)
    elif L:
        log(f"  spread lens r = {L['r']:+.2f} - too weak to show, skipping "
            "that exhibit")
    ctx["charts"] = c
    return c


# --- narrative items --------------------------------------------------------
# Every finding, recommendation and caveat is a dict with a tag, so a slide can
# ask for the lines it wants by name. Matching on substrings of the prose broke
# silently every time a sentence was reworded.

def _say(bucket, tag, text, value=None):
    bucket.append({"tag": tag, "text": text, "value": value})


def pick(items, *tags):
    """Texts carrying any of these tags, in the order the tags are given."""
    out = []
    for t in tags:
        out += [i["text"] for i in items if i["tag"] == t]
    return out


def texts(items):
    return [i["text"] for i in items]


def _and(names):
    """Join names the way a person reads them out."""
    names = list(names)
    if len(names) <= 1:
        return "".join(names)
    return ", ".join(names[:-1]) + " and " + names[-1]


def _band(name):
    """'0-1%' -> 'under 1% of a day\'s volume'."""
    n = str(name).strip()
    if n.startswith("0-"):
        return f"under {n.split('-', 1)[1]} of a day's volume"
    return f"at {n} of a day's volume"


def _money_from(bps, weight_pct, notional):
    """Currency effect of `bps` earned on `weight_pct` of `notional`."""
    if not all(np.isfinite(x) for x in (bps, weight_pct, notional)):
        return np.nan
    return bps * weight_pct / 100.0 * notional / 1e4


def build_simple_narrative(ctx: dict) -> dict:
    """Findings and advice, derived from the published tables.

    House style, because a client reads this once and often out loud: one idea
    per sentence, no sentence much longer than twenty words, and every
    recommendation names the thing to change rather than the place to look.
    """
    hl = ctx["headline"]
    algo = ctx["algo_table"]
    mkt = ctx["market_table"]
    adv = ctx["adv_table"]
    cap = ctx["cap_table"]
    ind = pd.DataFrame(ctx["industry"] or [],
                       columns=["name", "issues", "weight_pct", "bps",
                                "contrib_bps"]).set_index("name")
    seg = VENUE_SEGMENTS_REPORT or {}
    vc = ctx.get("venue_country")
    S = ctx.get("side_lens")
    L = ctx.get("spread_lens")
    NOT_ = hl["notional"]
    ahead = np.isfinite(hl["bps"]) and hl["bps"] >= 0
    n = {"findings": [], "recs": [], "notes": []}
    F, R, N = n["findings"], n["recs"], n["notes"]

    def _vc(m, col):
        try:
            v = float(vc.loc[m, col])
            return v if np.isfinite(v) else np.nan
        except Exception:
            return np.nan

    # ---- the headline ----------------------------------------------------
    _say(F, "headline",
         (f"You beat the benchmark by **{f_bps(hl['bps'])} bps**. That is "
          f"about **{f_money(abs(hl['pnl_ccy']))}** kept over the period."
          if ahead else
          f"You came in **{f_bps(hl['bps'])} bps** below the benchmark. That "
          f"cost about **{f_money(abs(hl['pnl_ccy']))}** over the period."))

    # ---- which algo actually decides it ----------------------------------
    if not algo.empty:
        big = algo[algo["n"] >= MIN_N_QUOTABLE]
        if not big.empty:
            k = big["contrib_bps"].idxmax() if ahead else big["contrib_bps"].idxmin()
            r = big.loc[k]
            _say(F, "algo",
                 f"**{k}** handled {f_money(r['notional_m'] * 1e6)} of the "
                 f"{f_money(NOT_)} you traded. It returned "
                 f"**{f_bps(r['bps'])} bps**, so your result is really the "
                 f"{k} result.")
        for k in algo[algo["bps"].abs() < 1e-9].index:
            _say(N, "zero_bench",
                 f"{k} trades in the closing auction and is scored against the "
                 f"close, so it always scores 0.0. That is "
                 f"{f_money(algo.loc[k, 'notional_m'] * 1e6)} of trading we "
                 "cannot judge. Ask for arrival-price scoring instead.")

    # ---- markets ---------------------------------------------------------
    good = bad = mkt.iloc[0:0]
    if not mkt.empty:
        good = mkt[mkt["contrib_bps"] > 0].sort_values("contrib_bps",
                                                       ascending=False)
        bad = mkt[mkt["contrib_bps"] < 0].sort_values("contrib_bps")
        if len(good):
            g, r = good.index[0], good.iloc[0]
            if ahead:
                _say(F, "market_good",
                     f"**{g}** carries the result. You traded "
                     f"{f_money(r['notional_m'] * 1e6)} there at "
                     f"**{f_bps(r['bps'])} bps**, worth "
                     f"**{f_money(r['pnl_ccy'])}** of your "
                     f"{f_money(abs(hl['pnl_ccy']))}.")
            else:
                _say(F, "market_good",
                     f"**{g}** is the one market that helps. You traded "
                     f"{f_money(r['notional_m'] * 1e6)} there at "
                     f"**{f_bps(r['bps'])} bps**, adding "
                     f"**{f_money(r['pnl_ccy'])}** back.")
        if len(bad):
            r = bad.iloc[0]
            two = " and ".join(f"**{x}**" for x in bad.head(2).index)
            _say(F, "market_bad",
                 f"{two} cost you the most. In {bad.index[0]} you traded "
                 f"{f_money(r['notional_m'] * 1e6)} at "
                 f"**{f_bps(r['bps'])} bps**, costing "
                 f"**{f_money(abs(r['pnl_ccy']))}**.")

    # ---- buys against sells ---------------------------------------------
    if S and S["worst"] is not None:
        w, b = S["worst"], S["best"]
        tot_ccy = abs(hl["pnl_ccy"]) if np.isfinite(hl["pnl_ccy"]) else np.nan
        _say(F, "side_worst",
             f"Split by buy and sell, one line stands out. Your "
             f"**{w['label']}** were {f_money(w['notional'])} at "
             f"**{f_bps(w['bps'])} bps**. That cost "
             f"**{f_money(abs(w['pnl_ccy']))}**"
             + (f", more than the whole shortfall of {f_money(tot_ccy)}."
                if np.isfinite(tot_ccy) and abs(w["pnl_ccy"]) > tot_ccy
                and hl["bps"] < 0 else "."))
        if b is not None and b["pnl_ccy"] > 0:
            _say(F, "side_best",
                 f"Your **{b['label']}** went the other way: "
                 f"{f_money(b['notional'])} at **{f_bps(b['bps'])} bps**, "
                 f"worth **{f_money(b['pnl_ccy'])}**. Same desk, same period.")
        if S["widest"] is not None:
            g = S["wide"].loc[S["widest"]]
            _say(F, "side_widest",
                 f"**{S['widest']}** has the widest split. Buys "
                 f"{g.get('BUY', np.nan):+.1f} bps, sells "
                 f"{g.get('SELL', np.nan):+.1f} bps. That is "
                 f"**{abs(g['gap']):.0f} bps** apart in one market.")
        _say(N, "side_caution",
             "A buy/sell gap lasting a whole period is often about market "
             "direction, not the algo. Confirm that before acting.")

    # ---- how each market trades, and therefore what to change ------------
    # This is the split that separates two opposite instructions: a market that
    # pays the spread on everything can be fixed with a setting, while one that
    # already rests in the queue cannot, and telling it to post more would
    # waste the desk's time.
    takers = []
    if vc is not None and not vc.empty and not mkt.empty:
        for m, r in mkt.iterrows():
            if m not in vc.index or r["weight_pct"] < 1.0:
                continue
            take, sp = _vc(m, "Visible Take"), r["spread_bps"]
            if not np.isfinite(take) or take < 50 or not np.isfinite(sp):
                continue
            if r["bps"] >= max(0.0, hl["bps"]):
                continue    # it is already beating the book; leave it alone
            takers.append((m, take, _money_from(sp / 2 * take / 100,
                                                r["weight_pct"], NOT_)))
        takers = [t for t in takers if np.isfinite(t[2])]
        takers.sort(key=lambda x: -x[2])

    # markets that get their own recommendation below, so the posting advice
    # does not repeat what has already been said
    covered = list(bad[bad["weight_pct"] >= 2.0].head(2).index) if len(bad) else []

    if takers:
        m, take, give = takers[0]
        _say(F, "venue_take",
             f"**{m} pays the spread on almost everything.** You take the "
             f"price on offer for {take:.0f}% of what you trade there, and "
             f"rest in the queue for {_vc(m, 'Visible Post'):.0f}%. At a "
             f"{mkt.loc[m, 'spread_bps']:.1f} bps spread that is about "
             f"**{f_money(give)}** given away.")
        left = [t for t in takers if t[0] not in covered][:3]
        if left:
            total_give = sum(t[2] for t in left)
            _say(R, "rec_post",
                 f"**Turn on passive posting in {_and(t[0] for t in left)}.** "
                 "Those orders take the price on offer for at least "
                 f"{min(t[1] for t in left):.0f}% of their volume. That hands "
                 f"over about {f_money(total_give)} of spread each period. It "
                 "is an algo setting, so it is entirely in your control.",
                 value=total_give)

    passives = []
    if vc is not None and not vc.empty and len(bad):
        for m, r in bad.iterrows():
            if m in vc.index and r["weight_pct"] >= 3.0 \
                    and _vc(m, "Visible Take") <= 20:
                passives.append(m)
    if passives:
        ex = passives[0]
        _say(F, "venue_passive",
             f"**{_and(passives)} already "
             + ("trades" if len(passives) == 1 else "trade")
             + f" the right way.** {ex} "
             f"rests in the queue for {_vc(ex, 'Visible Post'):.0f}% of its "
             f"volume and pays the spread on only {_vc(ex, 'Visible Take'):.0f}%. "
             "Whatever is costing you there, it is not the venue.")

    # ---- one recommendation per losing market ----------------------------
    side_tbl = S["table"] if S else None
    if len(bad):
        for m, r in bad[bad["weight_pct"] >= 2.0].head(2).iterrows():
            lost = abs(r["pnl_ccy"])
            bits = [f"**Work on {m}.** You traded "
                    f"{f_money(r['notional_m'] * 1e6)} there at "
                    f"{f_bps(r['bps'])} bps, costing {f_money(lost)}."]
            if side_tbl is not None:
                rows_m = side_tbl[side_tbl["market"] == m]
                # only worth naming a side when there are two of them
                if len(rows_m) > 1 and rows_m["pnl_ccy"].min() < 0:
                    ws = rows_m.loc[rows_m["pnl_ccy"].idxmin()]
                    bits.append(f"The {ws['side'].lower()} side carries it: "
                                f"{f_bps(ws['bps'])} bps, "
                                f"{f_money(abs(ws['pnl_ccy']))}.")
            take, post = _vc(m, "Visible Take"), _vc(m, "Visible Post")
            dark, sp = _vc(m, "Dark"), r["spread_bps"]
            if np.isfinite(take) and take >= 50:
                bits.append(f"It takes the price on offer for {take:.0f}% of "
                            "its volume and "
                            + ("never rests in the queue."
                               if np.isfinite(post) and post < 1 else
                               f"rests in the queue for only {post:.0f}%.")
                            + " Turn passive posting on there first.")
            elif np.isfinite(take) and take <= 20:
                bits.append(f"Venue is not the issue: {post:.0f}% already "
                            f"rests in the queue and {dark:.0f}% goes to the "
                            "midpoint. Change the pace instead. Run these "
                            "orders longer and start them earlier in the day.")
            elif np.isfinite(take) and np.isfinite(sp):
                half = _money_from(sp / 2 * take / 200, r["weight_pct"], NOT_)
                bits.append(f"It takes the price on offer for {take:.0f}% of "
                            f"its volume. Halving that is worth about "
                            f"{f_money(half)}.")
            _say(R, "rec_market", " ".join(bits), value=lost)

    # ---- what the spread explains ---------------------------------------
    if L and abs(L["r"]) >= 0.6:
        _say(F, "spread_fit",
             "**Your result tracks the spread.** Where the spread is wide you "
             "make money. Where it is tight you lose it. The turning point is "
             f"about **{L['breakeven']:.1f} bps** of spread.")
        if L["tight_names"] and np.isfinite(L["tight_bps"]):
            nm = ", ".join(f"**{x}**" for x in L["tight_names"])
            _say(F, "spread_tight",
                 f"{nm} have almost no spread to earn. That is "
                 f"{f_money(L['tight_share'] / 100 * NOT_)} of trading "
                 f"returning **{f_bps(L['tight_bps'])} bps**. Patience alone "
                 "will not fix them.")
        u = L["underperformer"]
        if u is not None:
            row = L["table"].loc[u]
            if row["vs_fit"] < -3:
                gap_ccy = abs(row["vs_fit"]) * row["notional_m"] * 1e6 / 1e4
                extra = ""
                if vc is not None and u in vc.index:
                    dk = _vc(u, "Dark")
                    extra = (f" It takes the price on offer for "
                             f"{_vc(u, 'Visible Take'):.0f}% of its volume"
                             + (" and never uses midpoint venues." if dk < 1
                                else f" and does {dk:.0f}% at the midpoint.")
                             + " Start there.")
                _say(R, "rec_outlier",
                     f"**{u} is your biggest untapped market.** On the "
                     "pattern your other markets follow, its "
                     f"{row['spread_bps']:.1f} bps spread should give about "
                     f"{row['predicted']:+.0f} bps. It gives {row['bps']:+.1f}. "
                     f"Closing that gap is worth roughly {f_money(gap_ccy)}."
                     + extra,
                     value=gap_ccy)

    # ---- the market worth copying ---------------------------------------
    # Ranking on result alone once crowned a market that takes the price on
    # offer 100% of the time, which is the opposite of the advice around it.
    if vc is not None and not vc.empty and "bps" in vc:
        model = vc[(vc["bps"] > 0) & (vc["weight_pct"] >= 3.0)
                   & (vc["Visible Post"] >= 40) & (vc["Visible Take"] <= 25)]
        if len(model):
            g = model["bps"].idxmax()
            _say(N, "model_market",
                 f"{g} is the market to copy: {_vc(g, 'Visible Post'):.0f}% "
                 f"resting in the queue, only {_vc(g, 'Visible Take'):.0f}% "
                 f"paying the spread, {_vc(g, 'Dark'):.0f}% at the midpoint, "
                 f"and {f_bps(vc.loc[g, 'bps'])} bps.")

    # ---- industry --------------------------------------------------------
    if not ind.empty:
        b_i, w_i = ind["contrib_bps"].idxmax(), ind["contrib_bps"].idxmin()
        if ind.loc[w_i, "contrib_bps"] < 0:
            w_not = ind.loc[w_i, "weight_pct"] / 100 * NOT_
            _say(F, "industry",
                 f"By sector, **{b_i}** is your best at "
                 f"{ind.loc[b_i, 'bps']:+.1f} bps. **{w_i}** is your worst: "
                 f"{f_money(w_not)} traded at {ind.loc[w_i, 'bps']:+.1f} bps, "
                 f"costing {f_money(abs(w_not * ind.loc[w_i, 'bps'] / 1e4))}.")
            if len(bad):
                _say(R, "rec_industry",
                     f"**Check whether {w_i} and {bad.index[0]} are the same "
                     "orders.** If they are, that is one problem to fix, not "
                     "two. Ask for the sector list split by market.",
                     value=0.0)

    # ---- order size ------------------------------------------------------
    # A band of forty orders can post a spectacular number by luck, so the
    # headline goes to the band that carries the money, not the best number.
    adv_q = adv[adv["n"] >= 20] if not adv.empty else adv
    if len(adv_q) >= 2:
        heavy = adv_q["weight_pct"].idxmax()
        best = adv_q["bps"].idxmax()
        hr = adv_q.loc[heavy]
        if hr["bps"] < adv_q.loc[best, "bps"]:
            tail = ("" if adv_q.loc[best, "n"] >= MIN_N_QUOTABLE
                    else f", though on only {int(adv_q.loc[best, 'n'])} orders")
            _say(F, "size",
                 "Your biggest group of orders is your weakest. Orders "
                 f"**{_band(heavy)}** were {f_money(hr['notional_m'] * 1e6)} "
                 f"and returned **{f_bps(hr['bps'])} bps**. Orders "
                 f"**{_band(best)}** returned "
                 f"{f_bps(adv_q.loc[best, 'bps'])} bps{tail}.")
        else:
            _say(F, "size",
                 f"Orders **{_band(heavy)}** were "
                 f"{f_money(hr['notional_m'] * 1e6)}, most of what you trade, "
                 f"and returned **{f_bps(hr['bps'])} bps**. That is also your "
                 "best band, so size is not working against you.")
        gap = hl["bps"] - hr["bps"]
        if gap > 0:
            prize = _money_from(gap, hr["weight_pct"], NOT_)
            _say(R, "rec_size",
                 f"**Give the small orders more time.** Orders "
                 f"{_band(heavy)} are {f_money(hr['notional_m'] * 1e6)} of "
                 "your trading and your weakest result. These are the easiest "
                 "orders you send, so there is "
                 "no reason to rush them. Lengthen the default end time. "
                 "Bringing this band up to your own average is worth about "
                 f"{f_money(prize)}.",
                 value=prize)
        thin = adv_q[adv_q["n"] < MIN_N_QUOTABLE]
        if len(thin):
            _say(N, "size_thin",
                 f"The {_and(thin.index)} bands hold fewer than "
                 f"{MIN_N_QUOTABLE} orders each. Treat those numbers as a "
                 "hint, not a result.")

    # ---- company size ----------------------------------------------------
    cap_q = cap[cap["n"] >= 20] if not cap.empty else cap
    if not cap_q.empty:
        w_c = cap_q["bps"].idxmin()
        if cap_q.loc[w_c, "bps"] < 0:
            lone = int((cap_q["bps"] < 0).sum()) == 1
            _say(F, "cap",
                 f"**{w_c}** is "
                 + ("the only company-size band losing money"
                    if lone else "your weakest company-size band")
                 + f". You traded "
                 f"{f_money(cap_q.loc[w_c, 'notional_m'] * 1e6)} there at "
                 f"{cap_q.loc[w_c, 'bps']:+.1f} bps, costing "
                 f"{f_money(abs(cap_q.loc[w_c, 'pnl_ccy']))}. At "
                 f"{cap_q.loc[w_c, 'adv_pct']:.1f}% of a day's volume these "
                 "are the easiest names you trade.")

    # ---- venue mix by algo ----------------------------------------------
    rows = {a: v for a, v in seg.items() if a != "TOTAL"}
    if rows:
        dark_tot = seg.get("TOTAL", {}).get("Dark", 0)
        weights = (algo["weight_pct"].to_dict() if not algo.empty else {})
        best_dark = max(rows, key=lambda a: rows[a].get("Dark", 0))
        if dark_tot >= 5:
            _say(F, "dark",
                 f"**{dark_tot:.0f}% of your volume trades at the midpoint**, "
                 f"nearly all of it through {best_dark} "
                 f"({rows[best_dark]['Dark']:.0f}%).")
            # only worth advising on an algo that carries real value: a change
            # to 0.2% of the book is not a recommendation
            cand = [a for a in rows
                    if rows[a].get("Dark", 0) < 1
                    and rows[a].get("Visible Take", 0) >= 50
                    and weights.get(a, 0) >= 2.0]
            if cand:
                a = max(cand, key=lambda x: weights.get(x, 0))
                prize = _money_from(hl["spread_bps"] / 2 *
                                    rows[a]["Visible Take"] / 100 * 0.5,
                                    weights.get(a, 0), NOT_)
                _say(R, "rec_algo_venue",
                     f"**Switch midpoint venues on for {a}.** It takes the "
                     f"price on offer for {rows[a]['Visible Take']:.0f}% of "
                     f"its volume and never tries the midpoint. That is "
                     f"{f_money(algo.loc[a, 'notional_m'] * 1e6)} of trading. "
                     f"This is one setting, and it is worth roughly "
                     f"{f_money(prize)}.",
                     value=prize)
        else:
            _say(F, "dark",
                 f"Only **{dark_tot:.1f}%** of your volume trades at the "
                 "midpoint. Venue is not what decides your result.")
            taker = max(rows, key=lambda a: rows[a].get("Visible Take", 0))
            if rows[taker].get("Visible Take", 0) >= 20 and \
                    weights.get(taker, 0) >= 2.0:
                tk = rows[taker]["Visible Take"]
                prize = _money_from(hl["spread_bps"] / 2 * tk / 100 * 0.5,
                                    weights.get(taker, 0), NOT_)
                _say(R, "rec_algo_venue",
                     f"**Get {taker} resting in the queue more and taking "
                     f"less.** It pays the spread on {tk:.0f}% of its volume. "
                     f"At a {hl['spread_bps']:.1f} bps spread that is about "
                     f"{hl['spread_bps'] / 2:.1f} bps each time. Halving it is "
                     f"worth roughly {f_money(prize)}.",
                     value=prize)
        auc = max(rows, key=lambda a: rows[a].get("Auction", 0))
        if rows[auc].get("Auction", 0) >= 50:
            _say(N, "auction",
                 f"{auc} puts {rows[auc]['Auction']:.0f}% of its volume into "
                 "the closing auction, so its result is about where the close "
                 "printed, not how the order was worked.")

    _say(N, "source",
         "Every number here comes from your published post-trade report, so "
         "it matches what you already hold.")

    # Biggest money first, so "what we would change" is honestly ordered.
    R.sort(key=lambda x: -(x["value"]
                           if x["value"] is not None and np.isfinite(x["value"])
                           else 0.0))
    return n


# --- slides ----------------------------------------------------------------
# Speaker notes are written for someone who has never seen a TCA report. Short
# sentences, plain words, and the jargon explained the first time it appears.

def s_cover(prs, ctx):
    s = new_slide(prs)
    hl = ctx["headline"]
    band = s.shapes.add_shape(1, Inches(0), Inches(0), SLIDE_W, Inches(2.45))
    band.fill.solid(); band.fill.fore_color.rgb = T_HEAD
    band.line.fill.background(); band.shadow.inherit = False
    # the azure band that sits under the charcoal panel on the house template
    stripe = s.shapes.add_shape(1, Inches(0), Inches(2.45), SLIDE_W, Inches(0.16))
    stripe.fill.solid(); stripe.fill.fore_color.rgb = T_ACCENT
    stripe.line.fill.background(); stripe.shadow.inherit = False

    box = s.shapes.add_textbox(MARGIN, Inches(0.66), Inches(11.5), Inches(1.5))
    tf = box.text_frame; tf.word_wrap = True
    _run(tf.paragraphs[0], "Transaction Cost Analysis", size=38, bold=True,
         color=T_WHITE)
    p2 = tf.add_paragraph()
    _run(p2, f"{CLIENT_NAME}  ·  {PERIOD_LABEL}", size=17, color=T_WHITE)

    ahead = np.isfinite(hl["bps"]) and hl["bps"] >= 0
    cards = [
        (f"{hl['orders']:,}", "orders", T_HEAD),
        (f_money(hl["notional"], 1), "traded", T_HEAD),
        (f_bps(hl["bps"]) + " bps", "vs benchmark", T_SAVE if ahead else T_COST),
        (f_money(abs(hl["pnl_ccy"])), "saved" if ahead else "cost",
         T_SAVE if ahead else T_COST),
        (f"{hl['markets']:.0f}", "markets", T_HEAD),
    ]
    w, gap = Inches(2.34), Inches(0.10)
    left = MARGIN
    for big, cap_, fill in cards:
        kpi(s, left, Inches(2.85), w, Inches(1.15), big, cap_, fill=fill,
            big_size=22 if len(big) <= 10 else 17)
        left += w + gap

    lines = []
    if np.isfinite(hl.get("spread_bps", np.nan)):
        lines.append(f"A typical order was **{hl['adv_pct']:.1f}%** of a normal "
                     "day's trading in that share. The average gap between the "
                     "buy and sell price was "
                     f"**{hl['spread_bps']:.1f} bps**.")
    if np.isfinite(hl.get("dark_share", np.nan)):
        lines.append(f"**{hl['dark_share']:.1f}%** of your volume traded at the "
                     "midpoint, in private venues.")
    lines.append("Every number comes from your published post-trade report.")
    bullets(s, lines, MARGIN, Inches(4.35), Inches(11.9), Inches(1.5), size=12.5)
    note(s, "1 bps = 0.01%. A plus means you saved money, a minus means it "
            "cost you. The summary and the advice are the last two slides.")
    notes(s, f"""WHAT THIS SLIDE SHOWS
The headline numbers for the whole period, in one place.

THE WORDS
"bps" means basis points. 1 bps is 0.01%. 100 bps is 1%. We use it because
these differences are small next to the amounts traded.

"Benchmark" is the fair reference price we measure you against. For most of
your orders it is the average price the rest of the market paid while your
order was running. Beat it and you saved money.

THE NUMBER THAT MATTERS
{_cover_line(hl)}

Everything here comes from your own post-trade report. Every number should
match what you already hold.""")
    footer(s)


def s_summary(prs, ctx, nar):
    """What we found.

    It sits second from last, next to the advice, because a reader who has
    just been through the evidence wants the wrap-up and the action together.
    The headline itself is already on the cover, so nothing is lost by
    holding this back.
    """
    s = new_slide(prs)
    hl = ctx["headline"]
    title(s, "Summary")
    ahead = hl["bps"] >= 0
    strapline(s, (f"{f_bps(hl['bps'])} bps "
                  + ("better than benchmark, worth " if ahead
                     else "worse than benchmark, costing ")
                  + f"{f_money(abs(hl['pnl_ccy']))}. "
                  + "The next slide says what to do about it."))

    left = pick(nar["findings"], "headline", "algo", "market_good",
                "market_bad")
    right = pick(nar["findings"], "side_worst", "side_best", "size",
                 "cap", "spread_tight")[:4]

    w = Inches(5.9)
    for x, head, items in [(MARGIN, "THE RESULT", left[:4]),
                           (Inches(6.85), "WHERE IT COMES FROM", right)]:
        box = s.shapes.add_textbox(x, Inches(1.62), w, Inches(0.32))
        _run(box.text_frame.paragraphs[0], head, size=11, bold=True,
             color=T_HEAD)
        bullets(s, items, x, Inches(1.98), w, Inches(4.4), size=11.5)
    note(s, "This deck cuts the same money several ways — by country, by "
            "sector, by order size, by company size. The cuts overlap, so do "
            "not add them together.")
    notes(s, """WHAT THIS SLIDE SHOWS
Everything the deck found, on one page.

HOW TO READ IT
Left: the result and which algo and which country produced it.
Right: where inside those countries the money actually sits.

Every figure here is measured, and every one has a chart earlier in the deck
behind it.

TWO NUMBERS, THAT IS ALL
Money is size and effect: what you traded, and what it made or cost.
Bps is quality per order: how well each order did against its benchmark.
A country can be large and mediocre, or small and terrible. The money tells
you which one is worth your time.

THE QUESTION THIS INVITES
"You told me one country cost us X, and one sector cost us Y. Which is it?"
Both. They are the same money sliced two different ways, and the same orders
appear in each. Every slide adds to the headline on its own. The slides must
not be added to each other.

WHAT COMES NEXT
The final slide turns this into four changes, ordered by money.""")
    footer(s)


def s_algo(prs, ctx):
    s = new_slide(prs)
    title(s, "By Algorithm")
    strapline(s, "How each algo did against its own benchmark, and how much "
                 "it moved your total.")
    picture(s, ctx["charts"].get("algo"), MARGIN, Inches(1.55),
            width=Inches(7.4))
    algo = ctx["algo_table"]
    if not algo.empty:
        data = [["Algo", "Orders", "Traded", "Spread", "Result",
                 "Money"]]
        colored, thin = {}, False
        for i, (idx, r) in enumerate(algo.iterrows(), start=1):
            tag = ""
            if r["n"] < MIN_N_QUOTABLE:
                tag, thin = " *", True
            data.append([f"{idx}{tag}", f"{int(r['n']):,}",
                         f_money(r["notional_m"] * 1e6),
                         f"{r['spread_bps']:.1f}",
                         f_bps(r["bps"]),
                         f_money(r["pnl_ccy"])])
            colored[(i, 4)] = "save" if r["bps"] >= 0 else "cost"
            colored[(i, 5)] = "save" if r["pnl_ccy"] >= 0 else "cost"
        table(s, data, Inches(8.15), Inches(1.70), Inches(4.6),
              row_h=Inches(0.30), colored=colored)
        y = Inches(1.74) + Inches(0.30) * len(data)
        if thin:
            box = s.shapes.add_textbox(Inches(8.15), y, Inches(4.6), Inches(0.3))
            _run(box.text_frame.paragraphs[0],
                 f"* under {MIN_N_QUOTABLE} orders — not a reliable number",
                 size=8, color=T_MUTED)
            y = y + Inches(0.30)
        cut = [r for r in (ALGO_REPORT or [])
               if ALGOS_STUDIED and r[0] not in ALGOS_STUDIED]
        if cut:
            box = s.shapes.add_textbox(Inches(8.15), y, Inches(4.6),
                                       Inches(0.5))
            box.text_frame.word_wrap = True
            _run(box.text_frame.paragraphs[0],
                 f"{_and([r[0] for r in cut])} are not reviewed here: "
                 f"{sum(r[1] for r in cut):,} orders, "
                 f"{f_money(sum(r[2] for r in cut))}, "
                 f"{f_money(sum(r[2] * r[8] / 1e4 for r in cut))} between "
                 "them. That is the gap to the cover figure.",
                 size=8, color=T_MUTED)
            y = y + Inches(0.42)
        zero = algo[algo["bps"].abs() < 1e-9]
        if len(zero):
            box = s.shapes.add_textbox(Inches(8.15), y + Inches(0.15),
                                       Inches(4.6), Inches(1.6))
            tf = box.text_frame; tf.word_wrap = True
            _run(tf.paragraphs[0],
                 f"{', '.join(zero.index)} trades in the closing auction and is "
                 "scored against the closing price. It is compared with itself, "
                 "so it always scores zero. Ask for it to be scored against the "
                 "arrival price instead.",
                 size=9.5, color=T_MUTED)
    notes(s, """WHAT THIS SLIDE SHOWS
How each of your algos performed.

THE WORDS
An "algo" is the automated strategy that works your order in the market. VWAP
spreads an order through the day to land near the day's average price. IIS
aims to trade at the closing price.

"Effect on total" is the column that matters. An algo only moves your number
if it does well AND handles a lot of your volume. Brilliant on 1% of the flow
changes nothing.

WATCH OUT FOR THIS
A score of exactly 0.0 is not a perfect score. It happens when an algo trades
in the closing auction and is then measured against the closing price. It is
being compared with itself, so the answer is always zero. It tells us nothing
about how well that algo worked. Ask for it to be measured against the price
when the order arrived instead.

A star means too few orders to trust. A handful of trades can produce a
spectacular number by luck.""")
    footer(s)


def s_market(prs, ctx, nar):
    s = new_slide(prs)
    title(s, "By Market")
    strapline(s, "A country moves your total only if you did well or badly "
                 "there and traded a lot there.")
    picture(s, ctx["charts"].get("attribution"), MARGIN, Inches(1.52),
            width=Inches(7.5))
    lines = pick(nar["findings"], "market_good", "market_bad")
    mkt = ctx["market_table"]
    if not mkt.empty:
        top3 = mkt.nlargest(3, "notional_m")
        lines.append("Your three biggest markets are "
                     + _and([f"**{x}**" for x in top3.index])
                     + f". Together that is "
                     f"{f_money(top3['notional_m'].sum() * 1e6)} of the "
                     f"{f_money(ctx['headline']['notional'])} you traded.")
    bullets(s, lines, Inches(8.20), Inches(1.72), Inches(4.55), Inches(4.3),
            size=11)
    notes(s, """WHAT THIS SLIDE SHOWS
Which countries helped your result and which hurt it.

HOW TO READ IT
Each bar is the money that country made or lost. That is how well you traded
there multiplied by how much you traded there.

Both halves matter. A country can be bad but tiny, and barely move your total.
Another can be slightly bad but huge, and dominate it. The label under each
name shows both: how much you traded, and how well.

Blue bars are money made. Red bars are money lost. Add all the bars together
and you get the figure on the cover.

WHY IT MATTERS
It shows where your attention is worth spending.""")
    footer(s)


def s_industry(prs, ctx, nar):
    if not ctx["charts"].get("industry"):
        return
    s = new_slide(prs)
    title(s, "By Industry")
    strapline(s, "The same result, grouped by sector instead of by country.")
    picture(s, ctx["charts"]["industry"], MARGIN, Inches(1.52),
            width=Inches(7.3))
    lines = pick(nar["findings"], "industry")
    lines += pick(nar["recs"], "rec_industry")
    lines.append("Read this against the market slide. If the same names appear "
                 "in both, the problem is the stock, not the country.")
    bullets(s, lines, Inches(8.05), Inches(1.72), Inches(4.7), Inches(4.3),
            size=10.5)
    note(s, "Sector figures come from the published industry breakdown.")
    notes(s, """WHAT THIS SLIDE SHOWS
The same money as the last slide, grouped by the kind of company rather than
the country.

HOW TO READ IT
Exactly like the country slide. Result multiplied by how much you traded.

WHY IT MATTERS
This is a cross-check. If a problem shows up in one country and in one sector,
those are probably the same orders described two ways. One thing to fix, not
two.

If a sector looks bad but is spread evenly across countries, the problem is
the shares themselves, not where or how you traded them.""")
    footer(s)


def s_venue_country(prs, ctx, nar):
    if not ctx["charts"].get("venue_country"):
        return
    s = new_slide(prs)
    title(s, "How Each Market Trades")
    strapline(s, "Where your shares actually change hands. This is what "
                 "decides which fix applies.")
    picture(s, ctx["charts"]["venue_country"], MARGIN, Inches(1.50),
            width=Inches(7.3))
    lines = pick(nar["findings"], "venue_take", "venue_passive")
    # when every taking market already has its own recommendation, show that
    # one here rather than leaving the slide with no action on it
    lines += pick(nar["recs"], "rec_post") or pick(nar["recs"], "rec_market")[:1]
    bullets(s, lines, Inches(8.05), Inches(1.70), Inches(4.7), Inches(4.6),
            size=10.5)
    note(s, "Each bar adds to 100% of what you traded in that market. "
            "Markets below 1% of your trading are left out.")
    notes(s, """WHAT THIS SLIDE SHOWS
Where your shares actually change hands in each country. There are four ways,
and the difference between them is money.

THE FOUR WAYS
AUCTION (blue). The single batch trade at the open or the close. Everyone
trades at one price.

RESTING IN THE QUEUE (green). You leave your order in the market and wait for
someone to come to you. You earn the spread for being patient. The risk is
that you do not get filled.

PAYING THE SPREAD (orange). You take the price on offer right now. It is
immediate and certain, but you give away about half the spread each time.

MIDPOINT (purple). A private venue where both sides meet in the middle. No
spread paid and none earned. You only trade if someone wants the other side.

HOW TO READ IT
Each bar adds up to 100% of what you traded in that country. More green and
purple is generally better. More orange means you are giving the spread away.

WHY THIS SLIDE DECIDES THE FIX
A country that loses money and pays the spread a lot can be fixed with a
setting. That is fully in your control.

A country that loses money but already rests in the queue and already uses
midpoint venues cannot be fixed that way. Its problem is pace: how long the
orders run, and what time of day they run. Telling that desk to post more
would waste their time.""")
    footer(s)


def s_side(prs, ctx, nar):
    if not ctx["charts"].get("side"):
        return
    s = new_slide(prs)
    title(s, "By Market and Side")
    strapline(s, "The same money, split into what you bought and what you "
                 "sold. This is the sharpest view in the deck.")
    picture(s, ctx["charts"]["side"], MARGIN, Inches(1.50),
            width=Inches(7.2))
    lines = pick(nar["findings"], "side_worst", "side_best", "side_widest")
    lines += pick(nar["recs"], "rec_market")[:1]
    bullets(s, lines, Inches(7.95), Inches(1.70), Inches(4.8), Inches(4.5),
            size=10.5)
    note(s, "Only groups worth more than 1% of what you traded are shown. "
            "Each bar is the money that group made or lost, so the bars add "
            "back to the headline.")
    notes(s, """WHAT THIS SLIDE SHOWS
The same money again, now split by whether you were buying or selling.

HOW TO READ IT
Each row is one country and one direction. "India - BUY" means every share you
bought in India over the period.

The bars are in dollars, not percentages. The longest bar really is the
largest amount of money.

WHY THIS IS THE SHARPEST SLIDE
Buying and selling in the same country often go in opposite directions. They
cancel each other out on every other slide. A country can look calm overall
while hiding a large loss on one side and a large gain on the other.

Splitting them lets you point at one specific group of orders instead of a
whole country.

CHECK THIS BEFORE ACTING
A buy/sell gap that lasts for months is often about market direction, not
trading. If prices rose all period and you were mostly buying, buys will look
worse. Worth confirming first.""")
    footer(s)


def s_spread(prs, ctx, nar):
    if not ctx["charts"].get("spread"):
        return
    s = new_slide(prs)
    title(s, "What Explains the Result")
    strapline(s, "Where the buy/sell gap is wide you make money. Where it is "
                 "tight you lose it.")
    picture(s, ctx["charts"]["spread"], MARGIN, Inches(1.52),
            width=Inches(7.5))
    lines = pick(nar["findings"], "spread_fit", "spread_tight")
    lines += pick(nar["recs"], "rec_outlier")
    bullets(s, lines, Inches(8.20), Inches(1.72), Inches(4.55), Inches(4.4),
            size=10.5)
    note(s, "Countries where you placed at least 50 orders. The line is the "
            "overall trend. Each bubble is sized by how much you trade there.")
    notes(s, """WHAT THIS SLIDE SHOWS
Why some countries work for you and others do not.

THE WORDS
The "spread" is the gap between the best buying price and the best selling
price. Wait, and you earn that gap. Hurry, and you pay it.

A wide spread means there is a lot to gain by being patient. A tight spread
means there is almost nothing there either way.

HOW TO READ IT
Each bubble is a country. Further left means a tighter spread. Higher up means
a better result. Bubble size is how much you trade there.

The dotted line is roughly where the two balance out. Countries to the left of
it have very little spread to earn.

WHY IT MATTERS
Your results line up with the spread almost exactly. That tells us your
strategy makes its money by earning the spread. Where the spread is tiny, that
strategy has nothing to work with. Those countries do not need more patience.
They need a different approach. The next slide shows which.""")
    footer(s)


def s_size(prs, ctx, nar):
    s = new_slide(prs)
    title(s, "By Order Size")
    strapline(s, "How you do on small orders against large ones.")
    picture(s, ctx["charts"].get("adv"), MARGIN, Inches(1.52),
            width=Inches(7.1))
    lines = pick(nar["findings"], "size", "cap")
    lines += pick(nar["recs"], "rec_size")
    bullets(s, lines, Inches(8.05), Inches(1.72), Inches(4.7), Inches(2.4),
            size=10.5)
    picture(s, ctx["charts"].get("cap"), Inches(8.05), Inches(4.35),
            width=Inches(4.75))
    notes(s, """WHAT THIS SLIDE SHOWS
Whether your small orders or your big orders get the better result.

THE WORDS
"% of a day's volume" is how big your order is next to what that share
normally trades in a whole day. An order that is 1% of a day is small and easy
to hide. An order that is 10% is large, and the market will notice you.

The smaller chart splits the same money by company size instead. The largest
companies are usually the easiest to trade.

HOW TO READ IT
You would normally expect big orders to do worse. You are asking the market to
absorb more, so the price moves against you.

WHY IT MATTERS
If your small, easy orders do worse than your large, hard ones, something is
backwards. It usually means the easy flow is being pushed through faster than
it needs to be. Easy flow is normally most of your money, so it is worth
getting right.

Bands with very few orders are flagged. A handful of trades can produce an
impressive number purely by luck.""")
    footer(s)


def s_venue(prs, ctx, nar):
    s = new_slide(prs)
    hl = ctx["headline"]
    title(s, "How Each Algo Trades")
    dark = hl.get("dark_share", np.nan)
    strapline(s, (f"{dark:.0f}% of your volume trades at the midpoint."
                  if np.isfinite(dark) and dark >= 5
                  else "Auction, resting in the queue, or paying the spread."))
    picture(s, ctx["charts"].get("venue"), MARGIN, Inches(1.55),
            width=Inches(7.5))
    seg = VENUE_SEGMENTS_REPORT or {}
    lines = []
    for a in [x for x in ALGO_ORDER if x in seg] or \
             [x for x in seg if x != "TOTAL"]:
        parts = []
        for k, lbl in [("Auction", "in the auction"),
                       ("Visible Post", "resting in the queue"),
                       ("Visible Take", "paying the spread"),
                       ("Dark", "at the midpoint")]:
            v = seg[a].get(k, 0)
            if v >= 5:
                parts.append(f"**{v:.0f}%** {lbl}")
        if parts:
            lines.append(f"**{a}** — " + ", ".join(parts) + ".")
    lines += pick(nar["recs"], "rec_algo_venue")
    bullets(s, lines, Inches(8.20), Inches(1.72), Inches(4.55), Inches(4.3),
            size=10.5)
    note(s, "Venue split comes from the published venue segment breakdown, "
            "as a share of what you traded.")
    notes(s, """WHAT THIS SLIDE SHOWS
The same four ways of trading as the market slide, split by algo instead of by
country.

A REMINDER OF THE FOUR
Auction. The batch trade at the open or the close.
Resting in the queue. You wait, others come to you, and you earn the spread.
Paying the spread. You take what is on offer now, and you give the spread up.
Midpoint. A private venue in the middle, where you neither pay nor earn it.

HOW TO READ IT
Each bar is 100% of what that algo traded.

WHY IT MATTERS
It shows whether each algo is using all the tools it has. An algo that only
ever pays the spread is leaving the cheap options untouched. That is normally
a setting, not a limit.""")
    footer(s)


def s_advice(prs, ctx, nar):
    s = new_slide(prs)
    title(s, "What To Change")
    strapline(s, "Biggest money first.")
    for i, r in enumerate(texts(nar["recs"])[:4]):
        y = Inches(1.55 + i * 1.05)
        num = s.shapes.add_shape(1, MARGIN, y, Inches(0.44), Inches(0.44))
        num.fill.solid(); num.fill.fore_color.rgb = T_HEAD
        num.line.fill.background(); num.shadow.inherit = False
        tf = num.text_frame; tf.margin_top = Inches(0.02)
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        _run(tf.paragraphs[0], str(i + 1), size=15, bold=True, color=T_WHITE)
        bullets(s, [r], Inches(1.22), y - Inches(0.05), Inches(11.4),
                Inches(1.0), size=11.5, marker="")

    if nar["notes"]:
        box = s.shapes.add_textbox(MARGIN, Inches(5.82), Inches(12.1),
                                   Inches(0.3))
        _run(box.text_frame.paragraphs[0], "WORTH KNOWING", size=10, bold=True,
             color=T_MUTED)
        bullets(s, ["~" + x for x in texts(nar["notes"])[:3]], MARGIN,
                Inches(6.10), Inches(12.1), Inches(0.85), size=9, gap=3)
    notes(s, """WHAT THIS SLIDE SHOWS
What we would change, in order of how much money is at stake.

HOW TO READ IT
Each item names a change and roughly what it is worth. Start with number one.

We have avoided "look into X". Everything here can be acted on.

THE SMALL PRINT
"Worth knowing" lists the things that could change these conclusions. Please
read it before acting. Where a number has an innocent explanation, we say so
rather than presenting it as settled.

WHAT WE WOULD DO NEXT
Run the same analysis next period and see whether the number moved. The
comparison that counts is against your own result on the same kind of flow.""")
    footer(s)


def build_simple_deck(ctx: dict, nar: dict, out: Path) -> Path:
    _PAGE[0] = 0
    prs = Presentation()
    prs.slide_width, prs.slide_height = SLIDE_W, SLIDE_H
    s_cover(prs, ctx)
    s_algo(prs, ctx)
    s_market(prs, ctx, nar)
    s_industry(prs, ctx, nar)
    s_side(prs, ctx, nar)
    s_spread(prs, ctx, nar)
    s_venue_country(prs, ctx, nar)
    s_size(prs, ctx, nar)
    s_venue(prs, ctx, nar)
    s_summary(prs, ctx, nar)
    s_advice(prs, ctx, nar)
    prs.save(str(out))
    log(f"  saved {out.name}  ({len(prs.slides._sldIdLst)} slides)")
    return out


# ===========================================================================
# 11. ORCHESTRATION
# ===========================================================================

def analyse(df: pd.DataFrame) -> dict:
    log("")
    log("-" * 74)
    log("ANALYSIS")
    log("-" * 74)
    df = add_control_buckets(df)
    dark_algos = derive_dark_algos(df)
    dpooled = pool_small_markets(df)

    rep_tot = REPORT.get("totals") or {}
    rep_bps = rep_tot.get("bps")

    # headline and every aggregate breakdown come from the report; the order
    # file is used below only for the dark work and the outlier list
    hl = headline_from_report(rep_tot) if rep_tot else headline(df)
    if rep_tot and "market" in df:
        hl["markets"] = int(df["market"].nunique())
        hl["symbols"] = int(df["symbol"].nunique()) if "symbol" in df else np.nan
        hl["dark_share"] = 100 * df["dark_value"].sum() / df["notional"].sum()
        hl["dark_order_share"] = 100 * float(df["has_dark"].mean())
        hl["dark_share_eligible"] = _dark_share_eligible(df)
    elif rep_tot:
        for k in ("markets", "symbols", "dark_share", "dark_order_share",
                  "dark_share_eligible"):
            hl.setdefault(k, np.nan)

    ctx = {
        "df": df,
        "headline": hl,
        "algo_table": (algo_table_from_report(ALGO_REPORT, ALGOS_STUDIED)
                       if ALGO_REPORT
                       else group_table(df, "algo", order=ALGO_ORDER)),
        "algo_computed": group_table(df, "algo", order=ALGO_ORDER),
        "algo_reconcile": reconcile_algo(df, ALGO_REPORT),
        "market_table": (table_from_report(REPORT.get("country"), "market",
                                           rep_bps)
                         if REPORT.get("country")
                         else (group_table(dpooled, "market_grp", min_n=1)
                               if "market_grp" in dpooled else pd.DataFrame())),
        "adv_table": (table_from_report(REPORT.get("adv"), "order size",
                                        rep_bps)
                      if REPORT.get("adv")
                      else (group_table(df, "adv_bucket", order=ADV_LABELS)
                            if "adv_bucket" in df else pd.DataFrame())),
        "cap_table": table_from_report(REPORT.get("marketcap"), "market cap",
                                       rep_bps),
        "side_table": table_from_report(REPORT.get("side"), "side", rep_bps),
        "benchmark_matrix": benchmark_matrix(df),
        "venue_table": venue_mix(df),
        "worst_orders": worst_orders(df),
        "dark_algos": dark_algos,
        "dark_markets": dark_markets_present(df),
    }
    if not ctx["market_table"].empty:
        ctx["market_table"].index.name = "market"

    if not DARK_STORY:
        log("  dark story is off for this client - skipping the dark section")
        empty = df.iloc[0:0].copy()
        ctx.update({"dark_df": empty, "dark_groups": ["No dark", "Any dark"],
                    "dark_split_note": "", "dark_zero_vs_any": pd.DataFrame(),
                    "dark_ladder": pd.DataFrame(),
                    "dark_completion": pd.DataFrame(),
                    "dark_by_market": pd.DataFrame(), "regression": None})
        return ctx

    d = dark_frame(df, dark_algos)
    ctx["dark_df"] = d
    ctx["dark_groups"] = d.attrs.get("groups", ["No dark", "Any dark"]) \
        if not d.empty else ["No dark", "Any dark"]
    ctx["dark_split_note"] = d.attrs.get(
        "split_note", "No dark-capable orders found.") if not d.empty \
        else "No dark-capable orders found in this dataset."
    ctx["dark_ok"] = dark_sanity(d, (REPORT.get("totals") or {}).get("bps"))
    ctx["dark_zero_vs_any"] = dark_zero_vs_any(d)
    ctx["dark_ladder"] = dark_ladder(d)
    ctx["dark_completion"] = dark_completion(d)
    ctx["dark_by_market"] = dark_by_market(d)
    ctx["regression"] = dark_regression(d)
    return ctx


def make_charts(ctx: dict, outdir: Path) -> dict:
    log("")
    log("-" * 74)
    log("CHARTS")
    log("-" * 74)
    apply_style()
    c = {}
    hl = ctx["headline"]
    if not ctx["algo_table"].empty:
        c["algo"] = chart_algo(ctx["algo_table"], outdir, hl)
    if not ctx["market_table"].empty:
        c["attribution"] = chart_attribution(ctx["market_table"], outdir,
                                             by="market",
                                             total_bps=hl["bps"])
    c["notional"] = chart_notional_by_country(
        ctx["df"], outdir, ctx.get("dark_markets"),
        report=ctx["market_table"] if REPORT.get("country") else None)
    if not ctx["adv_table"].empty:
        c["adv"] = chart_adv(ctx["adv_table"], outdir)
    if INDUSTRY_REPORT:
        c["industry"] = chart_industry(INDUSTRY_REPORT, outdir, hl["bps"])
    c["venue"] = chart_venue(outdir, ctx["venue_table"])
    if not ctx["dark_zero_vs_any"].empty:
        c["dark_zero"] = chart_dark_zero_vs_any(
            ctx["dark_zero_vs_any"], outdir, ctx["dark_split_note"])
        c["dark_controlled"] = chart_dark_controlled(
            ctx["dark_df"], DARK_CONTROLS, outdir, ctx["dark_groups"])
    if not ctx["dark_by_market"].empty:
        c["dark_market"] = chart_dark_by_market(ctx["dark_by_market"], outdir,
                                                ctx["dark_groups"])
    if not ctx["dark_completion"].empty:
        c["dark_completion"] = chart_dark_completion(ctx["dark_completion"],
                                                     outdir)
    ctx["charts"] = c
    return c


def write_tables(ctx: dict, path: Path) -> None:
    # --simple passes a ready-made {sheet: frame} mapping
    if "headline" in ctx and "algo_table" not in ctx:
        with pd.ExcelWriter(path, engine="openpyxl") as xl:
            for name, t in ctx.items():
                if t is not None and hasattr(t, "empty") and not t.empty:
                    t.to_excel(xl, sheet_name=name[:31])
        log(f"  saved {path.name}")
        return
    sheets = {
        "headline": pd.DataFrame([ctx["headline"]]).T.rename(columns={0: "value"}),
        "by_algo": ctx["algo_table"],
        "by_algo_computed": ctx.get("algo_computed", pd.DataFrame()),
        "algo_reconciliation": ctx.get("algo_reconcile", pd.DataFrame()),
        "by_market": ctx["market_table"],
        "by_adv_bucket": ctx["adv_table"],
        "by_market_cap": ctx.get("cap_table", pd.DataFrame()),
        "by_side": ctx.get("side_table", pd.DataFrame()),
        "benchmark_matrix": ctx["benchmark_matrix"],
        "venue_mix": ctx["venue_table"],
        "dark_zero_vs_any": ctx["dark_zero_vs_any"],
        "dark_ladder": ctx["dark_ladder"],
        "dark_completion": ctx["dark_completion"],
        "dark_by_market": ctx["dark_by_market"],
        "worst_orders": ctx["worst_orders"],
        "reference_check": ctx.get("reference_check", pd.DataFrame()),
    }
    for c in DARK_CONTROLS:
        t = dark_controlled(ctx["dark_df"], c)
        if not t.empty:
            sheets[f"dark_by_{c}"[:31]] = t
    with pd.ExcelWriter(path, engine="openpyxl") as xl:
        for name, t in sheets.items():
            if t is None or (hasattr(t, "empty") and t.empty):
                continue
            t.to_excel(xl, sheet_name=name[:31])
    log(f"  saved {path.name}  ({len(sheets)} sheets)")


def probe(path: Path) -> int:
    """Dump the shape of a data file and stop.

    Run this first on a new export and send me the output: it is enough to
    confirm the column mapping, the sign convention, the algo and market
    spellings and the dark coverage, without building a single slide.
    """
    log("=" * 74)
    log(f"PROBE  {path}")
    log("=" * 74)
    if path.suffix.lower() in (".xlsx", ".xlsm", ".xls"):
        xl = pd.ExcelFile(path)
        log(f"  sheets: {', '.join(map(str, xl.sheet_names))}")
        raw0 = pd.read_excel(path, sheet_name=SHEET, header=None, nrows=4)
        log("")
        log("  first 4 rows, unheadered (check which row is the header):")
        for i in range(len(raw0)):
            cells = [str(c) for c in raw0.iloc[i].tolist()[:9]]
            log(f"    row {i}: " + " | ".join(cells))
    log("")
    df_raw = (pd.read_excel(path, sheet_name=SHEET, header=HEADER_ROW)
              if path.suffix.lower() in (".xlsx", ".xlsm", ".xls")
              else pd.read_csv(path, header=HEADER_ROW, encoding="utf-8-sig"))
    log(f"  parsed with HEADER_ROW={HEADER_ROW}: {len(df_raw):,} rows x "
        f"{len(df_raw.columns)} cols")
    log("")
    log("  COLUMNS PRESENT")
    for c in df_raw.columns:
        nn = int(df_raw[c].notna().sum())
        # kept out of the f-string: a conditional spanning lines inside a
        # replacement field only parses on Python 3.12+
        sample = str(df_raw[c].dropna().iloc[0])[:28] if nn else "-"
        log(f"    {str(c):<24} {str(df_raw[c].dtype):<10} "
            f"{nn:>7,} non-null  e.g. {sample}")
    cols, missing = resolve_columns(df_raw)
    log("")
    log(f"  mapped {len(cols)}/{len(COLUMNS)} known fields")
    if missing:
        log(f"  NOT MAPPED: {', '.join(missing)}")

    df = load_orders(path)
    log("")
    log("  ALGOS")
    for a, n in df["algo"].value_counts().items():
        v = df.loc[df.algo == a, "notional"].sum()
        d = 100 * df.loc[df.algo == a, "has_dark"].mean()
        log(f"    {a:<10} {n:>6,} orders  {v/1e6:>10,.0f}m  "
            f"{d:>5.1f}% of orders with a dark fill")
    if "market" in df:
        log("")
        log("  MARKETS")
        g = df.groupby("market").agg(n=("notional", "size"),
                                     v=("notional", "sum"),
                                     dk=("has_dark", "mean"))
        for m, r in g.sort_values("v", ascending=False).iterrows():
            log(f"    {str(m):<16} {int(r.n):>6,} orders  {r.v/1e6:>10,.0f}m  "
                f"{100*r.dk:>5.1f}% with dark")
    log("")
    log("  SIGN CHECK (positive should mean a saving)")
    for c in ["slip_arrival", "slip_pvwap", "slip_vwap", "slip_close"]:
        if c in df:
            v = df[c].dropna()
            if len(v):
                log(f"    {c:<14} mean {v.mean():+8.2f}  median "
                    f"{v.median():+8.2f}  {100*(v>0).mean():>5.1f}% positive  "
                    f"{v.isna().mean()*100:.0f}% missing")
    log("")
    sanity_report(df)
    reference_check(df)
    log("")
    log("Probe complete - nothing was built. Send this output back if the")
    log("mapping or the sign convention looks wrong.")
    return 0


def main(argv=None) -> int:
    global HEADER_ROW, SHEET, CLIENT_NAME, CLIENT_CODE, PERIOD_LABEL
    global ALGOS_STUDIED, DARK_STORY, DARK_MARKETS, VENUE_SEGMENTS_REPORT
    global EXCLUDE_MARKETS
    global INDUSTRY_REPORT, ALGO_REPORT, REPORT, REFERENCE, ALGO_ORDER
    ap = argparse.ArgumentParser(
        description="Build the client TCA deck from the order export.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog=textwrap.dedent("""\
            examples:
              python tca_deck.py --client KIC --data KIC/symboldetailstable_full_KIC.xlsx
              python tca_deck.py --client NPS --data NPS/symboldetailstable_full_NPS.xlsx
              python tca_deck.py --client NPS --data NPS/orders.xlsx --probe
              python tca_deck.py --sample
            """))
    ap.add_argument("--client", default=DEFAULT_CLIENT,
                    choices=sorted(CLIENTS),
                    help=f"which client config to use (default {DEFAULT_CLIENT})")
    ap.add_argument("--data", type=Path,
                    help="the order export (.xlsx or .csv)")
    ap.add_argument("--out", type=Path, default=None,
                    help="output directory (default: <client>/output)")
    ap.add_argument("--probe", action="store_true",
                    help="inspect the data file and stop - no charts, no deck")
    ap.add_argument("--simple", action="store_true",
                    help="short client deck from the published report figures "
                         "only - no order file needed")
    ap.add_argument("--sample", action="store_true",
                    help="generate synthetic KIC-shaped data and run on it")
    ap.add_argument("--header-row", type=int, default=None,
                    help=f"0-indexed header row (default {HEADER_ROW}, "
                         "i.e. row 2 of the sheet)")
    ap.add_argument("--sheet", default=None,
                    help="worksheet name or index (default: first)")
    ap.add_argument("--no-deck", action="store_true",
                    help="tables and charts only")
    args = ap.parse_args(argv)

    if args.header_row is not None:
        HEADER_ROW = args.header_row
    if args.sheet is not None:
        SHEET = int(args.sheet) if str(args.sheet).isdigit() else args.sheet

    cfg = CLIENTS[args.client]
    CLIENT_NAME     = cfg["name"]
    CLIENT_CODE     = cfg["code"]
    PERIOD_LABEL    = cfg["period"] or ""
    ALGOS_STUDIED   = cfg["algos"]
    ALGO_ORDER      = cfg["algo_order"]
    DARK_STORY      = cfg["dark_story"]
    DARK_MARKETS    = cfg["dark_markets"]
    EXCLUDE_MARKETS = list(cfg.get("exclude_markets") or [])
    VENUE_SEGMENTS_REPORT = cfg["venue_segments"]
    INDUSTRY_REPORT = cfg["industry"]
    ALGO_REPORT     = cfg.get("algo_report")
    REPORT          = {k: cfg.get(k) for k in
                       ("country", "marketcap", "adv", "side", "country_side",
                        "venue_country", "totals")}
    REFERENCE       = cfg["reference"]

    out = args.out or Path(CLIENT_NAME) / "output"
    out.mkdir(parents=True, exist_ok=True)
    charts_dir = out / "charts"

    log("=" * 74)
    log(f"{CLIENT_NAME} TCA" + (f"  —  {PERIOD_LABEL}" if PERIOD_LABEL else ""))
    log("=" * 74)

    if args.simple:
        ctx = analyse_report_only()
        make_simple_charts(ctx, charts_dir)
        nar = build_simple_narrative(ctx)
        log("")
        log("-" * 74)
        log("DECK")
        log("-" * 74)
        build_simple_deck(ctx, nar, out / f"{CLIENT_NAME}_TCA.pptx")
        write_tables({"headline": pd.DataFrame([ctx["headline"]]).T
                      .rename(columns={0: "value"}),
                      "by_algo": ctx["algo_table"],
                      "by_market": ctx["market_table"],
                      "by_adv_bucket": ctx["adv_table"],
                      "by_market_cap": ctx["cap_table"],
                      "by_side": ctx["side_table"],
                      "venue_by_market": (ctx["venue_country"]
                                          if ctx.get("venue_country") is not None
                                          else pd.DataFrame()),
                      "by_market_side": (ctx["side_lens"]["table"]
                                         if ctx.get("side_lens")
                                         else pd.DataFrame())},
                     out / "tables.xlsx")
        (out / "run_log.txt").write_text("\n".join(LOG), encoding="utf-8")
        log("")
        log("Done. Everything is in " + str(out.resolve()))
        log("  Every figure comes from the published report - no order file "
            "was read.")
        return 0

    if args.sample:
        IS_SAMPLE[0] = True
        data_path = out / f"sample_{CLIENT_NAME}_orders.xlsx"
        make_sample(data_path)
    elif args.data:
        data_path = args.data
        if not data_path.exists():
            sys.exit(f"FATAL — file not found: {data_path}")
    elif not args.simple:
        ap.error("give --data PATH, --simple for the report-only deck, or "
                 "--sample to generate test data")

    if args.probe:
        return probe(data_path)

    df = load_orders(data_path)

    # The reference figures cover the whole book, so reconcile BEFORE any algo
    # filter - otherwise a scoped review always looks like a mismatch.
    ref = reference_check(df) if not args.sample else pd.DataFrame()
    if args.sample:
        log("")
        log("  (reference check skipped - synthetic data will not reconcile)")

    if EXCLUDE_MARKETS and "market" in df:
        out_ = df["market"].isin(EXCLUDE_MARKETS)
        if out_.any():
            log("")
            n_ord = int(out_.sum())
            log(f"  excluded from the review: {_and(EXCLUDE_MARKETS)} - "
                f"{n_ord:,} order{'' if n_ord == 1 else 's'}, "
                f"{f_money(df.loc[out_, 'notional'].sum())}")
            df = df[~out_].copy()

    if ALGOS_STUDIED:
        keep = df["algo"].isin(ALGOS_STUDIED)
        share = 100 * df.loc[keep, "notional"].sum() / df["notional"].sum()
        log("")
        log(f"  scope: {', '.join(ALGOS_STUDIED)} only -> {int(keep.sum()):,} "
            f"of {len(df):,} orders, {share:.1f}% of executed value")
        dropped = sorted(set(df.loc[~keep, "algo"].unique()))
        if dropped:
            log(f"  excluded from the review: {', '.join(dropped)}")
        df = df[keep].copy()
        if df.empty:
            sys.exit("FATAL - no orders left after the algo filter. Check the "
                     f"algo spellings in CLIENTS['{args.client}']['algos'].")

    sanity_report(df)
    # a client config with no period string gets it from the data
    if not PERIOD_LABEL and "date" in df and df["date"].notna().any():
        # %-d is not portable to Windows, so strip the leading zero by hand
        def _d(ts):
            return f"{ts.day} {ts:%B %Y}"
        PERIOD_LABEL = f"{_d(df['date'].min())} - {_d(df['date'].max())}"
        log(f"  period taken from the data: {PERIOD_LABEL}")

    ctx = analyse(df)
    ctx["reference_check"] = ref
    make_charts(ctx, charts_dir)
    write_tables(ctx, out / "tables.xlsx")

    if ctx["regression"]:
        (out / "dark_regression.txt").write_text(ctx["regression"][0],
                                                 encoding="utf-8")
        log(f"  saved dark_regression.txt")

    if not args.no_deck:
        log("")
        log("-" * 74)
        log("DECK")
        log("-" * 74)
        nar = build_narrative(ctx)
        build_deck(ctx, nar, out / f"{CLIENT_NAME}_TCA_deck.pptx")

    (out / "run_log.txt").write_text("\n".join(LOG), encoding="utf-8")
    log("")
    log("Done. Everything is in " + str(out.resolve()))
    if IS_SAMPLE[0]:
        log("")
        log("  !! THIS RUN USED SYNTHETIC DATA. Every number is invented.")
        log("     Re-run with --data <the real export> before showing anyone.")
    log("  charts/*.png are standalone — lift them into the corporate template.")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
